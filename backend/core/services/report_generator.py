"""
Report Generator Service
Handles PDF and PowerPoint report generation with scenario data integration
"""

import io
import json
import os
import uuid
import logging
from datetime import datetime, timedelta
from typing import Dict, List, Optional
from decimal import Decimal

from django.conf import settings
from django.core.files.base import ContentFile
from django.utils import timezone
from django.template.loader import render_to_string

# PDF Generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.linecharts import HorizontalLineChart
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.barcharts import VerticalBarChart

# PowerPoint Generation  
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..models import Client, Scenario, IncomeSource
from .data_service import ReportDataService
from .chart_service import ReportChartService

logger = logging.getLogger(__name__)


class ReportGenerator:
    """Main report generation service for PDF and PowerPoint exports"""
    
    def __init__(self):
        self.data_service = ReportDataService()
        self.chart_service = ReportChartService()
    
    def generate_report(self, report_config: Dict) -> Dict:
        """
        Main entry point for report generation.
        Generates both PDF and PowerPoint versions based on report configuration.
        
        Args:
            report_config: Dictionary containing:
                - client_id: ID of the client
                - scenario_ids: List of scenario IDs to include
                - template_type: Type of template to use
                - export_format: 'pdf', 'pptx', or 'both'
                - report_title: Title for the report
                - advisor_info: Advisor branding information
        
        Returns:
            Dictionary with file URLs and generation metadata
        """
        try:
            logger.info(f"Starting report generation for config: {report_config}")
            
            # Gather all necessary data
            report_data = self._prepare_report_data(report_config)
            
            results = {}
            
            # Generate PDF if requested
            if report_config.get('export_format') in ['pdf', 'both']:
                pdf_result = self._generate_pdf(report_config, report_data)
                results['pdf'] = pdf_result
            
            # Generate PowerPoint if requested  
            if report_config.get('export_format') in ['pptx', 'both']:
                pptx_result = self._generate_powerpoint(report_config, report_data)
                results['pptx'] = pptx_result
            
            logger.info(f"Report generation completed successfully")
            
            return {
                'success': True,
                'files': results,
                'generated_at': timezone.now().isoformat(),
                'report_title': report_config.get('report_title', 'Retirement Analysis Report')
            }
            
        except Exception as e:
            logger.error(f"Report generation failed: {str(e)}", exc_info=True)
            raise e
    
    def _prepare_report_data(self, config: Dict) -> Dict:
        """Prepare all data needed for report generation"""
        
        # Get client information
        client_id = config.get('client_id')
        if not client_id:
            raise ValueError("client_id is required for report generation")
        
        try:
            client = Client.objects.get(id=client_id)
        except Client.DoesNotExist:
            raise ValueError(f"Client with ID {client_id} not found")
        
        # Get scenarios
        scenario_ids = config.get('scenario_ids', [])
        scenarios = []
        scenario_data = {}
        
        if scenario_ids:
            scenarios = Scenario.objects.filter(id__in=scenario_ids, client=client)
        else:
            # Use all client scenarios if none specified
            scenarios = client.scenarios.all()[:3]  # Limit to 3 scenarios
        
        # Get scenario calculation data
        for scenario in scenarios:
            scenario_data[scenario.id] = self.data_service.get_scenario_data(scenario)
        
        # Prepare chart data
        chart_data = self.chart_service.prepare_chart_data(list(scenarios))
        
        return {
            'client': self.data_service.serialize_client(client),
            'scenarios': scenario_data,
            'charts': chart_data,
            'advisor': config.get('advisor_info', {}),
            'report_title': config.get('report_title', 'Retirement Analysis Report'),
            'generated_at': timezone.now(),
        }
    
    def _generate_pdf(self, config: Dict, data: Dict) -> Dict:
        """Generate PDF version of the report"""
        logger.info("Generating PDF report")
        
        buffer = io.BytesIO()
        
        # Create PDF document
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )
        
        # Get styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            alignment=1  # Center alignment
        )
        
        # Build content
        story = []
        
        # Add title page
        story.extend(self._create_pdf_title_page(data, title_style))
        
        # Add executive summary
        story.extend(self._create_pdf_executive_summary(data, styles))
        
        # Add scenario analysis
        story.extend(self._create_pdf_scenario_analysis(data, styles))
        
        # Add charts and visualizations
        story.extend(self._create_pdf_charts(data, styles))
        
        # Add recommendations
        story.extend(self._create_pdf_recommendations(data, styles))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        
        # Save to file system
        filename = f"report_{uuid.uuid4().hex[:8]}_{int(timezone.now().timestamp())}.pdf"
        file_content = ContentFile(buffer.getvalue(), name=filename)
        
        # For now, save to media directory
        file_path = f"reports/pdf/{filename}"
        full_path = os.path.join(settings.MEDIA_ROOT, "reports", "pdf", filename)
        
        # Ensure directory exists
        os.makedirs(os.path.dirname(full_path), exist_ok=True)
        
        # Save file
        with open(full_path, 'wb') as f:
            f.write(file_content.read())
        
        file_url = f"{settings.MEDIA_URL}reports/pdf/{filename}"
        
        logger.info(f"PDF generated successfully: {file_url}")
        
        return {
            'filename': filename,
            'path': str(full_path),
            'url': file_url,
            'size': len(buffer.getvalue()),
            'pages': len(story) // 20  # Rough estimate
        }
    
    def _generate_powerpoint(self, config: Dict, data: Dict) -> Dict:
        """Generate PowerPoint version of the report"""
        logger.info("Generating PowerPoint report")
        
        # Create presentation
        prs = Presentation()
        
        # Apply branding if available
        self._apply_pptx_branding(prs, data.get('advisor', {}))
        
        # Add title slide
        self._create_pptx_title_slide(prs, data)
        
        # Add executive summary slide
        self._create_pptx_executive_summary_slide(prs, data)
        
        # Add scenario slides
        self._create_pptx_scenario_slides(prs, data)
        
        # Add chart slides
        self._create_pptx_chart_slides(prs, data)
        
        # Add recommendations slide
        self._create_pptx_recommendations_slide(prs, data)
        
        # Save to buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        # Save to file system
        filename = f"report_{uuid.uuid4().hex[:8]}_{int(timezone.now().timestamp())}.pptx"
        file_content = ContentFile(buffer.getvalue(), name=filename)
        
        # For now, save to media directory
        file_path = f"reports/pptx/{filename}"
        full_path = os.path.join(settings.MEDIA_ROOT, "reports", "pptx", filename)
        
        # Ensure directory exists
        os.makedirs(os.path.dirname(full_path), exist_ok=True)
        
        # Save file
        with open(full_path, 'wb') as f:
            f.write(file_content.read())
        
        file_url = f"{settings.MEDIA_URL}reports/pptx/{filename}"
        
        logger.info(f"PowerPoint generated successfully: {file_url}")
        
        return {
            'filename': filename,
            'path': str(full_path),
            'url': file_url,
            'size': len(buffer.getvalue()),
            'slides': len(prs.slides)
        }
    
    # PDF Content Creation Methods
    def _create_pdf_title_page(self, data: Dict, title_style) -> List:
        """Create PDF title page content"""
        styles = getSampleStyleSheet()
        story = []
        
        # Title
        story.append(Paragraph(data['report_title'], title_style))
        story.append(Spacer(1, 0.5*inch))
        
        # Client information
        client = data['client']
        client_info = f"""
        <para fontSize=16>
        <b>Prepared for:</b><br/>
        {client['full_name']}<br/>
        Age: {client['age']}<br/>
        Report Date: {data['generated_at'].strftime('%B %d, %Y')}
        </para>
        """
        story.append(Paragraph(client_info, styles['Normal']))
        story.append(Spacer(1, 0.5*inch))
        
        # Advisor information
        advisor = data.get('advisor', {})
        if advisor:
            advisor_info = f"""
            <para fontSize=14>
            <b>Prepared by:</b><br/>
            {advisor.get('name', 'RetirementAdvisorPro')}<br/>
            {advisor.get('company', '')}<br/>
            {advisor.get('email', '')}<br/>
            {advisor.get('phone', '')}
            </para>
            """
            story.append(Paragraph(advisor_info, styles['Normal']))
        
        return story
    
    def _create_pdf_executive_summary(self, data: Dict, styles) -> List:
        """Create executive summary section"""
        story = []
        story.append(Paragraph("<b>Executive Summary</b>", styles['Heading2']))
        story.append(Spacer(1, 0.2*inch))
        
        # Generate summary based on scenario data
        scenarios = data['scenarios']
        if scenarios:
            first_scenario = list(scenarios.values())[0]
            summary_text = f"""
            This analysis presents retirement planning scenarios for {data['client']['full_name']}. 
            Based on the current financial position and retirement goals, we have evaluated 
            {len(scenarios)} scenario(s) to determine optimal strategies.
            
            Key findings include retirement income projections, tax optimization opportunities,
            and recommendations for achieving financial security in retirement.
            """
            
            story.append(Paragraph(summary_text, styles['Normal']))
        
        story.append(Spacer(1, 0.3*inch))
        return story
    
    def _create_pdf_scenario_analysis(self, data: Dict, styles) -> List:
        """Create scenario analysis section"""
        story = []
        story.append(Paragraph("<b>Scenario Analysis</b>", styles['Heading2']))
        story.append(Spacer(1, 0.2*inch))
        
        for scenario_id, scenario_data in data['scenarios'].items():
            scenario = scenario_data['scenario']
            
            # Scenario header
            story.append(Paragraph(f"<b>{scenario['name']}</b>", styles['Heading3']))
            
            # Scenario details
            scenario_text = f"""
            Retirement Age: {scenario.get('retirement_age', 'N/A')}<br/>
            Current Assets: ${scenario.get('total_assets', 0):,.2f}<br/>
            Projected Retirement Income: ${scenario.get('annual_income', 0):,.2f}<br/>
            """
            
            story.append(Paragraph(scenario_text, styles['Normal']))
            story.append(Spacer(1, 0.2*inch))
        
        return story
    
    def _create_pdf_charts(self, data: Dict, styles) -> List:
        """Create charts section"""
        story = []
        story.append(Paragraph("<b>Financial Projections</b>", styles['Heading2']))
        story.append(Spacer(1, 0.2*inch))
        
        # Add placeholder for charts
        chart_text = """
        Chart visualizations would be embedded here showing:
        - Asset growth projections
        - Income timeline analysis  
        - Tax impact analysis
        - Monte Carlo simulation results
        """
        
        story.append(Paragraph(chart_text, styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        return story
    
    def _create_pdf_recommendations(self, data: Dict, styles) -> List:
        """Create recommendations section"""
        story = []
        story.append(Paragraph("<b>Recommendations</b>", styles['Heading2']))
        story.append(Spacer(1, 0.2*inch))
        
        recommendations = [
            "Continue regular contributions to retirement accounts",
            "Consider Roth conversion strategies for tax optimization",
            "Review and adjust asset allocation based on risk tolerance",
            "Monitor IRMAA thresholds for Medicare premium management",
            "Schedule annual review to adjust strategy as needed"
        ]
        
        for i, rec in enumerate(recommendations, 1):
            story.append(Paragraph(f"{i}. {rec}", styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
        
        return story
    
    # PowerPoint Content Creation Methods
    def _apply_pptx_branding(self, prs: Presentation, advisor: Dict):
        """Apply advisor branding to PowerPoint presentation"""
        # This would apply custom colors, logos, etc.
        # For now, we'll use default styling
        pass
    
    def _create_pptx_title_slide(self, prs: Presentation, data: Dict):
        """Create PowerPoint title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data['report_title']
        
        client = data['client']
        subtitle_text = f"Prepared for {client['full_name']}\n"
        subtitle_text += f"Age: {client['age']}\n"
        subtitle_text += f"Report Date: {data['generated_at'].strftime('%B %d, %Y')}"
        
        subtitle.text = subtitle_text
    
    def _create_pptx_executive_summary(self, prs: Presentation, data: Dict):
        """Create executive summary slide"""
        slide_layout = prs.slide_layouts[1]  # Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Executive Summary"
        
        summary_points = [
            f"Analysis for {data['client']['full_name']}",
            f"Evaluated {len(data['scenarios'])} retirement scenario(s)",
            "Identified key opportunities for optimization",
            "Developed strategic recommendations"
        ]
        
        content.text = "\n".join(f"• {point}" for point in summary_points)
    
    def _create_pptx_scenario_slides(self, prs: Presentation, data: Dict):
        """Create scenario analysis slides"""
        for scenario_id, scenario_data in data['scenarios'].items():
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            scenario = scenario_data['scenario']
            title.text = f"Scenario: {scenario['name']}"
            
            scenario_points = [
                f"Retirement Age: {scenario.get('retirement_age', 'N/A')}",
                f"Current Assets: ${scenario.get('total_assets', 0):,.2f}",
                f"Projected Annual Income: ${scenario.get('annual_income', 0):,.2f}"
            ]
            
            content.text = "\n".join(f"• {point}" for point in scenario_points)
    
    def _create_pptx_chart_slides(self, prs: Presentation, data: Dict):
        """Create chart slides"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Financial Projections"
        content.text = "Chart visualizations would be embedded here showing asset growth, income timeline, and risk analysis."
    
    def _create_pptx_recommendations_slide(self, prs: Presentation, data: Dict):
        """Create recommendations slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Recommendations"
        
        recommendations = [
            "Continue regular retirement contributions",
            "Consider tax optimization strategies",
            "Review asset allocation annually",
            "Monitor Medicare premium thresholds",
            "Schedule regular strategy reviews"
        ]
        
        content.text = "\n".join(f"• {rec}" for rec in recommendations)
    
    def _generate_powerpoint(self, config: Dict, data: Dict) -> Dict:
        """Generate PowerPoint version of the report"""
        logger.info("Generating PowerPoint report")
        
        # Create presentation
        prs = Presentation()
        
        # Add title slide
        self._create_pptx_title_slide(prs, data)
        
        # Add executive summary slide
        self._create_pptx_executive_summary_slide(prs, data)
        
        # Add scenario slides
        for scenario_id, scenario_data in data['scenarios'].items():
            self._create_pptx_scenario_slide(prs, scenario_id, scenario_data)
        
        # Add chart slides
        self._create_pptx_chart_slides(prs, data)
        
        # Add recommendations slide
        self._create_pptx_recommendations_slide(prs, data)
        
        # Save to buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        # Save to file system
        filename = f"report_{uuid.uuid4().hex[:8]}_{int(timezone.now().timestamp())}.pptx"
        file_content = ContentFile(buffer.getvalue(), name=filename)
        
        # For now, save to media directory
        file_path = f"reports/pptx/{filename}"
        full_path = os.path.join(settings.MEDIA_ROOT, "reports", "pptx", filename)
        
        # Ensure directory exists
        os.makedirs(os.path.dirname(full_path), exist_ok=True)
        
        # Save file
        with open(full_path, 'wb') as f:
            f.write(file_content.read())
        
        file_size = os.path.getsize(full_path)
        
        logger.info(f"PowerPoint generated successfully: {full_path}")
        
        return {
            'filename': filename,
            'path': full_path,
            'url': f"/media/reports/pptx/{filename}",
            'size': file_size,
            'slides': len(prs.slides)
        }
    
    def _create_pptx_title_slide(self, prs, data):
        """Create title slide for PowerPoint presentation"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        client_name = data.get('client', {}).get('name', 'Client')
        title.text = f"Retirement Report - {client_name}"
        
        # Set subtitle
        if slide.placeholders[1]:
            subtitle = slide.placeholders[1]
            subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    
    def _create_pptx_executive_summary_slide(self, prs, data):
        """Create executive summary slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Add content
        if slide.placeholders[1]:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = "Retirement Planning Analysis Overview"
            
            # Add bullet points
            p = tf.add_paragraph()
            client_name = data.get('client', {}).get('name', 'Client')
            p.text = f"Client: {client_name}"
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = f"Scenarios Analyzed: {len(data.get('scenarios', {}))}"
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = "Comprehensive retirement income analysis completed"
            p.level = 1
    
    def _create_pptx_scenario_slide(self, prs, scenario_id, scenario_data):
        """Create a slide for a specific scenario"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = f"Scenario: {scenario_data.get('name', 'Analysis')}"
        
        # Add content
        if slide.placeholders[1]:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = "Scenario Details"
            
            # Add key metrics
            summary = scenario_data.get('summary', {})
            if summary:
                p = tf.add_paragraph()
                p.text = f"Retirement Age: {summary.get('retirement_age', 'N/A')}"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = f"Life Expectancy: {summary.get('mortality_age', 'N/A')}"
                p.level = 1
                
                if 'total_income' in summary:
                    p = tf.add_paragraph()
                    p.text = f"Total Projected Income: ${summary['total_income']:,.2f}"
                    p.level = 1
    
    def _create_pptx_scenario_slides(self, prs, data):
        """Create slides for all scenarios"""
        scenarios = data.get('scenarios', {})
        for scenario_id, scenario_data in scenarios.items():
            self._create_pptx_scenario_slide(prs, scenario_id, scenario_data)
    
    def _create_pptx_chart_slides(self, prs, data):
        """Create slides with charts and visualizations"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Financial Projections"
        
        # Add placeholder content (charts would need more complex implementation)
        if slide.placeholders[1]:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = "Key Financial Charts and Projections"
            
            p = tf.add_paragraph()
            p.text = "• Asset allocation over time"
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = "• Income vs. expenses projection"
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = "• Monte Carlo simulation results"
            p.level = 1
    
    def _create_pptx_recommendations_slide(self, prs, data):
        """Create recommendations slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Recommendations"
        
        # Add content
        if slide.placeholders[1]:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = "Key Recommendations"
            
            p = tf.add_paragraph()
            p.text = "Consider tax-efficient withdrawal strategies"
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = "Review Social Security claiming strategies"
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = "Optimize asset allocation based on risk tolerance"
            p.level = 1