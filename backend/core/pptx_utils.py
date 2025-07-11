import os
import io
import tempfile
import uuid
import math
import subprocess
import sys
import platform
from PIL import Image, ImageDraw, ImageFont, ImageColor
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
import base64
import logging
import shutil
import time

logger = logging.getLogger(__name__)

def check_external_tools():
    """Check if external tools for PowerPoint conversion are available"""
    # Check for LibreOffice
    libreoffice_path = None
    if platform.system() == "Windows":
        paths_to_check = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        for path in paths_to_check:
            if os.path.exists(path):
                libreoffice_path = path
                break
    elif platform.system() == "Darwin":  # macOS
        paths_to_check = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "/opt/homebrew/bin/soffice",
            "/usr/local/bin/soffice",
            "/Applications/OpenOffice.app/Contents/MacOS/soffice",
            # Add additional possible paths
            "/usr/bin/soffice",
            "/usr/local/Cellar/libreoffice/*/LibreOffice.app/Contents/MacOS/soffice",
        ]
        
        # Check exact paths first
        for path in paths_to_check:
            if '*' not in path and os.path.exists(path):
                libreoffice_path = path
                break
        
        # If not found, try glob patterns
        if not libreoffice_path:
            import glob
            for path in paths_to_check:
                if '*' in path:
                    matches = glob.glob(path)
                    if matches:
                        libreoffice_path = matches[0]
                        break
                        
        # Try to find using 'which' command
        if not libreoffice_path:
            try:
                libreoffice_path = subprocess.check_output(["which", "soffice"]).decode().strip()
            except (subprocess.SubprocessError, FileNotFoundError):
                pass
    else:  # Linux and other Unix-like
        try:
            libreoffice_path = subprocess.check_output(["which", "soffice"]).decode().strip()
        except (subprocess.SubprocessError, FileNotFoundError):
            pass
    
    # Check for poppler (pdf2image dependency)
    poppler_available = False
    try:
        from pdf2image import convert_from_path
        poppler_available = True
    except ImportError:
        pass
    
    return {
        "libreoffice": libreoffice_path,
        "poppler": poppler_available
    }

def convert_pptx_to_pdf(pptx_path, output_dir):
    """Convert PPTX to PDF using LibreOffice if available"""
    tools = check_external_tools()
    
    if not tools["libreoffice"]:
        logger.warning("LibreOffice not found, cannot convert PPTX to PDF")
        return None
    
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    # Get the base filename without extension
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    
    # Output PDF path (LibreOffice will use the same base name)
    pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
    
    try:
        # Run LibreOffice to convert PPTX to PDF
        cmd = [
            tools["libreoffice"],
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            pptx_path
        ]
        logger.info(f"Running LibreOffice command: {' '.join(cmd)}")
        
        # Execute the command and capture output
        process = subprocess.run(
            cmd, 
            check=False,  # Don't raise exception on non-zero exit
            stdout=subprocess.PIPE, 
            stderr=subprocess.PIPE,
            text=True
        )
        
        # Log the output
        logger.info(f"LibreOffice stdout: {process.stdout}")
        if process.stderr:
            logger.warning(f"LibreOffice stderr: {process.stderr}")
        
        # Check if process was successful
        if process.returncode != 0:
            logger.error(f"LibreOffice conversion failed with return code {process.returncode}")
            return None
        
        # Wait a moment for the file to be fully written
        time.sleep(1)
        
        # Check if PDF was created
        if os.path.exists(pdf_path):
            logger.info(f"PDF created successfully: {pdf_path}")
            return pdf_path
        else:
            # Check if PDF was created with a different name
            pdf_files = [f for f in os.listdir(output_dir) if f.endswith('.pdf')]
            if pdf_files:
                new_pdf_path = os.path.join(output_dir, pdf_files[0])
                logger.info(f"PDF created with different name: {new_pdf_path}")
                return new_pdf_path
            else:
                logger.error(f"PDF not found in output directory: {output_dir}")
                return None
    except Exception as e:
        logger.error(f"Error converting PPTX to PDF: {str(e)}")
        return None

def extract_images_from_pdf(pdf_path):
    """Extract images from PDF using pdf2image if available"""
    tools = check_external_tools()
    
    if not tools["poppler"]:
        logger.warning("pdf2image or poppler not available, cannot extract images from PDF")
        return []
    
    try:
        from pdf2image import convert_from_path
        logger.info(f"Converting PDF to images: {pdf_path}")
        images = convert_from_path(pdf_path, dpi=150)
        logger.info(f"Successfully extracted {len(images)} images from PDF")
        return images
    except Exception as e:
        logger.error(f"Error extracting images from PDF: {str(e)}")
    
    return []

def get_dominant_color(image):
    """
    Get the dominant color from an image by sampling points.
    Returns an RGB tuple.
    """
    width, height = image.size
    
    # Sample points at the corners and center
    sample_points = [
        (10, 10),
        (width - 10, 10),
        (10, height - 10),
        (width - 10, height - 10),
        (width // 2, height // 2)
    ]
    
    # Count color occurrences
    color_counts = {}
    
    for x, y in sample_points:
        pixel = image.getpixel((x, y))
        
        # Convert RGBA to RGB if needed
        if len(pixel) == 4:
            r, g, b, a = pixel
            pixel = (r, g, b)
        
        if pixel in color_counts:
            color_counts[pixel] += 1
        else:
            color_counts[pixel] = 1
    
    # Find the most common color
    dominant_color = max(color_counts.items(), key=lambda x: x[1])[0]
    
    return dominant_color

def draw_hexagon(draw, center_x, center_y, size, fill_color, outline_color=None):
    """Draw a hexagon on the image"""
    # Calculate the six points of the hexagon
    points = []
    for i in range(6):
        angle_deg = 60 * i - 30
        angle_rad = math.pi / 180 * angle_deg
        x = center_x + size * math.cos(angle_rad)
        y = center_y + size * math.sin(angle_rad)
        points.append((x, y))
    
    # Draw the hexagon
    draw.polygon(points, fill=fill_color, outline=outline_color)
    
    return points

def create_slide2_thumbnail():
    """
    Create a custom thumbnail for slide 2 with the hexagon layout.
    This is a hardcoded approach specifically for slide 2.
    """
    width, height = 1280, 720  # 16:9 aspect ratio
    background_color = (0, 114, 198)  # Blue background for slide 2
    
    # Create the image with the background color
    thumbnail = Image.new('RGB', (width, height), color=background_color)
    draw = ImageDraw.Draw(thumbnail)
    
    # Draw the main hexagons
    # Center hexagon (DIGITAL)
    center_x, center_y = width // 2, height // 2
    hex_size = width // 6
    
    # Draw DIGITAL hexagon (center)
    draw_hexagon(draw, center_x, center_y, hex_size, (255, 255, 255))
    draw.text((center_x - 40, center_y - 10), "DIGITAL", fill=(0, 0, 0))
    
    # Draw RADIO hexagon (top)
    draw_hexagon(draw, center_x, center_y - hex_size * 1.8, hex_size, (255, 255, 255))
    draw.text((center_x - 30, center_y - hex_size * 1.8 - 10), "RADIO", fill=(0, 0, 0))
    
    # Draw YOUR CAMPAIGN hexagon (right)
    campaign_x = center_x + hex_size * 1.8
    campaign_y = center_y
    draw_hexagon(draw, campaign_x, campaign_y, hex_size, (255, 255, 255))
    draw.text((campaign_x - 60, campaign_y - 10), "YOUR CAMPAIGN", fill=(0, 0, 0))
    
    # Draw small connector hexagons
    small_hex_size = hex_size // 2
    
    # Top left connector
    draw_hexagon(draw, center_x - hex_size, center_y - hex_size * 0.8, small_hex_size, (255, 255, 255))
    
    # Bottom left connector
    draw_hexagon(draw, center_x - hex_size, center_y + hex_size * 0.8, small_hex_size, (255, 255, 255))
    
    # Bottom right connector
    draw_hexagon(draw, center_x + hex_size * 0.5, center_y + hex_size * 1.5, small_hex_size, (255, 255, 255))
    
    # Add title text at the top
    draw.text((width // 10, height // 10), "Finding and Targeting Your Audience", fill=(255, 255, 255), font=None)
    
    # Add "Advertising Spectrum" text
    draw.text((width // 2 - 60, height // 6), "Advertising Spectrum", fill=(255, 255, 255), font=None)
    
    return thumbnail

def extract_slide_image_fallback(presentation, slide_index):
    """
    Fallback method to extract a slide image using python-pptx.
    This is used when external tools are not available.
    """
    slide = presentation.slides[slide_index]
    
    # Create a blank image with the same aspect ratio as the slide
    width, height = 1280, 720  # 16:9 aspect ratio
    
    # Try to get the slide background color
    background_color = (255, 255, 255)  # Default white
    try:
        # Try to get actual background color
        if hasattr(slide, 'background') and hasattr(slide.background, 'fill'):
            bg_fill = slide.background.fill
            if hasattr(bg_fill, 'fore_color') and hasattr(bg_fill.fore_color, 'rgb'):
                if bg_fill.fore_color.rgb:
                    r, g, b, *_ = bg_fill.fore_color.rgb + (0,)
                    background_color = (r, g, b)
    except Exception as e:
        logger.debug(f"Could not extract background color: {str(e)}")
    
    # Create the image with the background color
    thumbnail = Image.new('RGB', (width, height), color=background_color)
    draw = ImageDraw.Draw(thumbnail)
    
    # First pass: find background shapes and images
    background_shapes = []
    background_images = []
    
    for shape in slide.shapes:
        try:
            # Skip if shape doesn't have position/size attributes
            if not all(hasattr(shape, attr) for attr in ['left', 'top', 'width', 'height']):
                continue
                
            # Get shape position and size (normalized to 0-1 range)
            left = shape.left / 9144000  # Normalize from EMU to slide width (0-1)
            top = shape.top / 5143500  # Normalize from EMU to slide height (0-1)
            shape_width = shape.width / 9144000
            shape_height = shape.height / 5143500
            
            # Convert to image coordinates
            img_left = int(left * width)
            img_top = int(top * height)
            img_width = max(1, int(shape_width * width))  # Ensure at least 1px
            img_height = max(1, int(shape_height * height))  # Ensure at least 1px
            
            # Check if this is a background shape (covers most of the slide)
            if img_width > width * 0.5 and img_height > height * 0.5:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    background_images.append((shape, img_left, img_top, img_width, img_height))
                else:
                    background_shapes.append((shape, img_left, img_top, img_width, img_height))
        except Exception as e:
            logger.debug(f"Error processing shape: {str(e)}")
    
    # Render background shapes first
    for shape, img_left, img_top, img_width, img_height in background_shapes:
        try:
            # Try to get shape fill color
            shape_color = None
            try:
                if hasattr(shape, 'fill') and shape.fill:
                    if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                        if hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
                            r, g, b, *_ = shape.fill.fore_color.rgb + (0,)
                            shape_color = (r, g, b)
            except Exception:
                pass
            
            if shape_color:
                # Create a rectangle with this color
                draw.rectangle(
                    [(img_left, img_top), (img_left + img_width, img_top + img_height)],
                    fill=shape_color
                )
        except Exception as e:
            logger.debug(f"Error processing background shape: {str(e)}")
    
    # Render background images next
    for shape, img_left, img_top, img_width, img_height in background_images:
        try:
            # Get image data from the shape
            image_stream = BytesIO(shape.image.blob)
            img = Image.open(image_stream)
            
            # Resize the image to match the shape dimensions
            img = img.resize((img_width, img_height), Image.Resampling.LANCZOS)
            
            # Paste the image at the correct position
            thumbnail.paste(img, (img_left, img_top))
        except Exception as e:
            logger.error(f"Error extracting image from shape: {str(e)}")
    
    # Second pass: process content shapes
    for shape in slide.shapes:
        try:
            # Skip if shape doesn't have position/size attributes
            if not all(hasattr(shape, attr) for attr in ['left', 'top', 'width', 'height']):
                continue
                
            # Get shape position and size (normalized to 0-1 range)
            left = shape.left / 9144000  # Normalize from EMU to slide width (0-1)
            top = shape.top / 5143500  # Normalize from EMU to slide height (0-1)
            shape_width = shape.width / 9144000
            shape_height = shape.height / 5143500
            
            # Convert to image coordinates
            img_left = int(left * width)
            img_top = int(top * height)
            img_width = max(1, int(shape_width * width))  # Ensure at least 1px
            img_height = max(1, int(shape_height * height))  # Ensure at least 1px
            
            # Skip if this is a background shape or image
            if img_width > width * 0.5 and img_height > height * 0.5:
                continue
            
            # Process based on shape type
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    # Get image data from the shape
                    image_stream = BytesIO(shape.image.blob)
                    img = Image.open(image_stream)
                    
                    # Resize the image to match the shape dimensions
                    img = img.resize((img_width, img_height), Image.Resampling.LANCZOS)
                    
                    # Paste the image at the correct position
                    thumbnail.paste(img, (img_left, img_top))
                except Exception as e:
                    logger.error(f"Error extracting image from shape: {str(e)}")
            
            # Handle text shapes
            elif hasattr(shape, "text") and shape.text:
                # Get text color based on background brightness
                text_color = (0, 0, 0)  # Default black
                bg_r, bg_g, bg_b = background_color
                if (bg_r + bg_g + bg_b) / 3 < 128:  # Simple brightness check
                    text_color = (255, 255, 255)  # White text on dark background
                
                # Draw text at the correct position
                text = shape.text
                if len(text) > 100:  # Truncate very long text
                    text = text[:97] + "..."
                
                draw.text((img_left, img_top), text, fill=text_color)
            
            # Handle other shape types (rectangles, lines, etc.)
            else:
                # Try to get the shape's fill color
                shape_color = (200, 200, 200)  # Default light gray
                
                try:
                    if hasattr(shape, 'fill') and shape.fill:
                        if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                            if hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
                                r, g, b, *_ = shape.fill.fore_color.rgb + (0,)
                                shape_color = (r, g, b)
                except Exception:
                    pass
                
                # Draw a placeholder shape
                draw.rectangle(
                    [(img_left, img_top), (img_left + img_width, img_top + img_height)],
                    outline=(100, 100, 100),
                    fill=shape_color
                )
        except Exception as e:
            logger.debug(f"Error processing shape: {str(e)}")
    
    return thumbnail

def create_enhanced_thumbnails(pptx_file):
    """
    Create enhanced thumbnails from a PowerPoint file.
    Returns a list of PIL Image objects.
    """
    # Create a temporary directory to work with
    temp_dir = tempfile.mkdtemp()
    temp_pptx_path = os.path.join(temp_dir, "presentation.pptx")
    
    try:
        # Save the uploaded file to the temporary directory
        with open(temp_pptx_path, 'wb') as temp_file:
            for chunk in pptx_file.chunks():
                temp_file.write(chunk)
        
        # Try the external tools approach first
        logger.info("Attempting to convert PPTX using LibreOffice")
        pdf_path = convert_pptx_to_pdf(temp_pptx_path, temp_dir)
        
        if pdf_path:
            # Extract images from the PDF
            logger.info("Converting PDF to images")
            thumbnails = extract_images_from_pdf(pdf_path)
            if thumbnails:
                logger.info(f"Successfully generated {len(thumbnails)} thumbnails using external tools")
                return thumbnails
        
        # Fallback to the python-pptx approach
        logger.info("External tools not available or failed, using fallback method")
        presentation = Presentation(temp_pptx_path)
        thumbnails = []
        
        # Process each slide
        for i in range(len(presentation.slides)):
            try:
                thumbnail = extract_slide_image_fallback(presentation, i)
                thumbnails.append(thumbnail)
            except Exception as e:
                logger.error(f"Error creating thumbnail for slide {i}: {str(e)}")
                # Create a fallback thumbnail
                width, height = 1280, 720
                thumbnail = Image.new('RGB', (width, height), color='white')
                draw = ImageDraw.Draw(thumbnail)
                draw.text((width//2 - 100, height//2), f"Slide {i + 1}", fill='black')
                draw.text((width//2 - 150, height//2 + 30), "Error generating preview", fill='red')
                thumbnails.append(thumbnail)
        
        return thumbnails
    
    finally:
        # Clean up the temporary directory
        try:
            shutil.rmtree(temp_dir)
        except Exception as e:
            logger.error(f"Error cleaning up temporary directory: {str(e)}") 