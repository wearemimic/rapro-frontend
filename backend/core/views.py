import requests
import os
import sys
import tempfile
import shutil
from django.contrib.auth import authenticate
from rest_framework import status, generics, permissions
from rest_framework.decorators import api_view, permission_classes, throttle_classes, action
from rest_framework.permissions import AllowAny, IsAuthenticated
from rest_framework.response import Response
from rest_framework.views import APIView
from rest_framework_simplejwt.authentication import JWTAuthentication
from rest_framework_simplejwt.tokens import RefreshToken
from .throttles import LoginRateThrottle
from .serializers import CustomUserSerializer, UserSerializer, ClientSerializer
from django.shortcuts import render, get_object_or_404
from django.contrib.auth import get_user_model
from django.views.decorators.csrf import csrf_exempt
from django.utils.decorators import method_decorator
from .models import Scenario, Client, RealEstate, ReportTemplate, TemplateSlide, IncomeSource
from rest_framework.exceptions import PermissionDenied
from .serializers import ClientDetailSerializer, ClientEditSerializer, ClientCreateSerializer, RealEstateSerializer
from .serializers import ScenarioCreateSerializer, ScenarioUpdateSerializer, IncomeSourceUpdateSerializer
from .serializers import ReportTemplateSerializer, ReportTemplateDetailSerializer, TemplateSlideSerializer
from .scenario_processor import ScenarioProcessor
from .roth_conversion_processor import RothConversionProcessor
from django.http import HttpResponse, JsonResponse
import logging
from .serializers import AdvisorRegistrationSerializer
import stripe
from django.conf import settings
from rest_framework import viewsets
import json
from PIL import Image, ImageDraw
from io import BytesIO
import uuid
from .pptx_utils import create_enhanced_thumbnails
from pptx_to_png import convert_pptx_to_png
from decimal import Decimal, InvalidOperation
import copy

# Add the python-pptx library for handling PowerPoint files
try:
    from pptx import Presentation
except ImportError:
    # We'll handle this in the view with a proper error response
    pass

User = get_user_model()

# Initialize Stripe
stripe.api_key = settings.STRIPE_SECRET_KEY

def home_view(request):
    return render(request, 'home.html')

# @csrf_exempt
@api_view(['GET', 'POST'])
@throttle_classes([LoginRateThrottle])
@permission_classes([AllowAny])
def login_view(request):
    if request.method == 'GET':
        return Response({"message": "Use POST to login."}, status=200)

    email = request.data.get('email')
    password = request.data.get('password')

    if not email or not password:
        return Response({'message': 'Email and password are required'}, status=400)

    user_obj = User.objects.filter(email=email).first()
    if user_obj is None:
        return Response({'message': 'Invalid credentials'}, status=401)

    user = authenticate(request, username=user_obj.email, password=password)
    if user is None:
        return Response({'message': 'Invalid credentials'}, status=401)

    refresh = RefreshToken.for_user(user)
    return Response({
        'access': str(refresh.access_token),
        'refresh': str(refresh),
        'user': {
            'id': user.id,
            'email': user.email,
            'username': user.username
        }
    })

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def logout_view(request):
    try:
        # Blacklist the refresh token (if used)
        refresh_token = request.data.get("refresh_token")
        token = RefreshToken(refresh_token)
        token.blacklist()
    except Exception as e:
        return Response({"error": str(e)}, status=status.HTTP_400_BAD_REQUEST)

    return Response({"detail": "Successfully logged out."}, status=status.HTTP_200_OK)

@api_view(['POST'])
@permission_classes([AllowAny])
def register_view(request):
    data = request.data
    if User.objects.filter(email=data.get('email')).exists():
        return Response({'message': 'Email already in use'}, status=status.HTTP_400_BAD_REQUEST)

    user = User.objects.create_user(
        username=data.get('email'),
        email=data.get('email'),
        password=data.get('password'),
        first_name=data.get('firstName'),
        last_name=data.get('lastName')
    )
    # return Response({'message': 'User created'}, status=status.HTTP_201_CREATED)
    #user = User.objects.create_user(username=email, email=email, password=password)

    # ✅ Generate JWT token
    refresh = RefreshToken.for_user(user)
    return Response({
        'access': str(refresh.access_token),
        'refresh': str(refresh),
        'user': {
            'id': user.id,
            'email': user.email,
            'username': user.username
        }
    })

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def run_scenario_calculation(request, scenario_id):
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        # New instantiation: only pass scenario_id and debug
        processor = ScenarioProcessor(scenario_id=scenario.id, debug=True)
        result = processor.calculate()
        return Response(result)
    except Scenario.DoesNotExist:
        return Response({"error": "Scenario not found."}, status=404)
    
@api_view(['GET', 'PUT'])
@permission_classes([IsAuthenticated])
def profile_view(request):
    user = request.user
    if request.method == 'GET':
        serializer = UserSerializer(user)
        return Response(serializer.data)
    elif request.method == 'PUT':
        # Handle file uploads
        data = request.data.copy()
        
        # Handle empty string for website_url to prevent validation errors
        if 'website_url' in data and data['website_url'] == '':
            data['website_url'] = None
        
        # Create a serializer with the data and possible file
        serializer = UserSerializer(user, data=data, partial=True)
        if serializer.is_valid():
            # Handle logo file upload separately if present
            if 'logo' in request.FILES:
                user.logo = request.FILES['logo']
                user.save(update_fields=['logo'])
            
            # Save the rest of the data
            serializer.save()
            return Response(serializer.data, status=status.HTTP_200_OK)
        
        # Return detailed error messages
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
    
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def create_scenario(request, client_id):
    try:
        client = Client.objects.get(id=client_id, advisor=request.user)
    except Client.DoesNotExist:
        return Response({'error': 'Client not found or access denied.'}, status=404)

    data = request.data.copy()
    data['client'] = client.id

    serializer = ScenarioCreateSerializer(data=data, context={'request': request})
    if serializer.is_valid():
        scenario = serializer.save()
        return Response({'id': scenario.id}, status=201)
    else:
        return Response(serializer.errors, status=400)  

@csrf_exempt
def proxy_to_wealthbox(request):
    api_key = '020dfe66ac3f4213beefb36641a96721'  # Replace with your actual API key
    # Log the incoming request path
    print(f"Received request path: {request.path}")
    url = f"https://api.crmworkspace.com/{request.path.replace('/proxy', '')}"
    print(f"Proxying request to URL: {url}")
    headers = {'ACCESS_TOKEN': api_key}

    try:
        response = requests.get(url, headers=headers)
        # Log the response status code
        # self._log_debug(f"Response status code: {response.status_code}")
        return JsonResponse(response.json(), status=response.status_code)
    except requests.exceptions.RequestException as e:
        logging.error(f"RequestException: {str(e)}")
        return JsonResponse({'error': str(e)}, status=500)


class AdvisorClientListView(generics.ListAPIView):
    serializer_class = ClientSerializer
    authentication_classes = [JWTAuthentication]
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        return Client.objects.filter(advisor=self.request.user)
    
  
class ClientDetailView(generics.RetrieveAPIView):
    queryset = Client.objects.all()
    serializer_class = ClientDetailSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_object(self):
        client = get_object_or_404(Client, pk=self.kwargs['pk'])
        print(f"Client advisor ID: {client.advisor.id}")
        print(f"Requesting user ID: {self.request.user.id}")
        if client.advisor != self.request.user:
            raise PermissionDenied("You do not have permission to view this client.")
        return client
    
class ClientEditView(generics.RetrieveUpdateAPIView):
    queryset = Client.objects.all()
    serializer_class = ClientEditSerializer
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]

    def get_queryset(self):
        user = self.request.user
        if user.is_staff:
            return Client.objects.all()
        return Client.objects.filter(advisor=user)

    def get_object(self):
        client = get_object_or_404(self.get_queryset(), pk=self.kwargs["pk"])
        return client
    
class ClientCreateView(APIView):
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]

    def post(self, request):
        serializer = ClientCreateSerializer(data=self.request.data, context={'request': self.request})
        if serializer.is_valid():
            validated_data = serializer.validated_data.copy()
            validated_data.pop('advisor', None)  # Remove advisor if it's already in validated_data
            client = serializer.save(advisor=self.request.user)
            return Response(ClientDetailSerializer(client).data, status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
    
class ScenarioCreateView(APIView):
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]

    def post(self, request, client_id):
        serializer = ScenarioCreateSerializer(data=request.data, context={'request': request, 'client_id': client_id})
        if serializer.is_valid():
            scenario = serializer.save()
            return Response({'id': scenario.id, 'message': 'Scenario created successfully'}, status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def duplicate_scenario(request, scenario_id):
    """
    Duplicate a scenario with all its income sources and settings
    """
    try:
        # Get the original scenario
        original_scenario = Scenario.objects.get(id=scenario_id)
        
        # Check if the user owns the client associated with this scenario
        if original_scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Create a copy of the scenario
        duplicated_scenario = Scenario.objects.create(
            client=original_scenario.client,
            name=f"{original_scenario.name} (Copy)",
            description=original_scenario.description,
            retirement_age=original_scenario.retirement_age,
            medicare_age=original_scenario.medicare_age,
            spouse_retirement_age=original_scenario.spouse_retirement_age,
            spouse_medicare_age=original_scenario.spouse_medicare_age,
            mortality_age=original_scenario.mortality_age,
            spouse_mortality_age=original_scenario.spouse_mortality_age,
            retirement_year=original_scenario.retirement_year,
            share_with_client=original_scenario.share_with_client,
            part_b_inflation_rate=original_scenario.part_b_inflation_rate,
            part_d_inflation_rate=original_scenario.part_d_inflation_rate,
            FRA_amount=original_scenario.FRA_amount,
            roth_conversion_start_year=original_scenario.roth_conversion_start_year,
            roth_conversion_duration=original_scenario.roth_conversion_duration,
            roth_conversion_annual_amount=original_scenario.roth_conversion_annual_amount,
            apply_standard_deduction=original_scenario.apply_standard_deduction,
            income_vs_cost_percent=0,  # Reset to 0 for new scenario
            medicare_irmaa_percent=0   # Reset to 0 for new scenario
        )
        
        # Duplicate all income sources
        for income_source in original_scenario.income_sources.all():
            IncomeSource.objects.create(
                scenario=duplicated_scenario,
                owned_by=income_source.owned_by,
                income_type=income_source.income_type,
                income_name=income_source.income_name,
                current_asset_balance=income_source.current_asset_balance,
                monthly_amount=income_source.monthly_amount,
                monthly_contribution=income_source.monthly_contribution,
                age_to_begin_withdrawal=income_source.age_to_begin_withdrawal,
                age_to_end_withdrawal=income_source.age_to_end_withdrawal,
                rate_of_return=income_source.rate_of_return,
                cola=income_source.cola,
                exclusion_ratio=income_source.exclusion_ratio,
                tax_rate=income_source.tax_rate,
                max_to_convert=income_source.max_to_convert
            )
        
        return Response({
            'id': duplicated_scenario.id,
            'name': duplicated_scenario.name,
            'message': 'Scenario duplicated successfully'
        }, status=status.HTTP_201_CREATED)
        
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_scenario_detail(request, scenario_id):
    """
    Get detailed scenario data for editing/duplicating
    """
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Get all income sources for this scenario
        income_sources = scenario.income_sources.all()
        income_data = []
        
        for income in income_sources:
            income_dict = {
                'id': str(uuid.uuid4()),  # Generate new UUID for frontend
                'income_type': income.income_type,
                'income_name': income.income_name,  # Add the missing income_name field
                'owned_by': income.owned_by,
                'start_age': income.age_to_begin_withdrawal,
                'end_age': income.age_to_end_withdrawal,
                'current_balance': float(income.current_asset_balance or 0),
                'monthly_contribution': float(income.monthly_contribution or 0),
                'growth_rate': income.rate_of_return,
                'withdrawal_amount': float(income.monthly_amount or 0),
                'amount_per_month': float(income.monthly_amount or 0),
                'amount_at_fra': float(income.monthly_amount or 0),
                'cola': income.cola,
                'exclusion_ratio': income.exclusion_ratio,
                'tax_rate': income.tax_rate,
                'max_to_convert': float(income.max_to_convert or 0),
                'age_established': income.age_established,
                'is_contributing': income.is_contributing,
                'employer_match': income.employer_match,
                'age_last_contribution': income.age_last_contribution
            }
            
            # Add specific fields for different income types
            if income.income_type == 'social_security':
                income_dict['amount_at_fra'] = float(income.monthly_amount or 0)
            elif income.income_type == 'Life_Insurance':
                income_dict['loan_amount'] = float(income.monthly_amount or 0)
            elif income.income_type == 'Annuity':
                income_dict['percent_taxable'] = income.exclusion_ratio * 100
            
            income_data.append(income_dict)
        
        # Check if this is for duplication (add "Copy") or editing (keep original name)
        is_for_duplication = request.GET.get('mode') == 'duplicate'
        scenario_name = f"{scenario.name} (Copy)" if is_for_duplication else scenario.name
        
        scenario_data = {
            'name': scenario_name,
            'description': scenario.description,
            'primary_retirement_age': scenario.retirement_age,
            'primary_medicare_age': scenario.medicare_age,
            'primary_lifespan': scenario.mortality_age,
            'spouse_retirement_age': scenario.spouse_retirement_age,
            'spouse_medicare_age': scenario.spouse_medicare_age,
            'spouse_lifespan': scenario.spouse_mortality_age,
            'model_tax_change': '',  # Reset to default
            'reduction_2030_ss': scenario.reduction_2030_ss,
            'ss_adjustment_year': scenario.ss_adjustment_year,
            'ss_adjustment_direction': scenario.ss_adjustment_direction,
            'ss_adjustment_type': scenario.ss_adjustment_type,
            'ss_adjustment_amount': scenario.ss_adjustment_amount,
            'apply_standard_deduction': scenario.apply_standard_deduction,
            'federal_standard_deduction': float(scenario.federal_standard_deduction or 0),
            'state_standard_deduction': float(scenario.state_standard_deduction or 0),
            'custom_annual_deduction': float(scenario.custom_annual_deduction or 0),
            'primary_blind': scenario.primary_blind,
            'spouse_blind': scenario.spouse_blind,
            'is_dependent': scenario.is_dependent,
            'part_b_inflation_rate': str(int(scenario.part_b_inflation_rate)),
            'part_d_inflation_rate': str(int(scenario.part_d_inflation_rate)),
            'primary_state': scenario.primary_state or '',  # Use stored value or default
            'income': income_data
        }
        
        return Response(scenario_data, status=status.HTTP_200_OK)
        
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_scenario_for_editing(request, scenario_id):
    """
    Get scenario data formatted specifically for editing in ScenarioCreate.vue
    """
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Get all income sources for this scenario
        income_sources = scenario.income_sources.all()
        income_data = []
        
        for income in income_sources:
            income_dict = {
                'id': income.id,  # Use actual database ID
                'income_type': income.income_type,
                'income_name': income.income_name,
                'owned_by': income.owned_by,
                'start_age': income.age_to_begin_withdrawal,
                'end_age': income.age_to_end_withdrawal,
                'current_balance': float(income.current_asset_balance or 0),
                'monthly_contribution': float(income.monthly_contribution or 0),
                'growth_rate': income.rate_of_return,
                'withdrawal_amount': float(income.monthly_amount or 0),
                'amount_per_month': float(income.monthly_amount or 0),
                'amount_at_fra': float(income.monthly_amount or 0),
                'cola': income.cola,
                'exclusion_ratio': income.exclusion_ratio,
                'tax_rate': income.tax_rate,
                'max_to_convert': float(income.max_to_convert or 0),
                'age_established': income.age_established,
                'is_contributing': income.is_contributing,
                'employer_match': income.employer_match,
                'age_last_contribution': income.age_last_contribution
            }
            
            # Add specific fields for different income types
            if income.income_type == 'social_security':
                income_dict['amount_at_fra'] = float(income.monthly_amount or 0)
            elif income.income_type == 'Life_Insurance':
                income_dict['loan_amount'] = float(income.monthly_amount or 0)
            elif income.income_type == 'Annuity':
                income_dict['percent_taxable'] = income.exclusion_ratio * 100
            
            income_data.append(income_dict)
        
        # Check if this is for duplication (add "Copy") or keep original name for editing
        is_for_duplication = request.GET.get('mode') == 'duplicate'
        scenario_name = f"{scenario.name} (Copy)" if is_for_duplication else scenario.name
        
        scenario_data = {
            'name': scenario_name,
            'description': scenario.description,
            'primary_retirement_age': scenario.retirement_age,
            'primary_medicare_age': scenario.medicare_age,
            'primary_lifespan': scenario.mortality_age,
            'spouse_retirement_age': scenario.spouse_retirement_age,
            'spouse_medicare_age': scenario.spouse_medicare_age,
            'spouse_lifespan': scenario.spouse_mortality_age,
            'model_tax_change': '',  # Reset to default
            'reduction_2030_ss': scenario.reduction_2030_ss,
            'ss_adjustment_year': scenario.ss_adjustment_year,
            'ss_adjustment_direction': scenario.ss_adjustment_direction,
            'ss_adjustment_type': scenario.ss_adjustment_type,
            'ss_adjustment_amount': scenario.ss_adjustment_amount,
            'apply_standard_deduction': scenario.apply_standard_deduction,
            'federal_standard_deduction': float(scenario.federal_standard_deduction or 0),
            'state_standard_deduction': float(scenario.state_standard_deduction or 0),
            'custom_annual_deduction': float(scenario.custom_annual_deduction or 0),
            'primary_blind': scenario.primary_blind,
            'spouse_blind': scenario.spouse_blind,
            'is_dependent': scenario.is_dependent,
            'part_b_inflation_rate': str(int(scenario.part_b_inflation_rate)),
            'part_d_inflation_rate': str(int(scenario.part_d_inflation_rate)),
            'primary_state': scenario.primary_state or '',
            'income': income_data
        }
        
        return Response(scenario_data, status=status.HTTP_200_OK)
        
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_scenario_assets(request, scenario_id):
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        assets = scenario.income_sources.all()
        asset_details = [
            {
                'id': asset.id,
                'owned_by': asset.owned_by,
                'income_type': asset.income_type,
                'income_name': asset.income_name,
                'current_asset_balance': asset.current_asset_balance,
                'monthly_amount': asset.monthly_amount,
                'monthly_contribution': asset.monthly_contribution,
                'age_to_begin_withdrawal': asset.age_to_begin_withdrawal,
                'age_to_end_withdrawal': asset.age_to_end_withdrawal,
                'rate_of_return': asset.rate_of_return,
                'cola': asset.cola,
                'exclusion_ratio': asset.exclusion_ratio,
                'tax_rate': asset.tax_rate,
                'scenario_id': asset.scenario_id
            }
            for asset in assets
        ]
        return Response(asset_details)
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=404)

@api_view(['PUT'])
@permission_classes([IsAuthenticated])
def update_scenario(request, scenario_id):
    """
    Update a scenario with Roth conversion parameters and income sources
    """
    print(f"💰 INCOME_EDIT: *** UPDATE_SCENARIO FUNCTION CALLED *** scenario_id={scenario_id}")
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Extract income sources from request data
        income_sources_data = request.data.pop('income_sources', None)
        
        # DEBUG: Log the income sources data
        print(f"💰 INCOME_EDIT: Received {len(income_sources_data) if income_sources_data else 0} income sources")
        if income_sources_data:
            for i, income_data in enumerate(income_sources_data):
                print(f"💰 INCOME_EDIT: Source {i}: Type={income_data.get('income_type')}, Name={income_data.get('income_name')}")
                print(f"💰 INCOME_EDIT: Source {i}: monthly_amount={income_data.get('monthly_amount')}, amount_at_fra={income_data.get('amount_at_fra')}")
                print(f"💰 INCOME_EDIT: Source {i}: Raw data: {income_data}")
        
        # Update scenario fields
        serializer = ScenarioUpdateSerializer(scenario, data=request.data, partial=True)
        if serializer.is_valid():
            serializer.save()
            
            # Update income sources if provided
            if income_sources_data is not None:
                # Delete existing income sources for this scenario
                scenario.income_sources.all().delete()
                
                # Create new income sources
                for income_data in income_sources_data:
                    # Remove any frontend-specific fields
                    income_data.pop('id', None)  # Remove frontend ID
                    
                    # Map frontend field names to backend model field names
                    # Special handling for social security - prioritize amount_at_fra
                    if income_data.get('income_type') == 'social_security':
                        mapped_monthly_amount = income_data.get('amount_at_fra') or income_data.get('monthly_amount') or income_data.get('withdrawal_amount', 0)
                    else:
                        mapped_monthly_amount = income_data.get('monthly_amount') or income_data.get('withdrawal_amount') or income_data.get('amount_at_fra', 0)
                    
                    print(f"💰 INCOME_EDIT: Mapping {income_data.get('income_type', 'unknown')} - mapped_monthly_amount={mapped_monthly_amount}")
                    print(f"💰 INCOME_EDIT: Source values - monthly_amount={income_data.get('monthly_amount')}, withdrawal_amount={income_data.get('withdrawal_amount')}, amount_at_fra={income_data.get('amount_at_fra')}")
                    
                    mapped_data = {
                        'owned_by': income_data.get('owned_by'),
                        'income_type': income_data.get('income_type'),
                        'income_name': income_data.get('income_name', ''),
                        'current_asset_balance': income_data.get('current_asset_balance') or income_data.get('current_balance', 0),
                        'monthly_amount': mapped_monthly_amount,
                        'monthly_contribution': income_data.get('monthly_contribution', 0),
                        'age_to_begin_withdrawal': income_data.get('age_to_begin_withdrawal') or income_data.get('start_age'),
                        'age_to_end_withdrawal': income_data.get('age_to_end_withdrawal') or income_data.get('end_age'),
                        'rate_of_return': income_data.get('rate_of_return') or (float(income_data.get('growth_rate', 0)) / 100 if income_data.get('growth_rate') else 0),
                        'cola': income_data.get('cola', 0),
                        'exclusion_ratio': income_data.get('exclusion_ratio', 0),
                        'tax_rate': income_data.get('tax_rate', 0),
                        'max_to_convert': income_data.get('max_to_convert', 0),
                        'age_established': income_data.get('age_established'),
                        'is_contributing': income_data.get('is_contributing', False),
                        'employer_match': income_data.get('employer_match', 0),
                        'age_last_contribution': income_data.get('age_last_contribution')
                    }
                    
                    # Remove None values to use model defaults
                    mapped_data = {k: v for k, v in mapped_data.items() if v is not None}
                    
                    IncomeSource.objects.create(scenario=scenario, **mapped_data)
            
            # Return updated scenario with income sources
            return Response({
                'id': scenario.id,
                'message': 'Scenario updated successfully'
            })
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)

@api_view(['DELETE'])
@permission_classes([IsAuthenticated])
def delete_scenario(request, client_id, scenario_id):
    """
    Delete a scenario
    """
    try:
        # Get the scenario and verify ownership
        scenario = Scenario.objects.get(id=scenario_id, client_id=client_id)
        
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Delete the scenario
        scenario.delete()
        
        return Response({"message": "Scenario deleted successfully"}, status=status.HTTP_204_NO_CONTENT)
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)

@api_view(['PUT'])
@permission_classes([IsAuthenticated])
def update_scenario_percentages(request, scenario_id):
    """
    Update the income_vs_cost_percent and medicare_irmaa_percent fields for a scenario.
    """
    try:
        scenario = get_object_or_404(Scenario, pk=scenario_id)
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Extract percentages from request data
        income_vs_cost_percent = request.data.get('income_vs_cost_percent')
        medicare_irmaa_percent = request.data.get('medicare_irmaa_percent')
        
        # Update only if the fields are provided
        if income_vs_cost_percent is not None:
            scenario.income_vs_cost_percent = income_vs_cost_percent
        
        if medicare_irmaa_percent is not None:
            scenario.medicare_irmaa_percent = medicare_irmaa_percent
        
        scenario.save()
        
        return Response({
            "id": scenario.id,
            "income_vs_cost_percent": scenario.income_vs_cost_percent,
            "medicare_irmaa_percent": scenario.medicare_irmaa_percent
        }, status=status.HTTP_200_OK)
        
    except Exception as e:
        return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['PUT'])
@permission_classes([IsAuthenticated])
def update_scenario_assets(request, scenario_id):
    """
    Update assets associated with a scenario, particularly for Roth conversion
    """
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Update assets for the scenario
        if 'assets' not in request.data:
            return Response(
                {"error": "Assets data is required"},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        assets_data = request.data['assets']
        updated_assets = []
        
        for asset_data in assets_data:
            asset_id = asset_data.get('id')
            if not asset_id:
                continue
                
            try:
                asset = IncomeSource.objects.get(id=asset_id, scenario=scenario)
                
                # Only update max_to_convert field
                if 'max_to_convert' in asset_data:
                    asset.max_to_convert = asset_data['max_to_convert']
                    asset.save()
                    updated_assets.append({
                        'id': asset.id,
                        'income_type': asset.income_type,
                        'max_to_convert': asset.max_to_convert
                    })
            except IncomeSource.DoesNotExist:
                # Skip assets that don't exist or don't belong to this scenario
                continue
        
        return Response({
            'message': f'Updated {len(updated_assets)} assets',
            'assets': updated_assets
        })
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)

class RothConversionAPIView(APIView):
    """
    API endpoint for processing Roth conversion scenarios.
    
    This endpoint accepts scenario data, client data, and Roth conversion parameters,
    and returns a detailed analysis of the Roth conversion impact, including:
    - Baseline scenario (without Roth conversion)
    - Conversion scenario (with Roth conversion)
    - Comparison of key metrics (taxes, Medicare costs, RMDs, inheritance tax)
    - Asset balance projections for visualization
    """
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]
    
    def _format_comparison_metrics(self, comparison):
        """
        Format comparison metrics to match the expected frontend structure.
        
        The frontend expects metrics like:
        {
            'rmd_reduction': 100000.0,
            'rmd_reduction_pct': 10.0,
            'tax_savings': 50000.0,
            'tax_savings_pct': 5.0,
            ...
        }
        """
        formatted = {}
        
        # Format RMD reduction
        if 'total_rmds' in comparison:
            rmd_diff = comparison['total_rmds']['difference']
            formatted['rmd_reduction'] = -rmd_diff  # Negate to get reduction
            formatted['rmd_reduction_pct'] = -comparison['total_rmds']['percent_change']
        
        # Format tax savings
        if 'lifetime_tax' in comparison:
            tax_diff = comparison['lifetime_tax']['difference']
            formatted['tax_savings'] = -tax_diff  # Negate to get savings
            formatted['tax_savings_pct'] = -comparison['lifetime_tax']['percent_change']
        
        # Format Medicare savings
        if 'lifetime_medicare' in comparison:
            medicare_diff = comparison['lifetime_medicare']['difference']
            formatted['medicare_savings'] = -medicare_diff
            formatted['medicare_savings_pct'] = -comparison['lifetime_medicare']['percent_change']
        
        # Format IRMAA savings
        if 'total_irmaa' in comparison:
            irmaa_diff = comparison['total_irmaa']['difference']
            formatted['irmaa_savings'] = -irmaa_diff
            formatted['irmaa_savings_pct'] = -comparison['total_irmaa']['percent_change']
        
        # Format inheritance tax savings
        if 'inheritance_tax' in comparison:
            inheritance_diff = comparison['inheritance_tax']['difference']
            formatted['inheritance_tax_savings'] = -inheritance_diff
            formatted['inheritance_tax_savings_pct'] = -comparison['inheritance_tax']['percent_change']
        
        # Format net income increase
        if 'cumulative_net_income' in comparison:
            income_diff = comparison['cumulative_net_income']['difference']
            formatted['net_income_increase'] = income_diff
            formatted['net_income_increase_pct'] = comparison['cumulative_net_income']['percent_change']
        
        # Format Roth increase
        if 'final_roth' in comparison:
            roth_diff = comparison['final_roth']['difference']
            formatted['roth_increase'] = roth_diff
            formatted['roth_increase_pct'] = comparison['final_roth']['percent_change']
        
        # Format total savings
        if 'total_expenses' in comparison:
            total_diff = comparison['total_expenses']['difference']
            formatted['total_savings'] = -total_diff
            formatted['total_savings_pct'] = -comparison['total_expenses']['percent_change']
        
        return formatted
    
    def post(self, request):
        try:
            # Extract data from request
            scenario = request.data.get('scenario')
            client = request.data.get('client')
            spouse = request.data.get('spouse')
            assets = request.data.get('assets')
            conversion_params = request.data.get('optimizer_params')
            
            # Log received data
            print("Received data for Roth conversion:")
            print(f"Scenario: {scenario}")
            print(f"Client: {client}")
            print(f"Spouse: {spouse}")
            print(f"Assets count: {len(assets) if assets else 0}")
            print(f"Conversion params: {conversion_params}")
            
            # Validate required fields
            if not scenario or not client or assets is None or not conversion_params:
                missing = []
                if not scenario: missing.append('scenario')
                if not client: missing.append('client')
                if assets is None: missing.append('assets')
                if not conversion_params: missing.append('optimizer_params')
                error_msg = f"Missing required fields: {', '.join(missing)}"
                print(f"ERROR: {error_msg}")
                return Response({'error': error_msg}, status=status.HTTP_400_BAD_REQUEST)
            
            # Create and process the Roth conversion
            processor = RothConversionProcessor(
                scenario=scenario,
                client=client,
                spouse=spouse,
                assets=assets,
                conversion_params=conversion_params
            )
            
            # Process the conversion and get results
            result = processor.process()
            
            # Transform the result to match the expected frontend structure
            transformed_result = {
                'baseline': {
                    'metrics': result['metrics']['baseline'],
                    'year_by_year': result['baseline_results']
                },
                'conversion': {
                    'metrics': result['metrics']['conversion'],
                    'year_by_year': result['conversion_results']
                },
                'comparison': self._format_comparison_metrics(result['metrics']['comparison']),
                'optimal_schedule': {
                    'start_year': result['conversion_params']['conversion_start_year'],
                    'duration': result['conversion_params']['years_to_convert'],
                    'annual_amount': result['conversion_params']['annual_conversion'],
                    'total_amount': result['conversion_params']['total_conversion'],
                    'score_breakdown': result['metrics']['conversion']
                },
                'asset_balances': result['asset_balances'],
                'scenarioResults': result['conversion_results']  # Add scenarioResults for frontend
            }
            
            # Return the transformed results
            return Response(transformed_result, status=status.HTTP_200_OK)
            
        except Exception as e:
            print(f"ERROR in Roth conversion processing: {str(e)}")
            import traceback
            traceback.print_exc()
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['POST'])
@permission_classes([AllowAny])
def register_advisor(request):
    """
    Register a new advisor (Step 1 and 2)
    """
    serializer = AdvisorRegistrationSerializer(data=request.data)
    
    if serializer.is_valid():
        try:
            # Create the user
            user = serializer.save()
            
            # Create Stripe customer
            stripe_customer = stripe.Customer.create(
                email=user.email,
                name=f"{user.first_name} {user.last_name}",
                metadata={
                    'user_id': user.id,
                    'company_name': user.company_name
                }
            )
            
            # Save Stripe customer ID to user model
            user.stripe_customer_id = stripe_customer.id
            user.save()

            # Generate tokens for automatic login
            refresh = RefreshToken.for_user(user)
            
            return Response({
                'status': 'success',
                'message': 'Advisor registered successfully',
                'tokens': {
                    'refresh': str(refresh),
                    'access': str(refresh.access_token),
                },
                'user': {
                    'id': user.id,
                    'email': user.email,
                    'first_name': user.first_name,
                    'last_name': user.last_name,
                }
            }, status=status.HTTP_201_CREATED)
            
        except Exception as e:
            # If there's an error, delete the user if it was created
            if 'user' in locals():
                user.delete()
            return Response({
                'status': 'error',
                'message': str(e)
            }, status=status.HTTP_400_BAD_REQUEST)
    
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def complete_registration(request):
    """
    Complete advisor registration with payment (Step 3)
    """
    try:
        user = request.user
        payment_method_id = request.data.get('paymentMethodId')
        plan = request.data.get('plan')

        if not payment_method_id or not plan:
            return Response({
                'status': 'error',
                'message': 'Payment method and plan are required'
            }, status=status.HTTP_400_BAD_REQUEST)

        # Attach payment method to customer
        stripe.PaymentMethod.attach(
            payment_method_id,
            customer=user.stripe_customer_id,
        )

        # Set as default payment method
        stripe.Customer.modify(
            user.stripe_customer_id,
            invoice_settings={
                'default_payment_method': payment_method_id
            }
        )

        # Get the price ID based on the plan
        price_id = settings.STRIPE_MONTHLY_PRICE_ID if plan == 'monthly' else settings.STRIPE_ANNUAL_PRICE_ID

        # Create the subscription
        subscription = stripe.Subscription.create(
            customer=user.stripe_customer_id,
            items=[{'price': price_id}],
            payment_behavior='default_incomplete',
            expand=['latest_invoice.payment_intent'],
            metadata={
                'user_id': user.id,
                'plan': plan
            }
        )

        # Save subscription ID to user model
        user.stripe_subscription_id = subscription.id
        user.subscription_status = subscription.status
        user.save()

        return Response({
            'status': 'success',
            'subscription': {
                'id': subscription.id,
                'status': subscription.status,
                'client_secret': subscription.latest_invoice.payment_intent.client_secret
            }
        })

    except stripe.error.StripeError as e:
        return Response({
            'status': 'error',
            'message': str(e)
        }, status=status.HTTP_400_BAD_REQUEST)
    except Exception as e:
        return Response({
            'status': 'error',
            'message': str(e)
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# Real Estate API Views
class ListCreateRealEstateView(generics.ListCreateAPIView):
    serializer_class = RealEstateSerializer
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]

    def get_queryset(self):
        client_id = self.kwargs['client_id']
        return RealEstate.objects.filter(client_id=client_id)

    def perform_create(self, serializer):
        client_id = self.kwargs['client_id']
        client = get_object_or_404(Client, id=client_id, advisor=self.request.user)
        serializer.save(client=client)

class RealEstateDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = RealEstate.objects.all()
    serializer_class = RealEstateSerializer
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]

class ReportTemplateViewSet(viewsets.ViewSet):
    permission_classes = [IsAuthenticated]
    
    def list(self, request, client_pk=None):
        """List all templates for a specific client"""
        client = get_object_or_404(Client, pk=client_pk, advisor=request.user)
        templates = ReportTemplate.objects.filter(client=client)
        serializer = ReportTemplateSerializer(templates, many=True, context={'request': request})
        return Response(serializer.data)
    
    def create(self, request, client_pk=None):
        """Upload a new PowerPoint template and process its slides"""
        try:
            # Check if python-pptx is installed
            from pptx import Presentation
        except ImportError:
            return Response(
                {"error": "The python-pptx library is not installed on the server."},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
        client = get_object_or_404(Client, pk=client_pk, advisor=request.user)
        
        # Validate request data
        if 'name' not in request.data:
            return Response({"error": "Template name is required"}, status=status.HTTP_400_BAD_REQUEST)
        
        if 'file' not in request.FILES:
            return Response({"error": "PowerPoint file is required"}, status=status.HTTP_400_BAD_REQUEST)
        
        pptx_file = request.FILES['file']
        
        # Check file extension
        if not pptx_file.name.endswith('.pptx'):
            return Response(
                {"error": "Only PowerPoint files (.pptx) are supported"},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Create the template record
        template = ReportTemplate.objects.create(
            client=client,
            name=request.data['name'],
            file=pptx_file
        )
        
        # Process the PowerPoint file to extract slide thumbnails
        try:
            self._process_slides(template, pptx_file)
        except Exception as e:
            # If slide processing fails, delete the template and return error
            template.delete()
            return Response(
                {"error": f"Failed to process PowerPoint file: {str(e)}"},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
        serializer = ReportTemplateSerializer(template, context={'request': request})
        return Response(serializer.data, status=status.HTTP_201_CREATED)
    
    def retrieve(self, request, pk=None):
        """Get details of a specific template including all slides"""
        template = get_object_or_404(ReportTemplate, pk=pk)
        
        # Check if the user has access to this template
        if template.client.advisor != request.user:
            return Response(
                {"error": "You don't have permission to access this template"},
                status=status.HTTP_403_FORBIDDEN
            )
        
        serializer = ReportTemplateDetailSerializer(template, context={'request': request})
        return Response(serializer.data)
    
    def destroy(self, request, pk=None):
        """Delete a template and all its slides"""
        template = get_object_or_404(ReportTemplate, pk=pk)
        
        # Check if the user has access to this template
        if template.client.advisor != request.user:
            return Response(
                {"error": "You don't have permission to delete this template"},
                status=status.HTTP_403_FORBIDDEN
            )
        
        template.delete()
        return Response(status=status.HTTP_204_NO_CONTENT)
    
    @action(detail=True, methods=['get'])
    def slides(self, request, pk=None):
        """Get all slides for a specific template"""
        template = get_object_or_404(ReportTemplate, pk=pk)
        
        # Check if the user has access to this template
        if template.client.advisor != request.user:
            return Response(
                {"error": "You don't have permission to access this template"},
                status=status.HTTP_403_FORBIDDEN
            )
        
        slides = template.slides.all().order_by('order')
        serializer = TemplateSlideSerializer(slides, many=True, context={'request': request})
        
        return Response({"slides": serializer.data})
    
    @action(detail=True, methods=['post'])
    def update_slides(self, request, pk=None):
        """Update slide order and remove slides"""
        template = get_object_or_404(ReportTemplate, pk=pk)
        
        # Check if the user has access to this template
        if template.client.advisor != request.user:
            return Response(
                {"error": "You don't have permission to modify this template"},
                status=status.HTTP_403_FORBIDDEN
            )
        
        if 'slides' not in request.data:
            return Response(
                {"error": "Slides data is required"},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Get all current slides
        current_slides = {str(slide.id): slide for slide in template.slides.all()}
        
        # Process the update
        for idx, slide_data in enumerate(request.data['slides']):
            slide_id = slide_data.get('id')
            
            if not slide_id or slide_id not in current_slides:
                continue
                
            # Update the order
            slide = current_slides[slide_id]
            slide.order = idx
            slide.save()
            
            # Remove from the current_slides dict to track which ones to delete
            current_slides.pop(slide_id)
        
        # Delete any slides not included in the request
        for slide in current_slides.values():
            slide.delete()
        
        # Return the updated slides
        slides = template.slides.all().order_by('order')
        serializer = TemplateSlideSerializer(slides, many=True, context={'request': request})
        
        return Response({"slides": serializer.data})
    
    def _process_slides(self, template, pptx_file):
        """Process a PowerPoint file and extract slide thumbnails"""
        try:
            # Create a temporary file to save the uploaded file
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                for chunk in pptx_file.chunks():
                    temp_file.write(chunk)
                temp_file_path = temp_file.name
            
            # Create output directory for thumbnails
            thumbnails_dir = os.path.join(settings.MEDIA_ROOT, 'thumbnails')
            os.makedirs(thumbnails_dir, exist_ok=True)
            
            # Initialize png_paths to empty list
            png_paths = []
            
            try:
                # Convert PowerPoint to PNG using our new converter
                png_paths = convert_pptx_to_png(temp_file_path, thumbnails_dir, dpi=300)
                
                # Save each thumbnail as a slide
                for i, png_path in enumerate(png_paths):
                    # Open the PNG file
                    with open(png_path, 'rb') as f:
                        image_data = f.read()
                    
                    # Create the slide record
                    slide_filename = f"slide_{i+1}.png"
                    template_slide = TemplateSlide(
                        template=template,
                        order=i
                    )
                    template_slide.thumbnail.save(slide_filename, BytesIO(image_data), save=True)
            finally:
                # Clean up temporary files
                try:
                    os.unlink(temp_file_path)
                except:
                    pass
                
                for png_path in png_paths:
                    try:
                        os.unlink(png_path)
                    except:
                        pass
        
        except Exception as e:
            import logging
            logger = logging.getLogger(__name__)
            logger.error(f"Error processing PowerPoint file: {str(e)}")
            raise e

@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated])
def client_report_templates(request, client_id):
    """
    GET: List all report templates for a client
    POST: Create a new report template for a client
    """
    client = get_object_or_404(Client, pk=client_id, advisor=request.user)
    
    if request.method == 'GET':
        templates = ReportTemplate.objects.filter(client=client)
        serializer = ReportTemplateSerializer(templates, many=True, context={'request': request})
        return Response(serializer.data)
    
    elif request.method == 'POST':
        try:
            # Check if python-pptx is installed
            from pptx import Presentation
        except ImportError:
            return Response(
                {"error": "The python-pptx library is not installed on the server."},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
        # Validate request data
        if 'name' not in request.data:
            return Response({"error": "Template name is required"}, status=status.HTTP_400_BAD_REQUEST)
        
        if 'file' not in request.FILES:
            return Response({"error": "PowerPoint file is required"}, status=status.HTTP_400_BAD_REQUEST)
        
        pptx_file = request.FILES['file']
        
        # Check file extension
        if not pptx_file.name.endswith('.pptx'):
            return Response(
                {"error": "Only PowerPoint files (.pptx) are supported"},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Create the template record
        template = ReportTemplate.objects.create(
            client=client,
            name=request.data['name'],
            file=pptx_file
        )
        
        # Process the PowerPoint file to extract slide thumbnails
        try:
            process_pptx_slides(template, pptx_file)
        except Exception as e:
            # If slide processing fails, delete the template and return error
            template.delete()
            return Response(
                {"error": f"Failed to process PowerPoint file: {str(e)}"},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
        serializer = ReportTemplateSerializer(template, context={'request': request})
        return Response(serializer.data, status=status.HTTP_201_CREATED)

def process_pptx_slides(template, pptx_file):
    """Process a PowerPoint file and extract slide thumbnails"""
    try:
        # Create a temporary file to save the uploaded file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            for chunk in pptx_file.chunks():
                temp_file.write(chunk)
            temp_file_path = temp_file.name
        
        # Create output directory for thumbnails
        thumbnails_dir = os.path.join(settings.MEDIA_ROOT, 'thumbnails')
        os.makedirs(thumbnails_dir, exist_ok=True)
        
        # Initialize png_paths to empty list
        png_paths = []
        
        try:
            # Convert PowerPoint to PNG using our new converter
            png_paths = convert_pptx_to_png(temp_file_path, thumbnails_dir, dpi=300)
            
            # Save each thumbnail as a slide
            for i, png_path in enumerate(png_paths):
                # Open the PNG file
                with open(png_path, 'rb') as f:
                    image_data = f.read()
                
                # Create the slide record
                slide_filename = f"slide_{i+1}.png"
                template_slide = TemplateSlide(
                    template=template,
                    order=i
                )
                template_slide.thumbnail.save(slide_filename, BytesIO(image_data), save=True)
        finally:
            # Clean up temporary files
            try:
                os.unlink(temp_file_path)
            except:
                pass
            
            for png_path in png_paths:
                try:
                    os.unlink(png_path)
                except:
                    pass
    
    except Exception as e:
        import logging
        logger = logging.getLogger(__name__)
        logger.error(f"Error processing PowerPoint file: {str(e)}")
        raise e

@api_view(['GET', 'DELETE'])
@permission_classes([IsAuthenticated])
def report_template_detail(request, template_id):
    """
    GET: Retrieve details of a specific template
    DELETE: Delete a template
    """
    template = get_object_or_404(ReportTemplate, pk=template_id)
    
    # Check if the user has access to this template
    if template.client.advisor != request.user:
        return Response(
            {"error": "You don't have permission to access this template"},
            status=status.HTTP_403_FORBIDDEN
        )
    
    if request.method == 'GET':
        serializer = ReportTemplateDetailSerializer(template, context={'request': request})
        return Response(serializer.data)
    
    elif request.method == 'DELETE':
        template.delete()
        return Response(status=status.HTTP_204_NO_CONTENT)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def template_slides(request, template_id):
    """Get all slides for a specific template"""
    template = get_object_or_404(ReportTemplate, pk=template_id)
    
    # Check if the user has access to this template
    if template.client.advisor != request.user:
        return Response(
            {"error": "You don't have permission to access this template"},
            status=status.HTTP_403_FORBIDDEN
        )
    
    slides = template.slides.all().order_by('order')
    serializer = TemplateSlideSerializer(slides, many=True, context={'request': request})
    
    return Response({"slides": serializer.data})

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def update_template_slides(request, template_id):
    """Update slide order and remove slides"""
    template = get_object_or_404(ReportTemplate, pk=template_id)
    
    # Check if the user has access to this template
    if template.client.advisor != request.user:
        return Response(
            {"error": "You don't have permission to modify this template"},
            status=status.HTTP_403_FORBIDDEN
        )
    
    if 'slides' not in request.data:
        return Response(
            {"error": "Slides data is required"},
            status=status.HTTP_400_BAD_REQUEST
        )
    
    # Get all current slides
    current_slides = {str(slide.id): slide for slide in template.slides.all()}
    
    # Process the update
    for idx, slide_data in enumerate(request.data['slides']):
        slide_id = slide_data.get('id')
        
        if not slide_id or slide_id not in current_slides:
            continue
            
        # Update the order
        slide = current_slides[slide_id]
        slide.order = idx
        slide.save()
        
        # Remove from the current_slides dict to track which ones to delete
        current_slides.pop(slide_id)
    
    # Delete any slides not included in the request
    for slide in current_slides.values():
        slide.delete()
    
    # Return the updated slides
    slides = template.slides.all().order_by('order')
    serializer = TemplateSlideSerializer(slides, many=True, context={'request': request})
    
    return Response({"slides": serializer.data})

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_scenario_comparison_data(request, scenario_id):
    """Get scenario data formatted for comparison report"""
    logger = logging.getLogger(__name__)
    
    logger.info(f"COMPARISON ENDPOINT CALLED: scenario_id={scenario_id}, user={request.user}")
    
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Run calculation to get the latest data
        try:
            processor = ScenarioProcessor(scenario_id=scenario.id, debug=False)
            calculation_result = processor.calculate()
            logger.info(f"Calculation result type: {type(calculation_result)}")
            logger.info(f"Calculation result: {str(calculation_result)[:500]}...")  # Log first 500 chars
        except Exception as calc_error:
            logger.error(f"Error calculating scenario {scenario_id}: {str(calc_error)}")
            # Return basic data without calculation
            comparison_data = {
                'id': scenario.id,
                'name': scenario.name,
                'irmaa_reached': False,
                'medicare_cost': 0,
                'federal_taxes': 0,
                'solution_cost': 0,
                'total_costs': 0,
                'out_of_pocket': 0
            }
            return Response(comparison_data, status=status.HTTP_200_OK)
        
        # Handle both list and dict return types from calculate()
        summary = {}
        if isinstance(calculation_result, list):
            logger.info(f"Calculation result is list with {len(calculation_result)} items")
            # If it's a list, look for the Summary in the last element or iterate
            for i, item in enumerate(calculation_result):
                logger.info(f"Item {i}: type={type(item)}, keys={list(item.keys()) if isinstance(item, dict) else 'N/A'}")
                if isinstance(item, dict) and 'Summary' in item:
                    summary = item['Summary']
                    logger.info(f"Found Summary in item {i}: {summary}")
                    break
            # If no Summary found in list items, check if last item is Summary
            if not summary and calculation_result:
                last_item = calculation_result[-1]
                if isinstance(last_item, dict):
                    summary = last_item
                    logger.info(f"Using last item as summary: {summary}")
        else:
            logger.info(f"Calculation result is dict with keys: {list(calculation_result.keys()) if isinstance(calculation_result, dict) else 'N/A'}")
            # If it's a dict, extract Summary as before
            summary = calculation_result.get('Summary', {})
            logger.info(f"Extracted summary: {summary}")
        
        logger.info(f"Final summary data: {summary}")
        
        # Calculate totals from the actual calculation result data
        total_federal_taxes = 0
        total_medicare_costs = 0
        total_gross_income = 0
        irmaa_reached = False
        
        # The calculation result is a list of yearly data
        if isinstance(calculation_result, list):
            for year_data in calculation_result:
                if isinstance(year_data, dict):
                    # Extract values and convert Decimal to float
                    fed_tax_raw = year_data.get('federal_tax', 0)
                    medicare_raw = year_data.get('total_medicare', 0)
                    gross_raw = year_data.get('gross_income', 0)
                    irmaa_raw = year_data.get('irmaa_surcharge', 0)
                    
                    # Convert to float, handling Decimal objects
                    fed_tax_float = float(fed_tax_raw) if fed_tax_raw is not None else 0
                    medicare_float = float(medicare_raw) if medicare_raw is not None else 0
                    gross_float = float(gross_raw) if gross_raw is not None else 0
                    irmaa_float = float(irmaa_raw) if irmaa_raw is not None else 0
                    
                    # Add to totals
                    total_federal_taxes += fed_tax_float
                    total_medicare_costs += medicare_float
                    total_gross_income += gross_float
                    
                    # Check if IRMAA is reached (any year with IRMAA surcharge > 0)
                    if irmaa_float > 0:
                        irmaa_reached = True
        
        # Calculate derived values
        total_costs = total_federal_taxes + total_medicare_costs
        out_of_pocket = total_gross_income - total_costs  # Net income
        solution_cost = 0  # This might need to be calculated differently based on business logic
        
        comparison_data = {
            'id': scenario.id,
            'name': scenario.name,
            'irmaa_reached': irmaa_reached,
            'medicare_cost': total_medicare_costs,
            'federal_taxes': total_federal_taxes,
            'solution_cost': solution_cost,
            'total_costs': total_costs,
            'out_of_pocket': out_of_pocket
        }
        
        logger.info(f"Final comparison data: {comparison_data}")
        
        return Response(comparison_data, status=status.HTTP_200_OK)
        
    except Scenario.DoesNotExist:
        return Response({"error": "Scenario not found"}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        logger.error(f"Error in get_scenario_comparison_data: {str(e)}")
        return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated])
def comparison_preferences(request, client_id):
    """Get or update comparison preferences for a client"""
    from .models import ComparisonPreference, Client, Scenario
    
    try:
        client = Client.objects.get(id=client_id, advisor=request.user)
    except Client.DoesNotExist:
        return Response({"error": "Client not found or access denied"}, status=status.HTTP_404_NOT_FOUND)
    
    if request.method == 'GET':
        # Get existing preferences
        try:
            preference = ComparisonPreference.objects.get(user=request.user, client=client)
            return Response({
                'scenario1': preference.scenario1.id if preference.scenario1 else None,
                'scenario2': preference.scenario2.id if preference.scenario2 else None
            })
        except ComparisonPreference.DoesNotExist:
            return Response({
                'scenario1': None,
                'scenario2': None
            })
    
    elif request.method == 'POST':
        # Save or update preferences
        scenario1_id = request.data.get('scenario1')
        scenario2_id = request.data.get('scenario2')
        
        # Validate scenarios belong to the client
        scenario1 = None
        scenario2 = None
        
        if scenario1_id:
            try:
                scenario1 = Scenario.objects.get(id=scenario1_id, client=client)
            except Scenario.DoesNotExist:
                return Response({"error": "Scenario 1 not found"}, status=status.HTTP_400_BAD_REQUEST)
        
        if scenario2_id:
            try:
                scenario2 = Scenario.objects.get(id=scenario2_id, client=client)
            except Scenario.DoesNotExist:
                return Response({"error": "Scenario 2 not found"}, status=status.HTTP_400_BAD_REQUEST)
        
        # Update or create preference
        preference, created = ComparisonPreference.objects.update_or_create(
            user=request.user,
            client=client,
            defaults={
                'scenario1': scenario1,
                'scenario2': scenario2
            }
        )
        
        return Response({
            'message': 'Preferences saved successfully',
            'scenario1': preference.scenario1.id if preference.scenario1 else None,
            'scenario2': preference.scenario2.id if preference.scenario2 else None
        }, status=status.HTTP_200_OK)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_federal_standard_deduction(request):
    """Get federal standard deduction amount based on filing status"""
    from .tax_csv_loader import get_tax_loader
    
    filing_status = request.GET.get('filing_status', 'Single')
    
    try:
        loader = get_tax_loader()
        deduction_amount = loader.get_standard_deduction(filing_status)
        
        return Response({
            'filing_status': filing_status,
            'deduction_amount': float(deduction_amount)
        })
    except Exception as e:
        return Response({
            'error': f'Error retrieving standard deduction: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_irmaa_thresholds_for_years(request):
    """Get IRMAA thresholds inflated for each year in a scenario"""
    from .tax_csv_loader import get_tax_loader
    
    filing_status = request.GET.get('filing_status', 'Single')
    start_year = int(request.GET.get('start_year', 2025))
    end_year = int(request.GET.get('end_year', 2045))
    
    try:
        loader = get_tax_loader()
        
        # Get thresholds for each year
        year_thresholds = {}
        for year in range(start_year, end_year + 1):
            thresholds = loader.get_inflated_irmaa_thresholds(filing_status, year)
            year_thresholds[year] = [
                {
                    'magi_threshold': float(t['magi_threshold']),
                    'part_b_surcharge': float(t['part_b_surcharge']),
                    'part_d_surcharge': float(t['part_d_surcharge']),
                    'description': t['description']
                }
                for t in thresholds if t['magi_threshold'] > 0
            ]
        
        return Response({
            'filing_status': filing_status,
            'start_year': start_year,
            'end_year': end_year,
            'thresholds_by_year': year_thresholds
        })
    except Exception as e:
        return Response({
            'error': f'Error retrieving IRMAA thresholds: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

