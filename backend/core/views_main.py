import requests
import os
import sys
import tempfile
import shutil
from django.contrib.auth import authenticate
from rest_framework import status, generics, permissions
from rest_framework.decorators import api_view, permission_classes, throttle_classes, action
from rest_framework.permissions import AllowAny, IsAuthenticated
from rest_framework.response import Response
from rest_framework.views import APIView
from rest_framework_simplejwt.authentication import JWTAuthentication
from rest_framework_simplejwt.tokens import RefreshToken
from .throttles import LoginRateThrottle
from .serializers_main import (
    CustomUserSerializer, UserSerializer, ClientSerializer,
    EmailAccountSerializer, CommunicationSerializer, CommunicationCreateSerializer,
    LeadSerializer, LeadCreateSerializer, LeadSourceSerializer, SMSMessageSerializer,
    TwilioConfigurationSerializer, ActivityLogSerializer, EmailComposeSerializer,
    TaskSerializer, TaskCreateSerializer, TaskDetailSerializer, TaskTemplateSerializer, TaskCommentSerializer,
    CalendarAccountSerializer, CalendarEventSerializer, CalendarEventCreateSerializer, 
    MeetingTemplateSerializer, CalendarEventReminderSerializer, MeetingScheduleSerializer
)
from django.shortcuts import render, get_object_or_404
from django.contrib.auth import get_user_model
from django.views.decorators.csrf import csrf_exempt
from django.utils.decorators import method_decorator
from .models import (
    Scenario, Client, RealEstate, ReportTemplate, TemplateSlide, IncomeSource,
    EmailAccount, Communication, Lead, LeadSource, SMSMessage, TwilioConfiguration, ActivityLog,
    Task, TaskTemplate, TaskComment, CalendarAccount, CalendarEvent, MeetingTemplate, CalendarEventReminder
)
from rest_framework.exceptions import PermissionDenied
from .serializers import ClientDetailSerializer, ClientEditSerializer, ClientCreateSerializer, RealEstateSerializer
from .serializers import ScenarioCreateSerializer, ScenarioUpdateSerializer, IncomeSourceUpdateSerializer
from .serializers import ReportTemplateSerializer, ReportTemplateDetailSerializer, TemplateSlideSerializer
from .scenario_processor import ScenarioProcessor
from .roth_conversion_processor import RothConversionProcessor
from .tasks import calculate_scenario_async
from celery.result import AsyncResult
from django.http import HttpResponse, JsonResponse
from django.utils import timezone
from django.db import models
import logging
from .serializers import AdvisorRegistrationSerializer
import stripe
from django.conf import settings
from rest_framework import viewsets
import json
from PIL import Image, ImageDraw
from io import BytesIO
import uuid
from .pptx_utils import create_enhanced_thumbnails
from pptx_to_png import convert_pptx_to_png
from decimal import Decimal, InvalidOperation
import copy

# Add the python-pptx library for handling PowerPoint files
try:
    from pptx import Presentation
except ImportError:
    # We'll handle this in the view with a proper error response
    pass

User = get_user_model()

# Initialize Stripe
stripe.api_key = settings.STRIPE_SECRET_KEY

def home_view(request):
    return render(request, 'home.html')

# @csrf_exempt
@api_view(['GET', 'POST'])
@throttle_classes([LoginRateThrottle])
@permission_classes([AllowAny])
def login_view(request):
    if request.method == 'GET':
        return Response({"message": "Use POST to login."}, status=200)

    email = request.data.get('email')
    password = request.data.get('password')

    if not email or not password:
        return Response({'message': 'Email and password are required'}, status=400)

    user_obj = User.objects.filter(email=email).first()
    if user_obj is None:
        return Response({'message': 'Invalid credentials'}, status=401)

    user = authenticate(request, username=user_obj.email, password=password)
    if user is None:
        return Response({'message': 'Invalid credentials'}, status=401)

    refresh = RefreshToken.for_user(user)
    return Response({
        'access': str(refresh.access_token),
        'refresh': str(refresh),
        'user': {
            'id': user.id,
            'email': user.email,
            'username': user.username
        }
    })

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def logout_view(request):
    try:
        # Blacklist the refresh token (if used)
        refresh_token = request.data.get("refresh_token")
        token = RefreshToken(refresh_token)
        token.blacklist()
    except Exception as e:
        return Response({"error": str(e)}, status=status.HTTP_400_BAD_REQUEST)

    return Response({"detail": "Successfully logged out."}, status=status.HTTP_200_OK)

@api_view(['POST'])
@permission_classes([AllowAny])
def register_view(request):
    data = request.data
    if User.objects.filter(email=data.get('email')).exists():
        return Response({'message': 'Email already in use'}, status=status.HTTP_400_BAD_REQUEST)

    user = User.objects.create_user(
        username=data.get('email'),
        email=data.get('email'),
        password=data.get('password'),
        first_name=data.get('firstName'),
        last_name=data.get('lastName')
    )
    # return Response({'message': 'User created'}, status=status.HTTP_201_CREATED)
    #user = User.objects.create_user(username=email, email=email, password=password)

    # ✅ Generate JWT token
    refresh = RefreshToken.for_user(user)
    return Response({
        'access': str(refresh.access_token),
        'refresh': str(refresh),
        'user': {
            'id': user.id,
            'email': user.email,
            'username': user.username
        }
    })

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def run_scenario_calculation(request, scenario_id):
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        # New instantiation: only pass scenario_id and debug
        processor = ScenarioProcessor(scenario_id=scenario.id, debug=True)
        result = processor.calculate()
        return Response(result)
    except Scenario.DoesNotExist:
        return Response({"error": "Scenario not found."}, status=404)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def start_scenario_calculation_async(request, scenario_id):
    """
    Start async scenario calculation and return task ID for progress tracking
    """
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Start async task
        task = calculate_scenario_async.delay(scenario_id, request.user.id)
        
        return Response({
            "task_id": task.id,
            "status": "PENDING",
            "message": "Calculation started in background",
            "scenario_id": scenario_id
        }, status=202)
        
    except Scenario.DoesNotExist:
        return Response({"error": "Scenario not found."}, status=404)
    except Exception as e:
        return Response({"error": str(e)}, status=500)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_task_status(request, task_id):
    """
    Get the status and result of an async task
    """
    try:
        task = AsyncResult(task_id)
        
        if task.state == 'PENDING':
            response = {
                'task_id': task_id,
                'state': task.state,
                'status': 'PENDING',
                'message': 'Task is waiting to be processed',
                'progress': 0
            }
        elif task.state == 'PROGRESS':
            response = {
                'task_id': task_id,
                'state': task.state,
                'status': 'PROGRESS',
                'message': task.info.get('message', 'Processing...'),
                'progress': task.info.get('progress', 0),
                'current_step': task.info.get('current_step', ''),
                'current_year': task.info.get('current_year'),
                'processed_years': task.info.get('processed_years', 0),
                'total_years': task.info.get('total_years', 0)
            }
        elif task.state == 'SUCCESS':
            result = task.result
            response = {
                'task_id': task_id,
                'state': task.state,
                'status': 'SUCCESS',
                'message': result.get('message', 'Calculation completed successfully'),
                'progress': 100,
                'results': result.get('results', []),
                'summary': result.get('summary', {}),
                'calculation_time': result.get('calculation_time')
            }
        else:
            # FAILURE or other states
            response = {
                'task_id': task_id,
                'state': task.state,
                'status': 'FAILURE',
                'message': str(task.info) if task.info else 'Task failed with unknown error',
                'progress': 0,
                'error': str(task.info) if task.info else 'Unknown error'
            }
        
        return Response(response)
        
    except Exception as e:
        return Response({
            'task_id': task_id,
            'status': 'ERROR',
            'error': str(e)
        }, status=500)


@api_view(['DELETE'])
@permission_classes([IsAuthenticated])
def cancel_task(request, task_id):
    """
    Cancel a running async task
    """
    try:
        task = AsyncResult(task_id)
        task.revoke(terminate=True)
        
        return Response({
            'task_id': task_id,
            'status': 'CANCELLED',
            'message': 'Task cancelled successfully'
        })
        
    except Exception as e:
        return Response({
            'task_id': task_id,
            'status': 'ERROR',
            'error': str(e)
        }, status=500)
    
@api_view(['GET', 'PUT'])
@permission_classes([IsAuthenticated])
def profile_view(request):
    from .authentication import get_enhanced_user_data
    user = request.user
    if request.method == 'GET':
        # Return enhanced user data including admin fields
        return Response(get_enhanced_user_data(user))
    elif request.method == 'PUT':
        # Handle file uploads
        data = request.data.copy()
        
        # Handle empty string for website_url to prevent validation errors
        if 'website_url' in data and data['website_url'] == '':
            data['website_url'] = None
        
        # Create a serializer with the data and possible file
        serializer = UserSerializer(user, data=data, partial=True)
        if serializer.is_valid():
            # Handle logo file upload separately if present
            if 'logo' in request.FILES:
                user.logo = request.FILES['logo']
                user.save(update_fields=['logo'])
            
            # Save the rest of the data
            serializer.save()
            return Response(serializer.data, status=status.HTTP_200_OK)
        
        # Return detailed error messages
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
    
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def create_scenario(request, client_id):
    try:
        client = Client.objects.get(id=client_id, advisor=request.user)
    except Client.DoesNotExist:
        return Response({'error': 'Client not found or access denied.'}, status=404)

    data = request.data.copy()
    data['client'] = client.id

    serializer = ScenarioCreateSerializer(data=data, context={'request': request})
    if serializer.is_valid():
        scenario = serializer.save()
        return Response({'id': scenario.id}, status=201)
    else:
        return Response(serializer.errors, status=400)  

@csrf_exempt
def proxy_to_wealthbox(request):
    api_key = '020dfe66ac3f4213beefb36641a96721'  # Replace with your actual API key
    # Log the incoming request path
    print(f"Received request path: {request.path}")
    url = f"https://api.crmworkspace.com/{request.path.replace('/proxy', '')}"
    print(f"Proxying request to URL: {url}")
    headers = {'ACCESS_TOKEN': api_key}

    try:
        response = requests.get(url, headers=headers)
        # Log the response status code
        # self._log_debug(f"Response status code: {response.status_code}")
        return JsonResponse(response.json(), status=response.status_code)
    except requests.exceptions.RequestException as e:
        logging.error(f"RequestException: {str(e)}")
        return JsonResponse({'error': str(e)}, status=500)


class AdvisorClientListView(generics.ListAPIView):
    serializer_class = ClientSerializer
    authentication_classes = [JWTAuthentication]
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        queryset = Client.objects.filter(advisor=self.request.user)
        
        # Add search functionality
        search = self.request.query_params.get('search', None)
        if search:
            # Search in first name, last name, and email fields
            queryset = queryset.filter(
                models.Q(first_name__icontains=search) |
                models.Q(last_name__icontains=search) |
                models.Q(email__icontains=search)
            )
        
        # Add limit parameter for search results
        limit = self.request.query_params.get('limit', None)
        if limit:
            try:
                limit = int(limit)
                queryset = queryset[:limit]
            except ValueError:
                pass  # Ignore invalid limit values
                
        return queryset
    
  
class ClientDetailView(generics.RetrieveAPIView):
    queryset = Client.objects.all()
    serializer_class = ClientDetailSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_object(self):
        client = get_object_or_404(Client, pk=self.kwargs['pk'])
        print(f"Client advisor ID: {client.advisor.id}")
        print(f"Requesting user ID: {self.request.user.id}")
        if client.advisor != self.request.user:
            raise PermissionDenied("You do not have permission to view this client.")
        return client
    
class ClientEditView(generics.RetrieveUpdateAPIView):
    queryset = Client.objects.all()
    serializer_class = ClientEditSerializer
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]

    def get_queryset(self):
        user = self.request.user
        if user.is_staff:
            return Client.objects.all()
        return Client.objects.filter(advisor=user)

    def get_object(self):
        client = get_object_or_404(self.get_queryset(), pk=self.kwargs["pk"])
        return client
    
class ClientCreateView(APIView):
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]

    def post(self, request):
        serializer = ClientCreateSerializer(data=self.request.data, context={'request': self.request})
        if serializer.is_valid():
            validated_data = serializer.validated_data.copy()
            validated_data.pop('advisor', None)  # Remove advisor if it's already in validated_data
            client = serializer.save(advisor=self.request.user)
            return Response(ClientDetailSerializer(client).data, status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
    
class ScenarioCreateView(APIView):
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]

    def post(self, request, client_id):
        serializer = ScenarioCreateSerializer(data=request.data, context={'request': request, 'client_id': client_id})
        if serializer.is_valid():
            scenario = serializer.save()
            return Response({'id': scenario.id, 'message': 'Scenario created successfully'}, status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def duplicate_scenario(request, scenario_id):
    """
    Duplicate a scenario with all its income sources and settings
    """
    try:
        # Get the original scenario
        original_scenario = Scenario.objects.get(id=scenario_id)
        
        # Check if the user owns the client associated with this scenario
        if original_scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Create a copy of the scenario
        duplicated_scenario = Scenario.objects.create(
            client=original_scenario.client,
            name=f"{original_scenario.name} (Copy)",
            description=original_scenario.description,
            retirement_age=original_scenario.retirement_age,
            medicare_age=original_scenario.medicare_age,
            spouse_retirement_age=original_scenario.spouse_retirement_age,
            spouse_medicare_age=original_scenario.spouse_medicare_age,
            mortality_age=original_scenario.mortality_age,
            spouse_mortality_age=original_scenario.spouse_mortality_age,
            retirement_year=original_scenario.retirement_year,
            share_with_client=original_scenario.share_with_client,
            part_b_inflation_rate=original_scenario.part_b_inflation_rate,
            part_d_inflation_rate=original_scenario.part_d_inflation_rate,
            FRA_amount=original_scenario.FRA_amount,
            roth_conversion_start_year=original_scenario.roth_conversion_start_year,
            roth_conversion_duration=original_scenario.roth_conversion_duration,
            roth_conversion_annual_amount=original_scenario.roth_conversion_annual_amount,
            apply_standard_deduction=original_scenario.apply_standard_deduction,
            income_vs_cost_percent=0,  # Reset to 0 for new scenario
            medicare_irmaa_percent=0   # Reset to 0 for new scenario
        )
        
        # Duplicate all income sources
        for income_source in original_scenario.income_sources.all():
            IncomeSource.objects.create(
                scenario=duplicated_scenario,
                owned_by=income_source.owned_by,
                income_type=income_source.income_type,
                income_name=income_source.income_name,
                current_asset_balance=income_source.current_asset_balance,
                monthly_amount=income_source.monthly_amount,
                monthly_contribution=income_source.monthly_contribution,
                age_to_begin_withdrawal=income_source.age_to_begin_withdrawal,
                age_to_end_withdrawal=income_source.age_to_end_withdrawal,
                rate_of_return=income_source.rate_of_return,
                cola=income_source.cola,
                exclusion_ratio=income_source.exclusion_ratio,
                tax_rate=income_source.tax_rate,
                max_to_convert=income_source.max_to_convert
            )
        
        return Response({
            'id': duplicated_scenario.id,
            'name': duplicated_scenario.name,
            'message': 'Scenario duplicated successfully'
        }, status=status.HTTP_201_CREATED)
        
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_scenario_detail(request, scenario_id):
    """
    Get detailed scenario data for editing/duplicating
    """
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Get all income sources for this scenario
        income_sources = scenario.income_sources.all()
        income_data = []
        
        for income in income_sources:
            income_dict = {
                'id': str(uuid.uuid4()),  # Generate new UUID for frontend
                'income_type': income.income_type,
                'income_name': income.income_name,  # Add the missing income_name field
                'owned_by': income.owned_by,
                'start_age': income.age_to_begin_withdrawal,
                'end_age': income.age_to_end_withdrawal,
                'current_balance': float(income.current_asset_balance or 0),
                'monthly_contribution': float(income.monthly_contribution or 0),
                'growth_rate': income.rate_of_return,
                'withdrawal_amount': float(income.monthly_amount or 0),
                'amount_per_month': float(income.monthly_amount or 0),
                'amount_at_fra': float(income.monthly_amount or 0),
                'cola': income.cola,
                'exclusion_ratio': income.exclusion_ratio,
                'tax_rate': income.tax_rate,
                'max_to_convert': float(income.max_to_convert or 0),
                'age_established': income.age_established,
                'is_contributing': income.is_contributing,
                'employer_match': income.employer_match,
                'age_last_contribution': income.age_last_contribution
            }
            
            # Add specific fields for different income types
            if income.income_type == 'social_security':
                income_dict['amount_at_fra'] = float(income.monthly_amount or 0)
            elif income.income_type == 'Life_Insurance':
                income_dict['loan_amount'] = float(income.monthly_amount or 0)
            elif income.income_type == 'Annuity':
                income_dict['percent_taxable'] = income.exclusion_ratio * 100
            
            income_data.append(income_dict)
        
        # Check if this is for duplication (add "Copy") or editing (keep original name)
        is_for_duplication = request.GET.get('mode') == 'duplicate'
        scenario_name = f"{scenario.name} (Copy)" if is_for_duplication else scenario.name
        
        scenario_data = {
            'name': scenario_name,
            'description': scenario.description,
            'primary_retirement_age': scenario.retirement_age,
            'primary_medicare_age': scenario.medicare_age,
            'primary_lifespan': scenario.mortality_age,
            'spouse_retirement_age': scenario.spouse_retirement_age,
            'spouse_medicare_age': scenario.spouse_medicare_age,
            'spouse_lifespan': scenario.spouse_mortality_age,
            'model_tax_change': '',  # Reset to default
            'reduction_2030_ss': scenario.reduction_2030_ss,
            'ss_adjustment_year': scenario.ss_adjustment_year,
            'ss_adjustment_direction': scenario.ss_adjustment_direction,
            'ss_adjustment_type': scenario.ss_adjustment_type,
            'ss_adjustment_amount': scenario.ss_adjustment_amount,
            'apply_standard_deduction': scenario.apply_standard_deduction,
            'federal_standard_deduction': float(scenario.federal_standard_deduction or 0),
            'state_standard_deduction': float(scenario.state_standard_deduction or 0),
            'custom_annual_deduction': float(scenario.custom_annual_deduction or 0),
            'primary_blind': scenario.primary_blind,
            'spouse_blind': scenario.spouse_blind,
            'is_dependent': scenario.is_dependent,
            'part_b_inflation_rate': str(int(scenario.part_b_inflation_rate)),
            'part_d_inflation_rate': str(int(scenario.part_d_inflation_rate)),
            'primary_state': scenario.primary_state or '',  # Use stored value or default
            'income': income_data
        }
        
        return Response(scenario_data, status=status.HTTP_200_OK)
        
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_scenario_for_editing(request, scenario_id):
    """
    Get scenario data formatted specifically for editing in ScenarioCreate.vue
    """
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Get all income sources for this scenario
        income_sources = scenario.income_sources.all()
        income_data = []
        
        for income in income_sources:
            income_dict = {
                'id': income.id,  # Use actual database ID
                'income_type': income.income_type,
                'income_name': income.income_name,
                'owned_by': income.owned_by,
                'start_age': income.age_to_begin_withdrawal,
                'end_age': income.age_to_end_withdrawal,
                'current_balance': float(income.current_asset_balance or 0),
                'monthly_contribution': float(income.monthly_contribution or 0),
                'growth_rate': income.rate_of_return,
                'withdrawal_amount': float(income.monthly_amount or 0),
                'amount_per_month': float(income.monthly_amount or 0),
                'amount_at_fra': float(income.monthly_amount or 0),
                'cola': income.cola,
                'exclusion_ratio': income.exclusion_ratio,
                'tax_rate': income.tax_rate,
                'max_to_convert': float(income.max_to_convert or 0),
                'age_established': income.age_established,
                'is_contributing': income.is_contributing,
                'employer_match': income.employer_match,
                'age_last_contribution': income.age_last_contribution
            }
            
            # Add specific fields for different income types
            if income.income_type == 'social_security':
                income_dict['amount_at_fra'] = float(income.monthly_amount or 0)
            elif income.income_type == 'Life_Insurance':
                income_dict['loan_amount'] = float(income.monthly_amount or 0)
            elif income.income_type == 'Annuity':
                income_dict['percent_taxable'] = income.exclusion_ratio * 100
            
            income_data.append(income_dict)
        
        # Check if this is for duplication (add "Copy") or keep original name for editing
        is_for_duplication = request.GET.get('mode') == 'duplicate'
        scenario_name = f"{scenario.name} (Copy)" if is_for_duplication else scenario.name
        
        scenario_data = {
            'name': scenario_name,
            'description': scenario.description,
            'primary_retirement_age': scenario.retirement_age,
            'primary_medicare_age': scenario.medicare_age,
            'primary_lifespan': scenario.mortality_age,
            'spouse_retirement_age': scenario.spouse_retirement_age,
            'spouse_medicare_age': scenario.spouse_medicare_age,
            'spouse_lifespan': scenario.spouse_mortality_age,
            'model_tax_change': '',  # Reset to default
            'reduction_2030_ss': scenario.reduction_2030_ss,
            'ss_adjustment_year': scenario.ss_adjustment_year,
            'ss_adjustment_direction': scenario.ss_adjustment_direction,
            'ss_adjustment_type': scenario.ss_adjustment_type,
            'ss_adjustment_amount': scenario.ss_adjustment_amount,
            'apply_standard_deduction': scenario.apply_standard_deduction,
            'federal_standard_deduction': float(scenario.federal_standard_deduction or 0),
            'state_standard_deduction': float(scenario.state_standard_deduction or 0),
            'custom_annual_deduction': float(scenario.custom_annual_deduction or 0),
            'primary_blind': scenario.primary_blind,
            'spouse_blind': scenario.spouse_blind,
            'is_dependent': scenario.is_dependent,
            'part_b_inflation_rate': str(int(scenario.part_b_inflation_rate)),
            'part_d_inflation_rate': str(int(scenario.part_d_inflation_rate)),
            'primary_state': scenario.primary_state or '',
            'income': income_data
        }
        
        return Response(scenario_data, status=status.HTTP_200_OK)
        
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_scenario_assets(request, scenario_id):
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        assets = scenario.income_sources.all()
        asset_details = [
            {
                'id': asset.id,
                'owned_by': asset.owned_by,
                'income_type': asset.income_type,
                'income_name': asset.income_name,
                'current_asset_balance': asset.current_asset_balance,
                'monthly_amount': asset.monthly_amount,
                'monthly_contribution': asset.monthly_contribution,
                'age_to_begin_withdrawal': asset.age_to_begin_withdrawal,
                'age_to_end_withdrawal': asset.age_to_end_withdrawal,
                'rate_of_return': asset.rate_of_return,
                'cola': asset.cola,
                'exclusion_ratio': asset.exclusion_ratio,
                'tax_rate': asset.tax_rate,
                'scenario_id': asset.scenario_id
            }
            for asset in assets
        ]
        return Response(asset_details)
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=404)

@api_view(['PUT'])
@permission_classes([IsAuthenticated])
def update_scenario(request, scenario_id):
    """
    Update a scenario with Roth conversion parameters and income sources
    """
    print(f"💰 INCOME_EDIT: *** UPDATE_SCENARIO FUNCTION CALLED *** scenario_id={scenario_id}")
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Extract income sources from request data
        income_sources_data = request.data.pop('income_sources', None)
        
        # DEBUG: Log the income sources data
        print(f"💰 INCOME_EDIT: Received {len(income_sources_data) if income_sources_data else 0} income sources")
        if income_sources_data:
            for i, income_data in enumerate(income_sources_data):
                print(f"💰 INCOME_EDIT: Source {i}: Type={income_data.get('income_type')}, Name={income_data.get('income_name')}")
                print(f"💰 INCOME_EDIT: Source {i}: monthly_amount={income_data.get('monthly_amount')}, amount_at_fra={income_data.get('amount_at_fra')}")
                print(f"💰 INCOME_EDIT: Source {i}: Raw data: {income_data}")
        
        # Update scenario fields
        serializer = ScenarioUpdateSerializer(scenario, data=request.data, partial=True)
        if serializer.is_valid():
            serializer.save()
            
            # Update income sources if provided
            if income_sources_data is not None:
                # Delete existing income sources for this scenario
                scenario.income_sources.all().delete()
                
                # Create new income sources
                for income_data in income_sources_data:
                    # Remove any frontend-specific fields
                    income_data.pop('id', None)  # Remove frontend ID
                    
                    # Map frontend field names to backend model field names
                    # Special handling for social security - prioritize amount_at_fra
                    if income_data.get('income_type') == 'social_security':
                        mapped_monthly_amount = income_data.get('amount_at_fra') or income_data.get('monthly_amount') or income_data.get('withdrawal_amount', 0)
                    else:
                        mapped_monthly_amount = income_data.get('monthly_amount') or income_data.get('withdrawal_amount') or income_data.get('amount_at_fra', 0)
                    
                    print(f"💰 INCOME_EDIT: Mapping {income_data.get('income_type', 'unknown')} - mapped_monthly_amount={mapped_monthly_amount}")
                    print(f"💰 INCOME_EDIT: Source values - monthly_amount={income_data.get('monthly_amount')}, withdrawal_amount={income_data.get('withdrawal_amount')}, amount_at_fra={income_data.get('amount_at_fra')}")
                    
                    mapped_data = {
                        'owned_by': income_data.get('owned_by'),
                        'income_type': income_data.get('income_type'),
                        'income_name': income_data.get('income_name', ''),
                        'current_asset_balance': income_data.get('current_asset_balance') or income_data.get('current_balance', 0),
                        'monthly_amount': mapped_monthly_amount,
                        'monthly_contribution': income_data.get('monthly_contribution', 0),
                        'age_to_begin_withdrawal': income_data.get('age_to_begin_withdrawal') or income_data.get('start_age'),
                        'age_to_end_withdrawal': income_data.get('age_to_end_withdrawal') or income_data.get('end_age'),
                        'rate_of_return': income_data.get('rate_of_return') or (float(income_data.get('growth_rate', 0)) / 100 if income_data.get('growth_rate') else 0),
                        'cola': income_data.get('cola', 0),
                        'exclusion_ratio': income_data.get('exclusion_ratio', 0),
                        'tax_rate': income_data.get('tax_rate', 0),
                        'max_to_convert': income_data.get('max_to_convert', 0),
                        'age_established': income_data.get('age_established'),
                        'is_contributing': income_data.get('is_contributing', False),
                        'employer_match': income_data.get('employer_match', 0),
                        'age_last_contribution': income_data.get('age_last_contribution')
                    }
                    
                    # Remove None values to use model defaults
                    mapped_data = {k: v for k, v in mapped_data.items() if v is not None}
                    
                    IncomeSource.objects.create(scenario=scenario, **mapped_data)
            
            # Return updated scenario with income sources
            return Response({
                'id': scenario.id,
                'message': 'Scenario updated successfully'
            })
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)

@api_view(['DELETE'])
@permission_classes([IsAuthenticated])
def delete_scenario(request, client_id, scenario_id):
    """
    Delete a scenario
    """
    try:
        # Get the scenario and verify ownership
        scenario = Scenario.objects.get(id=scenario_id, client_id=client_id)
        
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Delete the scenario
        scenario.delete()
        
        return Response({"message": "Scenario deleted successfully"}, status=status.HTTP_204_NO_CONTENT)
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)

@api_view(['PATCH'])
@permission_classes([IsAuthenticated])
def toggle_scenario_sharing(request, scenario_id):
    """
    Toggle the share_with_client setting for a scenario
    """
    try:
        scenario = get_object_or_404(Scenario, pk=scenario_id)
        
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Get the new sharing state from request data
        share_with_client = request.data.get('share_with_client', not scenario.share_with_client)
        
        # Update the scenario
        scenario.share_with_client = share_with_client
        scenario.save()
        
        return Response({
            'id': scenario.id,
            'share_with_client': scenario.share_with_client,
            'message': f"Scenario {'shared with' if share_with_client else 'removed from'} client portal successfully"
        }, status=status.HTTP_200_OK)
        
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_400_BAD_REQUEST)

@api_view(['PUT'])
@permission_classes([IsAuthenticated])
def update_scenario_percentages(request, scenario_id):
    """
    Update the income_vs_cost_percent and medicare_irmaa_percent fields for a scenario.
    """
    try:
        scenario = get_object_or_404(Scenario, pk=scenario_id)
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Extract percentages from request data
        income_vs_cost_percent = request.data.get('income_vs_cost_percent')
        medicare_irmaa_percent = request.data.get('medicare_irmaa_percent')
        
        # Update only if the fields are provided
        if income_vs_cost_percent is not None:
            scenario.income_vs_cost_percent = income_vs_cost_percent
        
        if medicare_irmaa_percent is not None:
            scenario.medicare_irmaa_percent = medicare_irmaa_percent
        
        scenario.save()
        
        return Response({
            "id": scenario.id,
            "income_vs_cost_percent": scenario.income_vs_cost_percent,
            "medicare_irmaa_percent": scenario.medicare_irmaa_percent
        }, status=status.HTTP_200_OK)
        
    except Exception as e:
        return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['PUT'])
@permission_classes([IsAuthenticated])
def update_scenario_assets(request, scenario_id):
    """
    Update assets associated with a scenario, particularly for Roth conversion
    """
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Update assets for the scenario
        if 'assets' not in request.data:
            return Response(
                {"error": "Assets data is required"},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        assets_data = request.data['assets']
        updated_assets = []
        
        for asset_data in assets_data:
            asset_id = asset_data.get('id')
            if not asset_id:
                continue
                
            try:
                asset = IncomeSource.objects.get(id=asset_id, scenario=scenario)
                
                # Only update max_to_convert field
                if 'max_to_convert' in asset_data:
                    asset.max_to_convert = asset_data['max_to_convert']
                    asset.save()
                    updated_assets.append({
                        'id': asset.id,
                        'income_type': asset.income_type,
                        'max_to_convert': asset.max_to_convert
                    })
            except IncomeSource.DoesNotExist:
                # Skip assets that don't exist or don't belong to this scenario
                continue
        
        return Response({
            'message': f'Updated {len(updated_assets)} assets',
            'assets': updated_assets
        })
    except Scenario.DoesNotExist:
        return Response({'error': 'Scenario not found.'}, status=status.HTTP_404_NOT_FOUND)

class RothConversionAPIView(APIView):
    """
    API endpoint for processing Roth conversion scenarios.
    
    This endpoint accepts scenario data, client data, and Roth conversion parameters,
    and returns a detailed analysis of the Roth conversion impact, including:
    - Baseline scenario (without Roth conversion)
    - Conversion scenario (with Roth conversion)
    - Comparison of key metrics (taxes, Medicare costs, RMDs, inheritance tax)
    - Asset balance projections for visualization
    """
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]
    
    def _format_comparison_metrics(self, comparison):
        """
        Format comparison metrics to match the expected frontend structure.
        
        The frontend expects metrics like:
        {
            'rmd_reduction': 100000.0,
            'rmd_reduction_pct': 10.0,
            'tax_savings': 50000.0,
            'tax_savings_pct': 5.0,
            ...
        }
        """
        formatted = {}
        
        # Format RMD reduction
        if 'total_rmds' in comparison:
            rmd_diff = comparison['total_rmds']['difference']
            formatted['rmd_reduction'] = -rmd_diff  # Negate to get reduction
            formatted['rmd_reduction_pct'] = -comparison['total_rmds']['percent_change']
        
        # Format tax savings
        if 'lifetime_tax' in comparison:
            tax_diff = comparison['lifetime_tax']['difference']
            formatted['tax_savings'] = -tax_diff  # Negate to get savings
            formatted['tax_savings_pct'] = -comparison['lifetime_tax']['percent_change']
        
        # Format Medicare savings
        if 'lifetime_medicare' in comparison:
            medicare_diff = comparison['lifetime_medicare']['difference']
            formatted['medicare_savings'] = -medicare_diff
            formatted['medicare_savings_pct'] = -comparison['lifetime_medicare']['percent_change']
        
        # Format IRMAA savings
        if 'total_irmaa' in comparison:
            irmaa_diff = comparison['total_irmaa']['difference']
            formatted['irmaa_savings'] = -irmaa_diff
            formatted['irmaa_savings_pct'] = -comparison['total_irmaa']['percent_change']
        
        # Format inheritance tax savings
        if 'inheritance_tax' in comparison:
            inheritance_diff = comparison['inheritance_tax']['difference']
            formatted['inheritance_tax_savings'] = -inheritance_diff
            formatted['inheritance_tax_savings_pct'] = -comparison['inheritance_tax']['percent_change']
        
        # Format net income increase
        if 'cumulative_net_income' in comparison:
            income_diff = comparison['cumulative_net_income']['difference']
            formatted['net_income_increase'] = income_diff
            formatted['net_income_increase_pct'] = comparison['cumulative_net_income']['percent_change']
        
        # Format Roth increase
        if 'final_roth' in comparison:
            roth_diff = comparison['final_roth']['difference']
            formatted['roth_increase'] = roth_diff
            formatted['roth_increase_pct'] = comparison['final_roth']['percent_change']
        
        # Format total savings
        if 'total_expenses' in comparison:
            total_diff = comparison['total_expenses']['difference']
            formatted['total_savings'] = -total_diff
            formatted['total_savings_pct'] = -comparison['total_expenses']['percent_change']
        
        return formatted
    
    def post(self, request):
        try:
            # Extract data from request
            scenario = request.data.get('scenario')
            client = request.data.get('client')
            spouse = request.data.get('spouse')
            assets = request.data.get('assets')
            conversion_params = request.data.get('optimizer_params')
            
            # Log received data
            print("Received data for Roth conversion:")
            print(f"Scenario: {scenario}")
            print(f"Client: {client}")
            print(f"Spouse: {spouse}")
            print(f"Assets count: {len(assets) if assets else 0}")
            print(f"Conversion params: {conversion_params}")
            
            # Validate required fields
            if not scenario or not client or assets is None or not conversion_params:
                missing = []
                if not scenario: missing.append('scenario')
                if not client: missing.append('client')
                if assets is None: missing.append('assets')
                if not conversion_params: missing.append('optimizer_params')
                error_msg = f"Missing required fields: {', '.join(missing)}"
                print(f"ERROR: {error_msg}")
                return Response({'error': error_msg}, status=status.HTTP_400_BAD_REQUEST)
            
            # Create and process the Roth conversion
            processor = RothConversionProcessor(
                scenario=scenario,
                client=client,
                spouse=spouse,
                assets=assets,
                conversion_params=conversion_params
            )
            
            # Process the conversion and get results
            result = processor.process()
            
            # Transform the result to match the expected frontend structure
            transformed_result = {
                'baseline': {
                    'metrics': result['metrics']['baseline'],
                    'year_by_year': result['baseline_results']
                },
                'conversion': {
                    'metrics': result['metrics']['conversion'],
                    'year_by_year': result['conversion_results']
                },
                'comparison': self._format_comparison_metrics(result['metrics']['comparison']),
                'optimal_schedule': {
                    'start_year': result['conversion_params']['conversion_start_year'],
                    'duration': result['conversion_params']['years_to_convert'],
                    'annual_amount': result['conversion_params']['annual_conversion'],
                    'total_amount': result['conversion_params']['total_conversion'],
                    'score_breakdown': result['metrics']['conversion']
                },
                'asset_balances': result['asset_balances'],
                'scenarioResults': result['conversion_results']  # Add scenarioResults for frontend
            }
            
            # Return the transformed results
            return Response(transformed_result, status=status.HTTP_200_OK)
            
        except Exception as e:
            print(f"ERROR in Roth conversion processing: {str(e)}")
            import traceback
            traceback.print_exc()
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['POST'])
@permission_classes([AllowAny])
def register_advisor(request):
    """
    Register a new advisor (Step 1 and 2)
    """
    serializer = AdvisorRegistrationSerializer(data=request.data)
    
    if serializer.is_valid():
        try:
            # Create the user
            user = serializer.save()
            
            # Create Stripe customer
            stripe_customer = stripe.Customer.create(
                email=user.email,
                name=f"{user.first_name} {user.last_name}",
                metadata={
                    'user_id': user.id,
                    'company_name': user.company_name
                }
            )
            
            # Save Stripe customer ID to user model
            user.stripe_customer_id = stripe_customer.id
            user.save()

            # Generate tokens for automatic login
            refresh = RefreshToken.for_user(user)
            
            return Response({
                'status': 'success',
                'message': 'Advisor registered successfully',
                'tokens': {
                    'refresh': str(refresh),
                    'access': str(refresh.access_token),
                },
                'user': {
                    'id': user.id,
                    'email': user.email,
                    'first_name': user.first_name,
                    'last_name': user.last_name,
                }
            }, status=status.HTTP_201_CREATED)
            
        except Exception as e:
            # If there's an error, delete the user if it was created
            if 'user' in locals():
                user.delete()
            return Response({
                'status': 'error',
                'message': str(e)
            }, status=status.HTTP_400_BAD_REQUEST)
    
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def complete_registration(request):
    """
    Complete advisor registration with payment (Step 3)
    """
    try:
        user = request.user
        payment_method_id = request.data.get('paymentMethodId')
        plan = request.data.get('plan')

        if not payment_method_id or not plan:
            return Response({
                'status': 'error',
                'message': 'Payment method and plan are required'
            }, status=status.HTTP_400_BAD_REQUEST)

        # Attach payment method to customer
        stripe.PaymentMethod.attach(
            payment_method_id,
            customer=user.stripe_customer_id,
        )

        # Set as default payment method
        stripe.Customer.modify(
            user.stripe_customer_id,
            invoice_settings={
                'default_payment_method': payment_method_id
            }
        )

        # Get the price ID based on the plan
        price_id = settings.STRIPE_MONTHLY_PRICE_ID if plan == 'monthly' else settings.STRIPE_ANNUAL_PRICE_ID

        # Create the subscription
        subscription = stripe.Subscription.create(
            customer=user.stripe_customer_id,
            items=[{'price': price_id}],
            payment_behavior='default_incomplete',
            expand=['latest_invoice.payment_intent'],
            metadata={
                'user_id': user.id,
                'plan': plan
            }
        )

        # Save subscription ID to user model
        user.stripe_subscription_id = subscription.id
        user.subscription_status = subscription.status
        user.save()

        return Response({
            'status': 'success',
            'subscription': {
                'id': subscription.id,
                'status': subscription.status,
                'client_secret': subscription.latest_invoice.payment_intent.client_secret
            }
        })

    except stripe.error.StripeError as e:
        return Response({
            'status': 'error',
            'message': str(e)
        }, status=status.HTTP_400_BAD_REQUEST)
    except Exception as e:
        return Response({
            'status': 'error',
            'message': str(e)
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# Real Estate API Views
class ListCreateRealEstateView(generics.ListCreateAPIView):
    serializer_class = RealEstateSerializer
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]

    def get_queryset(self):
        client_id = self.kwargs['client_id']
        return RealEstate.objects.filter(client_id=client_id)

    def perform_create(self, serializer):
        client_id = self.kwargs['client_id']
        client = get_object_or_404(Client, id=client_id, advisor=self.request.user)
        serializer.save(client=client)

class RealEstateDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = RealEstate.objects.all()
    serializer_class = RealEstateSerializer
    permission_classes = [IsAuthenticated]
    authentication_classes = [JWTAuthentication]

class ReportTemplateViewSet(viewsets.ViewSet):
    permission_classes = [IsAuthenticated]
    
    def list(self, request, client_pk=None):
        """List all templates for a specific client"""
        client = get_object_or_404(Client, pk=client_pk, advisor=request.user)
        templates = ReportTemplate.objects.filter(client=client)
        serializer = ReportTemplateSerializer(templates, many=True, context={'request': request})
        return Response(serializer.data)
    
    def create(self, request, client_pk=None):
        """Upload a new PowerPoint template and process its slides"""
        try:
            # Check if python-pptx is installed
            from pptx import Presentation
        except ImportError:
            return Response(
                {"error": "The python-pptx library is not installed on the server."},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
        client = get_object_or_404(Client, pk=client_pk, advisor=request.user)
        
        # Validate request data
        if 'name' not in request.data:
            return Response({"error": "Template name is required"}, status=status.HTTP_400_BAD_REQUEST)
        
        if 'file' not in request.FILES:
            return Response({"error": "PowerPoint file is required"}, status=status.HTTP_400_BAD_REQUEST)
        
        pptx_file = request.FILES['file']
        
        # Check file extension
        if not pptx_file.name.endswith('.pptx'):
            return Response(
                {"error": "Only PowerPoint files (.pptx) are supported"},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Create the template record
        template = ReportTemplate.objects.create(
            client=client,
            name=request.data['name'],
            file=pptx_file
        )
        
        # Process the PowerPoint file to extract slide thumbnails
        try:
            self._process_slides(template, pptx_file)
        except Exception as e:
            # If slide processing fails, delete the template and return error
            template.delete()
            return Response(
                {"error": f"Failed to process PowerPoint file: {str(e)}"},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
        serializer = ReportTemplateSerializer(template, context={'request': request})
        return Response(serializer.data, status=status.HTTP_201_CREATED)
    
    def retrieve(self, request, pk=None):
        """Get details of a specific template including all slides"""
        template = get_object_or_404(ReportTemplate, pk=pk)
        
        # Check if the user has access to this template
        if template.client.advisor != request.user:
            return Response(
                {"error": "You don't have permission to access this template"},
                status=status.HTTP_403_FORBIDDEN
            )
        
        serializer = ReportTemplateDetailSerializer(template, context={'request': request})
        return Response(serializer.data)
    
    def destroy(self, request, pk=None):
        """Delete a template and all its slides"""
        template = get_object_or_404(ReportTemplate, pk=pk)
        
        # Check if the user has access to this template
        if template.client.advisor != request.user:
            return Response(
                {"error": "You don't have permission to delete this template"},
                status=status.HTTP_403_FORBIDDEN
            )
        
        template.delete()
        return Response(status=status.HTTP_204_NO_CONTENT)
    
    @action(detail=True, methods=['get'])
    def slides(self, request, pk=None):
        """Get all slides for a specific template"""
        template = get_object_or_404(ReportTemplate, pk=pk)
        
        # Check if the user has access to this template
        if template.client.advisor != request.user:
            return Response(
                {"error": "You don't have permission to access this template"},
                status=status.HTTP_403_FORBIDDEN
            )
        
        slides = template.slides.all().order_by('order')
        serializer = TemplateSlideSerializer(slides, many=True, context={'request': request})
        
        return Response({"slides": serializer.data})
    
    @action(detail=True, methods=['post'])
    def update_slides(self, request, pk=None):
        """Update slide order and remove slides"""
        template = get_object_or_404(ReportTemplate, pk=pk)
        
        # Check if the user has access to this template
        if template.client.advisor != request.user:
            return Response(
                {"error": "You don't have permission to modify this template"},
                status=status.HTTP_403_FORBIDDEN
            )
        
        if 'slides' not in request.data:
            return Response(
                {"error": "Slides data is required"},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Get all current slides
        current_slides = {str(slide.id): slide for slide in template.slides.all()}
        
        # Process the update
        for idx, slide_data in enumerate(request.data['slides']):
            slide_id = slide_data.get('id')
            
            if not slide_id or slide_id not in current_slides:
                continue
                
            # Update the order
            slide = current_slides[slide_id]
            slide.order = idx
            slide.save()
            
            # Remove from the current_slides dict to track which ones to delete
            current_slides.pop(slide_id)
        
        # Delete any slides not included in the request
        for slide in current_slides.values():
            slide.delete()
        
        # Return the updated slides
        slides = template.slides.all().order_by('order')
        serializer = TemplateSlideSerializer(slides, many=True, context={'request': request})
        
        return Response({"slides": serializer.data})
    
    def _process_slides(self, template, pptx_file):
        """Process a PowerPoint file and extract slide thumbnails"""
        try:
            # Create a temporary file to save the uploaded file
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                for chunk in pptx_file.chunks():
                    temp_file.write(chunk)
                temp_file_path = temp_file.name
            
            # Create output directory for thumbnails
            thumbnails_dir = os.path.join(settings.MEDIA_ROOT, 'thumbnails')
            os.makedirs(thumbnails_dir, exist_ok=True)
            
            # Initialize png_paths to empty list
            png_paths = []
            
            try:
                # Convert PowerPoint to PNG using our new converter
                png_paths = convert_pptx_to_png(temp_file_path, thumbnails_dir, dpi=300)
                
                # Save each thumbnail as a slide
                for i, png_path in enumerate(png_paths):
                    # Open the PNG file
                    with open(png_path, 'rb') as f:
                        image_data = f.read()
                    
                    # Create the slide record
                    slide_filename = f"slide_{i+1}.png"
                    template_slide = TemplateSlide(
                        template=template,
                        order=i
                    )
                    template_slide.thumbnail.save(slide_filename, BytesIO(image_data), save=True)
            finally:
                # Clean up temporary files
                try:
                    os.unlink(temp_file_path)
                except:
                    pass
                
                for png_path in png_paths:
                    try:
                        os.unlink(png_path)
                    except:
                        pass
        
        except Exception as e:
            import logging
            logger = logging.getLogger(__name__)
            logger.error(f"Error processing PowerPoint file: {str(e)}")
            raise e

@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated])
def client_report_templates(request, client_id):
    """
    GET: List all report templates for a client
    POST: Create a new report template for a client
    """
    client = get_object_or_404(Client, pk=client_id, advisor=request.user)
    
    if request.method == 'GET':
        templates = ReportTemplate.objects.filter(client=client)
        serializer = ReportTemplateSerializer(templates, many=True, context={'request': request})
        return Response(serializer.data)
    
    elif request.method == 'POST':
        try:
            # Check if python-pptx is installed
            from pptx import Presentation
        except ImportError:
            return Response(
                {"error": "The python-pptx library is not installed on the server."},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
        # Validate request data
        if 'name' not in request.data:
            return Response({"error": "Template name is required"}, status=status.HTTP_400_BAD_REQUEST)
        
        if 'file' not in request.FILES:
            return Response({"error": "PowerPoint file is required"}, status=status.HTTP_400_BAD_REQUEST)
        
        pptx_file = request.FILES['file']
        
        # Check file extension
        if not pptx_file.name.endswith('.pptx'):
            return Response(
                {"error": "Only PowerPoint files (.pptx) are supported"},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Create the template record
        template = ReportTemplate.objects.create(
            client=client,
            name=request.data['name'],
            file=pptx_file
        )
        
        # Process the PowerPoint file to extract slide thumbnails
        try:
            process_pptx_slides(template, pptx_file)
        except Exception as e:
            # If slide processing fails, delete the template and return error
            template.delete()
            return Response(
                {"error": f"Failed to process PowerPoint file: {str(e)}"},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
        serializer = ReportTemplateSerializer(template, context={'request': request})
        return Response(serializer.data, status=status.HTTP_201_CREATED)

def process_pptx_slides(template, pptx_file):
    """Process a PowerPoint file and extract slide thumbnails"""
    try:
        # Create a temporary file to save the uploaded file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            for chunk in pptx_file.chunks():
                temp_file.write(chunk)
            temp_file_path = temp_file.name
        
        # Create output directory for thumbnails
        thumbnails_dir = os.path.join(settings.MEDIA_ROOT, 'thumbnails')
        os.makedirs(thumbnails_dir, exist_ok=True)
        
        # Initialize png_paths to empty list
        png_paths = []
        
        try:
            # Convert PowerPoint to PNG using our new converter
            png_paths = convert_pptx_to_png(temp_file_path, thumbnails_dir, dpi=300)
            
            # Save each thumbnail as a slide
            for i, png_path in enumerate(png_paths):
                # Open the PNG file
                with open(png_path, 'rb') as f:
                    image_data = f.read()
                
                # Create the slide record
                slide_filename = f"slide_{i+1}.png"
                template_slide = TemplateSlide(
                    template=template,
                    order=i
                )
                template_slide.thumbnail.save(slide_filename, BytesIO(image_data), save=True)
        finally:
            # Clean up temporary files
            try:
                os.unlink(temp_file_path)
            except:
                pass
            
            for png_path in png_paths:
                try:
                    os.unlink(png_path)
                except:
                    pass
    
    except Exception as e:
        import logging
        logger = logging.getLogger(__name__)
        logger.error(f"Error processing PowerPoint file: {str(e)}")
        raise e

@api_view(['GET', 'DELETE'])
@permission_classes([IsAuthenticated])
def report_template_detail(request, template_id):
    """
    GET: Retrieve details of a specific template
    DELETE: Delete a template
    """
    template = get_object_or_404(ReportTemplate, pk=template_id)
    
    # Check if the user has access to this template
    if template.client.advisor != request.user:
        return Response(
            {"error": "You don't have permission to access this template"},
            status=status.HTTP_403_FORBIDDEN
        )
    
    if request.method == 'GET':
        serializer = ReportTemplateDetailSerializer(template, context={'request': request})
        return Response(serializer.data)
    
    elif request.method == 'DELETE':
        template.delete()
        return Response(status=status.HTTP_204_NO_CONTENT)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def template_slides(request, template_id):
    """Get all slides for a specific template"""
    template = get_object_or_404(ReportTemplate, pk=template_id)
    
    # Check if the user has access to this template
    if template.client.advisor != request.user:
        return Response(
            {"error": "You don't have permission to access this template"},
            status=status.HTTP_403_FORBIDDEN
        )
    
    slides = template.slides.all().order_by('order')
    serializer = TemplateSlideSerializer(slides, many=True, context={'request': request})
    
    return Response({"slides": serializer.data})

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def update_template_slides(request, template_id):
    """Update slide order and remove slides"""
    template = get_object_or_404(ReportTemplate, pk=template_id)
    
    # Check if the user has access to this template
    if template.client.advisor != request.user:
        return Response(
            {"error": "You don't have permission to modify this template"},
            status=status.HTTP_403_FORBIDDEN
        )
    
    if 'slides' not in request.data:
        return Response(
            {"error": "Slides data is required"},
            status=status.HTTP_400_BAD_REQUEST
        )
    
    # Get all current slides
    current_slides = {str(slide.id): slide for slide in template.slides.all()}
    
    # Process the update
    for idx, slide_data in enumerate(request.data['slides']):
        slide_id = slide_data.get('id')
        
        if not slide_id or slide_id not in current_slides:
            continue
            
        # Update the order
        slide = current_slides[slide_id]
        slide.order = idx
        slide.save()
        
        # Remove from the current_slides dict to track which ones to delete
        current_slides.pop(slide_id)
    
    # Delete any slides not included in the request
    for slide in current_slides.values():
        slide.delete()
    
    # Return the updated slides
    slides = template.slides.all().order_by('order')
    serializer = TemplateSlideSerializer(slides, many=True, context={'request': request})
    
    return Response({"slides": serializer.data})

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_scenario_comparison_data(request, scenario_id):
    """Get scenario data formatted for comparison report"""
    logger = logging.getLogger(__name__)
    
    logger.info(f"COMPARISON ENDPOINT CALLED: scenario_id={scenario_id}, user={request.user}")
    
    try:
        scenario = Scenario.objects.get(id=scenario_id)
        
        # Check if the user owns the client associated with this scenario
        if scenario.client.advisor != request.user:
            return Response({"error": "Access denied."}, status=403)
        
        # Run calculation to get the latest data
        try:
            processor = ScenarioProcessor(scenario_id=scenario.id, debug=False)
            calculation_result = processor.calculate()
            logger.info(f"Calculation result type: {type(calculation_result)}")
            logger.info(f"Calculation result: {str(calculation_result)[:500]}...")  # Log first 500 chars
        except Exception as calc_error:
            logger.error(f"Error calculating scenario {scenario_id}: {str(calc_error)}")
            # Return basic data without calculation
            comparison_data = {
                'id': scenario.id,
                'name': scenario.name,
                'irmaa_reached': False,
                'medicare_cost': 0,
                'federal_taxes': 0,
                'solution_cost': 0,
                'total_costs': 0,
                'out_of_pocket': 0
            }
            return Response(comparison_data, status=status.HTTP_200_OK)
        
        # Handle both list and dict return types from calculate()
        summary = {}
        if isinstance(calculation_result, list):
            logger.info(f"Calculation result is list with {len(calculation_result)} items")
            # If it's a list, look for the Summary in the last element or iterate
            for i, item in enumerate(calculation_result):
                logger.info(f"Item {i}: type={type(item)}, keys={list(item.keys()) if isinstance(item, dict) else 'N/A'}")
                if isinstance(item, dict) and 'Summary' in item:
                    summary = item['Summary']
                    logger.info(f"Found Summary in item {i}: {summary}")
                    break
            # If no Summary found in list items, check if last item is Summary
            if not summary and calculation_result:
                last_item = calculation_result[-1]
                if isinstance(last_item, dict):
                    summary = last_item
                    logger.info(f"Using last item as summary: {summary}")
        else:
            logger.info(f"Calculation result is dict with keys: {list(calculation_result.keys()) if isinstance(calculation_result, dict) else 'N/A'}")
            # If it's a dict, extract Summary as before
            summary = calculation_result.get('Summary', {})
            logger.info(f"Extracted summary: {summary}")
        
        logger.info(f"Final summary data: {summary}")
        
        # Calculate totals from the actual calculation result data
        total_federal_taxes = 0
        total_medicare_costs = 0
        total_gross_income = 0
        irmaa_reached = False
        
        # The calculation result is a list of yearly data
        if isinstance(calculation_result, list):
            for year_data in calculation_result:
                if isinstance(year_data, dict):
                    # Extract values and convert Decimal to float
                    fed_tax_raw = year_data.get('federal_tax', 0)
                    medicare_raw = year_data.get('total_medicare', 0)
                    gross_raw = year_data.get('gross_income', 0)
                    irmaa_raw = year_data.get('irmaa_surcharge', 0)
                    
                    # Convert to float, handling Decimal objects
                    fed_tax_float = float(fed_tax_raw) if fed_tax_raw is not None else 0
                    medicare_float = float(medicare_raw) if medicare_raw is not None else 0
                    gross_float = float(gross_raw) if gross_raw is not None else 0
                    irmaa_float = float(irmaa_raw) if irmaa_raw is not None else 0
                    
                    # Add to totals
                    total_federal_taxes += fed_tax_float
                    total_medicare_costs += medicare_float
                    total_gross_income += gross_float
                    
                    # Check if IRMAA is reached (any year with IRMAA surcharge > 0)
                    if irmaa_float > 0:
                        irmaa_reached = True
        
        # Calculate derived values
        total_costs = total_federal_taxes + total_medicare_costs
        out_of_pocket = total_gross_income - total_costs  # Net income
        solution_cost = 0  # This might need to be calculated differently based on business logic
        
        comparison_data = {
            'id': scenario.id,
            'name': scenario.name,
            'irmaa_reached': irmaa_reached,
            'medicare_cost': total_medicare_costs,
            'federal_taxes': total_federal_taxes,
            'solution_cost': solution_cost,
            'total_costs': total_costs,
            'out_of_pocket': out_of_pocket
        }
        
        logger.info(f"Final comparison data: {comparison_data}")
        
        return Response(comparison_data, status=status.HTTP_200_OK)
        
    except Scenario.DoesNotExist:
        return Response({"error": "Scenario not found"}, status=status.HTTP_404_NOT_FOUND)
    except Exception as e:
        logger.error(f"Error in get_scenario_comparison_data: {str(e)}")
        return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated])
def comparison_preferences(request, client_id):
    """Get or update comparison preferences for a client"""
    from .models import ComparisonPreference, Client, Scenario
    
    try:
        client = Client.objects.get(id=client_id, advisor=request.user)
    except Client.DoesNotExist:
        return Response({"error": "Client not found or access denied"}, status=status.HTTP_404_NOT_FOUND)
    
    if request.method == 'GET':
        # Get existing preferences
        try:
            preference = ComparisonPreference.objects.get(user=request.user, client=client)
            return Response({
                'scenario1': preference.scenario1.id if preference.scenario1 else None,
                'scenario2': preference.scenario2.id if preference.scenario2 else None
            })
        except ComparisonPreference.DoesNotExist:
            return Response({
                'scenario1': None,
                'scenario2': None
            })
    
    elif request.method == 'POST':
        # Save or update preferences
        scenario1_id = request.data.get('scenario1')
        scenario2_id = request.data.get('scenario2')
        
        # Validate scenarios belong to the client
        scenario1 = None
        scenario2 = None
        
        if scenario1_id:
            try:
                scenario1 = Scenario.objects.get(id=scenario1_id, client=client)
            except Scenario.DoesNotExist:
                return Response({"error": "Scenario 1 not found"}, status=status.HTTP_400_BAD_REQUEST)
        
        if scenario2_id:
            try:
                scenario2 = Scenario.objects.get(id=scenario2_id, client=client)
            except Scenario.DoesNotExist:
                return Response({"error": "Scenario 2 not found"}, status=status.HTTP_400_BAD_REQUEST)
        
        # Update or create preference
        preference, created = ComparisonPreference.objects.update_or_create(
            user=request.user,
            client=client,
            defaults={
                'scenario1': scenario1,
                'scenario2': scenario2
            }
        )
        
        return Response({
            'message': 'Preferences saved successfully',
            'scenario1': preference.scenario1.id if preference.scenario1 else None,
            'scenario2': preference.scenario2.id if preference.scenario2 else None
        }, status=status.HTTP_200_OK)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_federal_standard_deduction(request):
    """Get federal standard deduction amount based on filing status"""
    from .tax_csv_loader import get_tax_loader
    
    filing_status = request.GET.get('filing_status', 'Single')
    
    try:
        loader = get_tax_loader()
        deduction_amount = loader.get_standard_deduction(filing_status)
        
        return Response({
            'filing_status': filing_status,
            'deduction_amount': float(deduction_amount)
        })
    except Exception as e:
        return Response({
            'error': f'Error retrieving standard deduction: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_irmaa_thresholds_for_years(request):
    """Get IRMAA thresholds inflated for each year in a scenario"""
    from .tax_csv_loader import get_tax_loader
    
    filing_status = request.GET.get('filing_status', 'Single')
    start_year = int(request.GET.get('start_year', 2025))
    end_year = int(request.GET.get('end_year', 2045))
    
    try:
        loader = get_tax_loader()
        
        # Get thresholds for each year
        year_thresholds = {}
        for year in range(start_year, end_year + 1):
            thresholds = loader.get_inflated_irmaa_thresholds(filing_status, year)
            year_thresholds[year] = [
                {
                    'magi_threshold': float(t['magi_threshold']),
                    'part_b_surcharge': float(t['part_b_surcharge']),
                    'part_d_surcharge': float(t['part_d_surcharge']),
                    'description': t['description']
                }
                for t in thresholds if t['magi_threshold'] > 0
            ]
        
        return Response({
            'filing_status': filing_status,
            'start_year': start_year,
            'end_year': end_year,
            'thresholds_by_year': year_thresholds
        })
    except Exception as e:
        return Response({
            'error': f'Error retrieving IRMAA thresholds: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


# =============================================================================
# CRM VIEWS - Phase 1 Implementation
# =============================================================================

# =============================================================================
# EMAIL ACCOUNT MANAGEMENT
# =============================================================================

class EmailAccountViewSet(viewsets.ModelViewSet):
    """ViewSet for managing email accounts"""
    serializer_class = EmailAccountSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        return EmailAccount.objects.filter(user=self.request.user)
    
    def perform_create(self, serializer):
        serializer.save(user=self.request.user)
    
    @action(detail=True, methods=['post'])
    def sync(self, request, pk=None):
        """Trigger email sync for this account"""
        email_account = self.get_object()
        
        try:
            from .services.email_service import EmailService
            
            service = EmailService(email_account)
            stats = service.sync_emails_to_crm()
            
            return Response({
                'success': True,
                'stats': stats,
                'message': f"Synced {stats['created']} new emails, updated {stats['updated']}"
            })
            
        except Exception as e:
            logger.error(f"Email sync failed for account {pk}: {str(e)}")
            return Response({
                'error': f'Sync failed: {str(e)}'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
    @action(detail=True, methods=['post'])
    def validate(self, request, pk=None):
        """Validate email account connection"""
        email_account = self.get_object()
        
        try:
            from .services.oauth_service import EmailOAuthService
            
            oauth_service = EmailOAuthService(self.request.user, self.request)
            result = oauth_service.validate_account(email_account)
            
            return Response(result)
            
        except Exception as e:
            logger.error(f"Email account validation failed for {pk}: {str(e)}")
            return Response({
                'is_valid': False,
                'error': str(e)
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
    @action(detail=True, methods=['post'])
    def disconnect(self, request, pk=None):
        """Disconnect email account"""
        email_account = self.get_object()
        
        try:
            from .services.oauth_service import EmailOAuthService
            
            oauth_service = EmailOAuthService(self.request.user, self.request)
            success = oauth_service.disconnect_account(email_account)
            
            if success:
                return Response({'message': 'Email account disconnected successfully'})
            else:
                return Response({
                    'error': 'Failed to disconnect account'
                }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
                
        except Exception as e:
            logger.error(f"Email account disconnect failed for {pk}: {str(e)}")
            return Response({
                'error': str(e)
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def gmail_auth_url(request):
    """Get Gmail OAuth authorization URL"""
    try:
        from .services.oauth_service import EmailOAuthService
        
        oauth_service = EmailOAuthService(request.user, request)
        auth_url = oauth_service.get_gmail_auth_url()
        
        return Response({
            'auth_url': auth_url,
            'provider': 'gmail'
        })
        
    except Exception as e:
        logger.error(f"Failed to generate Gmail auth URL: {str(e)}")
        return Response({
            'error': f'Failed to generate auth URL: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def gmail_oauth_callback(request):
    """Handle Gmail OAuth callback"""
    code = request.data.get('code')
    state = request.data.get('state')
    
    if not code:
        return Response({
            'error': 'Authorization code is required'
        }, status=status.HTTP_400_BAD_REQUEST)
    
    try:
        from .services.oauth_service import EmailOAuthService
        
        oauth_service = EmailOAuthService(request.user, request)
        email_account = oauth_service.handle_gmail_callback(code, state)
        
        serializer = EmailAccountSerializer(email_account)
        return Response({
            'success': True,
            'account': serializer.data,
            'message': f'Gmail account {email_account.email_address} linked successfully'
        })
        
    except Exception as e:
        logger.error(f"Gmail OAuth callback failed: {str(e)}")
        return Response({
            'error': f'Failed to link Gmail account: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def outlook_auth_url(request):
    """Get Outlook OAuth authorization URL"""
    try:
        from .services.oauth_service import EmailOAuthService
        
        oauth_service = EmailOAuthService(request.user, request)
        auth_url = oauth_service.get_outlook_auth_url()
        
        return Response({
            'auth_url': auth_url,
            'provider': 'outlook'
        })
        
    except Exception as e:
        logger.error(f"Failed to generate Outlook auth URL: {str(e)}")
        return Response({
            'error': f'Failed to generate auth URL: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def outlook_oauth_callback(request):
    """Handle Outlook OAuth callback"""
    code = request.data.get('code')
    state = request.data.get('state')
    
    if not code:
        return Response({
            'error': 'Authorization code is required'
        }, status=status.HTTP_400_BAD_REQUEST)
    
    try:
        from .services.oauth_service import EmailOAuthService
        
        oauth_service = EmailOAuthService(request.user, request)
        email_account = oauth_service.handle_outlook_callback(code, state)
        
        serializer = EmailAccountSerializer(email_account)
        return Response({
            'success': True,
            'account': serializer.data,
            'message': f'Outlook account {email_account.email_address} linked successfully'
        })
        
    except Exception as e:
        logger.error(f"Outlook OAuth callback failed: {str(e)}")
        return Response({
            'error': f'Failed to link Outlook account: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def send_email(request):
    """Send an email through connected email account"""
    serializer = EmailComposeSerializer(data=request.data, context={'request': request})
    
    if not serializer.is_valid():
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
    
    try:
        from .services.email_service import EmailService
        
        # Get email account
        email_account_id = serializer.validated_data['email_account_id']
        email_account = EmailAccount.objects.get(
            id=email_account_id, 
            user=request.user, 
            is_active=True
        )
        
        # Send email
        service = EmailService(email_account)
        message_id = service.send_email(
            to_addresses=serializer.validated_data['to_addresses'],
            subject=serializer.validated_data['subject'],
            body=serializer.validated_data['body'],
            cc_addresses=serializer.validated_data.get('cc_addresses'),
            attachments=serializer.validated_data.get('attachments')
        )
        
        # Create communication record
        communication_data = {
            'communication_type': 'email',
            'direction': 'outbound',
            'subject': serializer.validated_data['subject'],
            'content': serializer.validated_data['body'],
            'from_address': email_account.email_address,
            'to_addresses': serializer.validated_data['to_addresses'],
            'cc_addresses': serializer.validated_data.get('cc_addresses', []),
            'email_account': email_account,
            'provider_message_id': message_id,
            'sent_at': timezone.now(),
            'sync_status': 'synced',
            'sync_direction': 'from_crm'
        }
        
        # Link to client or lead if specified
        client_id = serializer.validated_data.get('client_id')
        lead_id = serializer.validated_data.get('lead_id')
        
        if client_id:
            communication_data['client_id'] = client_id
        elif lead_id:
            communication_data['lead_id'] = lead_id
        
        comm_serializer = CommunicationCreateSerializer(
            data=communication_data,
            context={'request': request}
        )
        
        if comm_serializer.is_valid():
            communication = comm_serializer.save()
        
        return Response({
            'success': True,
            'message_id': message_id,
            'message': 'Email sent successfully'
        })
        
    except EmailAccount.DoesNotExist:
        return Response({
            'error': 'Invalid email account'
        }, status=status.HTTP_400_BAD_REQUEST)
    except Exception as e:
        logger.error(f"Email send failed: {str(e)}")
        return Response({
            'error': f'Failed to send email: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


# =============================================================================
# COMMUNICATION MANAGEMENT
# =============================================================================

class CommunicationViewSet(viewsets.ModelViewSet):
    """ViewSet for managing communications with advanced search, filtering, and bulk operations"""
    serializer_class = CommunicationSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        queryset = Communication.objects.filter(advisor=self.request.user).select_related(
            'client', 'lead', 'email_account'
        ).prefetch_related('client', 'lead')
        
        # Search functionality
        search = self.request.query_params.get('search')
        if search:
            queryset = queryset.filter(
                models.Q(subject__icontains=search) |
                models.Q(content__icontains=search) |
                models.Q(from_address__icontains=search) |
                models.Q(client__first_name__icontains=search) |
                models.Q(client__last_name__icontains=search) |
                models.Q(lead__first_name__icontains=search) |
                models.Q(lead__last_name__icontains=search) |
                models.Q(lead__email__icontains=search)
            )
        
        # Filter by type
        comm_type = self.request.query_params.get('type')
        if comm_type:
            queryset = queryset.filter(communication_type=comm_type)
        
        # Filter by direction
        direction = self.request.query_params.get('direction')
        if direction:
            queryset = queryset.filter(direction=direction)
        
        # Filter by client
        client_id = self.request.query_params.get('client_id')
        if client_id:
            queryset = queryset.filter(client_id=client_id)
        
        # Filter by lead
        lead_id = self.request.query_params.get('lead_id')
        if lead_id:
            queryset = queryset.filter(lead_id=lead_id)
        
        # Filter by read status
        is_read = self.request.query_params.get('is_read')
        if is_read is not None:
            is_read_bool = is_read.lower() in ('true', '1', 'yes')
            queryset = queryset.filter(is_read=is_read_bool)
        
        # Filter by AI sentiment
        sentiment = self.request.query_params.get('sentiment')
        if sentiment:
            queryset = queryset.filter(ai_sentiment_label=sentiment)
        
        # Filter by AI priority score range
        priority_min = self.request.query_params.get('priority_min')
        priority_max = self.request.query_params.get('priority_max')
        if priority_min:
            queryset = queryset.filter(ai_priority_score__gte=priority_min)
        if priority_max:
            queryset = queryset.filter(ai_priority_score__lte=priority_max)
        
        # Filter by date range
        date_from = self.request.query_params.get('date_from')
        date_to = self.request.query_params.get('date_to')
        if date_from:
            queryset = queryset.filter(created_at__gte=date_from)
        if date_to:
            queryset = queryset.filter(created_at__lte=date_to)
        
        # Filter by email account
        email_account_id = self.request.query_params.get('email_account_id')
        if email_account_id:
            queryset = queryset.filter(email_account_id=email_account_id)
        
        # Sorting
        ordering = self.request.query_params.get('ordering', '-created_at')
        if ordering:
            queryset = queryset.order_by(ordering)
        
        return queryset
    
    def get_serializer_class(self):
        if self.action == 'create':
            return CommunicationCreateSerializer
        return CommunicationSerializer
    
    def perform_create(self, serializer):
        serializer.save(advisor=self.request.user)
    
    @action(detail=False, methods=['post'])
    def bulk_mark_read(self, request):
        """Mark multiple communications as read"""
        communication_ids = request.data.get('ids', [])
        if not communication_ids:
            return Response({
                'error': 'No communication IDs provided'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        updated_count = Communication.objects.filter(
            id__in=communication_ids,
            advisor=request.user
        ).update(
            is_read=True,
            read_at=timezone.now()
        )
        
        return Response({
            'success': True,
            'message': f'Marked {updated_count} communications as read',
            'updated_count': updated_count
        })
    
    @action(detail=False, methods=['post'])
    def bulk_mark_unread(self, request):
        """Mark multiple communications as unread"""
        communication_ids = request.data.get('ids', [])
        if not communication_ids:
            return Response({
                'error': 'No communication IDs provided'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        updated_count = Communication.objects.filter(
            id__in=communication_ids,
            advisor=request.user
        ).update(
            is_read=False,
            read_at=None
        )
        
        return Response({
            'success': True,
            'message': f'Marked {updated_count} communications as unread',
            'updated_count': updated_count
        })
    
    @action(detail=False, methods=['delete'])
    def bulk_delete(self, request):
        """Delete multiple communications"""
        communication_ids = request.data.get('ids', [])
        if not communication_ids:
            return Response({
                'error': 'No communication IDs provided'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        deleted_count = Communication.objects.filter(
            id__in=communication_ids,
            advisor=request.user
        ).delete()[0]
        
        return Response({
            'success': True,
            'message': f'Deleted {deleted_count} communications',
            'deleted_count': deleted_count
        })
    
    @action(detail=False, methods=['post'])
    def sync_emails(self, request):
        """Trigger email sync for all email accounts"""
        from .tasks import sync_emails_batch
        from .models import EmailAccount
        
        email_accounts = EmailAccount.objects.filter(
            advisor=request.user,
            is_active=True
        )
        
        if not email_accounts.exists():
            return Response({
                'error': 'No active email accounts found'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        # Queue email sync tasks
        task_ids = []
        for account in email_accounts:
            task = sync_emails_batch.apply_async(
                args=[account.id],
                queue='email_sync'
            )
            task_ids.append(task.id)
        
        return Response({
            'success': True,
            'message': f'Email sync queued for {email_accounts.count()} accounts',
            'task_ids': task_ids,
            'account_count': email_accounts.count()
        })
    
    @action(detail=False, methods=['get'])
    def sync_status(self, request):
        """Get email sync status for all accounts"""
        from .models import EmailAccount
        
        accounts = EmailAccount.objects.filter(
            advisor=request.user
        ).values(
            'id', 'email', 'provider', 'is_active', 'last_sync_at', 'sync_status'
        )
        
        return Response({
            'accounts': list(accounts),
            'total_accounts': accounts.count(),
            'active_accounts': sum(1 for a in accounts if a['is_active'])
        })
    
    @action(detail=False, methods=['get'])
    def analytics(self, request):
        """Get communication analytics"""
        from django.db.models import Count, Q, Avg
        from datetime import datetime, timedelta
        
        queryset = Communication.objects.filter(advisor=request.user)
        
        # Date range filter for analytics
        days = int(request.query_params.get('days', 30))
        start_date = timezone.now() - timedelta(days=days)
        queryset = queryset.filter(created_at__gte=start_date)
        
        # Basic counts
        total_communications = queryset.count()
        unread_count = queryset.filter(is_read=False).count()
        
        # By type
        by_type = list(queryset.values('communication_type').annotate(
            count=Count('id')
        ).order_by('-count'))
        
        # By direction
        by_direction = list(queryset.values('direction').annotate(
            count=Count('id')
        ).order_by('-count'))
        
        # AI sentiment analysis
        sentiment_stats = list(queryset.exclude(
            ai_sentiment_label__isnull=True
        ).values('ai_sentiment_label').annotate(
            count=Count('id')
        ).order_by('-count'))
        
        # Average AI scores
        ai_averages = queryset.aggregate(
            avg_sentiment=Avg('ai_sentiment_score'),
            avg_urgency=Avg('ai_urgency_score'),
            avg_priority=Avg('ai_priority_score')
        )
        
        # High priority communications (priority > 0.7)
        high_priority_count = queryset.filter(
            ai_priority_score__gt=0.7
        ).count()
        
        # Daily communication counts for the period
        daily_counts = []
        for i in range(days):
            date = start_date + timedelta(days=i)
            count = queryset.filter(
                created_at__date=date.date()
            ).count()
            daily_counts.append({
                'date': date.date().isoformat(),
                'count': count
            })
        
        return Response({
            'period_days': days,
            'total_communications': total_communications,
            'unread_count': unread_count,
            'read_percentage': round((total_communications - unread_count) / max(total_communications, 1) * 100, 1),
            'by_type': by_type,
            'by_direction': by_direction,
            'sentiment_stats': sentiment_stats,
            'ai_averages': ai_averages,
            'high_priority_count': high_priority_count,
            'daily_counts': daily_counts
        })


# =============================================================================
# LEAD MANAGEMENT
# =============================================================================

class LeadSourceViewSet(viewsets.ModelViewSet):
    """ViewSet for managing lead sources"""
    serializer_class = LeadSourceSerializer
    permission_classes = [IsAuthenticated]
    queryset = LeadSource.objects.all()


class LeadViewSet(viewsets.ModelViewSet):
    """ViewSet for managing leads"""
    serializer_class = LeadSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        queryset = Lead.objects.filter(advisor=self.request.user)
        
        # Filter by status
        status_filter = self.request.query_params.get('status')
        if status_filter:
            queryset = queryset.filter(status=status_filter)
        
        # Filter by lead source
        source_id = self.request.query_params.get('source_id')
        if source_id:
            queryset = queryset.filter(lead_source_id=source_id)
        
        # Search by name or email
        search = self.request.query_params.get('search')
        if search:
            queryset = queryset.filter(
                models.Q(first_name__icontains=search) |
                models.Q(last_name__icontains=search) |
                models.Q(email__icontains=search)
            )
        
        return queryset.order_by('-created_at')
    
    def get_serializer_class(self):
        if self.action == 'create':
            return LeadCreateSerializer
        return LeadSerializer
    
    def perform_create(self, serializer):
        serializer.save(advisor=self.request.user)
    
    @action(detail=True, methods=['post'])
    def convert_to_client(self, request, pk=None):
        """Convert lead to client"""
        lead = self.get_object()
        
        if lead.status == 'converted':
            return Response({
                'error': 'Lead is already converted'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        try:
            from django.utils import timezone
            
            # Create client from lead data
            client_data = {
                'first_name': lead.first_name,
                'last_name': lead.last_name,
                'email': lead.email,
                'birthdate': request.data.get('birthdate'),
                'gender': request.data.get('gender', 'Other'),
                'tax_status': request.data.get('tax_status', 'Single'),
            }
            
            client_serializer = ClientCreateSerializer(
                data=client_data,
                context={'request': request}
            )
            
            if client_serializer.is_valid():
                client = client_serializer.save(advisor=request.user)
                
                # Update lead status
                lead.status = 'converted'
                lead.converted_client = client
                lead.conversion_date = timezone.now()
                lead.save()
                
                # Log activity
                ActivityLog.objects.create(
                    activity_type='lead_converted',
                    user=request.user,
                    lead=lead,
                    client=client,
                    description=f'Converted lead {lead} to client {client}',
                    metadata={
                        'lead_id': lead.id,
                        'client_id': client.id,
                        'lead_email': lead.email
                    }
                )
                
                return Response({
                    'success': True,
                    'client': ClientSerializer(client).data,
                    'message': f'Lead converted to client successfully'
                })
            else:
                return Response(client_serializer.errors, status=status.HTTP_400_BAD_REQUEST)
                
        except Exception as e:
            logger.error(f"Lead conversion failed: {str(e)}")
            return Response({
                'error': f'Failed to convert lead: {str(e)}'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


# =============================================================================
# ACTIVITY LOG
# =============================================================================

class ActivityLogViewSet(viewsets.ReadOnlyModelViewSet):
    """ViewSet for viewing activity logs"""
    serializer_class = ActivityLogSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        queryset = ActivityLog.objects.filter(user=self.request.user)
        
        # Filter by activity type
        activity_type = self.request.query_params.get('type')
        if activity_type:
            queryset = queryset.filter(activity_type=activity_type)
        
        # Filter by client
        client_id = self.request.query_params.get('client_id')
        if client_id:
            queryset = queryset.filter(client_id=client_id)
        
        # Filter by lead
        lead_id = self.request.query_params.get('lead_id')
        if lead_id:
            queryset = queryset.filter(lead_id=lead_id)
        
        # Limit results for performance
        limit = min(int(self.request.query_params.get('limit', 50)), 200)
        
        return queryset.order_by('-created_at')[:limit]


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def sync_all_emails(request):
    """Trigger email sync for all user's accounts"""
    try:
        from .services.email_service import EmailSyncManager
        
        results = EmailSyncManager.sync_all_accounts(user_id=request.user.id)
        
        return Response({
            'success': True,
            'results': results,
            'message': f"Processed {results['accounts_processed']} accounts, "
                      f"created {results['total_created']} new communications"
        })
        
    except Exception as e:
        logger.error(f"Bulk email sync failed: {str(e)}")
        return Response({
            'error': f'Sync failed: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def oauth_settings_status(request):
    """Get OAuth provider configuration status"""
    try:
        from .services.oauth_service import validate_oauth_settings
        
        status_info = validate_oauth_settings()
        return Response(status_info)
        
    except Exception as e:
        logger.error(f"OAuth settings validation failed: {str(e)}")
        return Response({
            'error': str(e)
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


# =============================================================================
# AI ANALYSIS ENDPOINTS
# =============================================================================

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def analyze_communication(request, communication_id):
    """Trigger AI analysis for a specific communication"""
    try:
        from .tasks import process_email_with_ai
        
        # Check if communication exists and belongs to user
        communication = Communication.objects.filter(
            id=communication_id,
            advisor=request.user
        ).first()
        
        if not communication:
            return Response({
                'error': 'Communication not found'
            }, status=status.HTTP_404_NOT_FOUND)
        
        # Queue AI analysis task
        task = process_email_with_ai.apply_async(
            args=[communication_id],
            queue='ai_processing'
        )
        
        return Response({
            'success': True,
            'message': 'AI analysis queued for processing',
            'task_id': task.id,
            'communication_id': communication_id
        })
            
    except Exception as e:
        logger.error(f"Communication analysis endpoint failed: {str(e)}")
        return Response({
            'error': f'Failed to queue analysis: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def bulk_analyze_communications(request):
    """Trigger AI analysis for multiple communications"""
    try:
        from .tasks import batch_process_communications
        
        communication_ids = request.data.get('communication_ids', [])
        if not communication_ids or not isinstance(communication_ids, list):
            return Response({
                'error': 'communication_ids list is required'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        # Filter to only include user's communications
        user_comms = Communication.objects.filter(
            id__in=communication_ids,
            advisor=request.user
        ).values_list('id', flat=True)
        
        valid_ids = list(user_comms)
        if not valid_ids:
            return Response({
                'error': 'No valid communications found'
            }, status=status.HTTP_404_NOT_FOUND)
        
        # Queue batch analysis
        task = batch_process_communications.apply_async(
            args=[valid_ids],
            queue='ai_processing'
        )
        
        return Response({
            'success': True,
            'message': f"Queued {len(valid_ids)} communications for AI analysis",
            'task_id': task.id,
            'communication_count': len(valid_ids)
        })
        
    except Exception as e:
        logger.error(f"Bulk analysis endpoint failed: {str(e)}")
        return Response({
            'error': f'Failed to queue bulk analysis: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def ai_analysis_stats(request):
    """Get AI analysis statistics for the user"""
    try:
        from .tasks.ai_tasks import get_ai_analysis_stats
        
        stats = get_ai_analysis_stats(advisor_id=request.user.id)
        
        if 'error' in stats:
            return Response({
                'error': stats['error']
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
        
        return Response(stats)
        
    except Exception as e:
        logger.error(f"AI stats endpoint failed: {str(e)}")
        return Response({
            'error': f'Failed to get stats: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def high_priority_communications(request):
    """Get high priority communications that need attention"""
    try:
        from .services.ai_email_service import get_high_priority_communications
        
        limit = min(int(request.query_params.get('limit', 20)), 100)
        
        high_priority_comms = get_high_priority_communications(
            advisor_id=request.user.id,
            limit=limit
        )
        
        serializer = CommunicationSerializer(high_priority_comms, many=True)
        
        return Response({
            'communications': serializer.data,
            'count': len(high_priority_comms),
            'message': f"Found {len(high_priority_comms)} high priority communications"
        })
        
    except Exception as e:
        logger.error(f"High priority communications endpoint failed: {str(e)}")
        return Response({
            'error': f'Failed to get high priority communications: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def trigger_auto_analysis(request):
    """Trigger analysis of new unprocessed emails"""
    try:
        from .tasks import analyze_new_inbound_emails
        
        # Queue the auto-analysis task
        task = analyze_new_inbound_emails.apply_async(queue='ai_processing')
        
        return Response({
            'success': True,
            'message': 'Auto-analysis of new emails queued for processing',
            'task_id': task.id
        })
        
    except Exception as e:
        logger.error(f"Auto analysis trigger failed: {str(e)}")
        return Response({
            'error': f'Failed to queue auto analysis: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


# =============================================================================
# CELERY MONITORING & HEALTH CHECK ENDPOINTS
# =============================================================================

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def celery_health_check(request):
    """Check Celery worker health and queue status"""
    try:
        from retirementadvisorpro.celery import celery_health_check, get_queue_stats, get_worker_status
        from .tasks import health_check_task
        
        # Test basic Celery connectivity
        health_status = celery_health_check()
        
        # Get queue statistics
        queue_stats = get_queue_stats()
        
        # Get worker status
        worker_status = get_worker_status()
        
        # Test task execution
        task_test_success = False
        try:
            task = health_check_task.apply_async(expires=10)
            result = task.get(timeout=10)
            task_test_success = result.get('status') == 'healthy'
        except Exception as e:
            logger.warning(f"Task test failed: {str(e)}")
        
        return Response({
            'celery_healthy': health_status,
            'task_execution_test': task_test_success,
            'queue_stats': queue_stats,
            'worker_status': worker_status,
            'timestamp': timezone.now().isoformat()
        })
        
    except Exception as e:
        logger.error(f"Celery health check failed: {str(e)}")
        return Response({
            'celery_healthy': False,
            'error': str(e),
            'timestamp': timezone.now().isoformat()
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def task_status(request, task_id):
    """Get status of a specific Celery task"""
    try:
        from celery.result import AsyncResult
        
        task_result = AsyncResult(task_id)
        
        response = {
            'task_id': task_id,
            'status': task_result.status,
            'ready': task_result.ready(),
            'successful': task_result.successful() if task_result.ready() else None,
            'timestamp': timezone.now().isoformat()
        }
        
        if task_result.ready():
            if task_result.successful():
                response['result'] = task_result.result
            else:
                response['error'] = str(task_result.result)
                response['traceback'] = task_result.traceback
        
        return Response(response)
        
    except Exception as e:
        logger.error(f"Failed to get task status: {str(e)}")
        return Response({
            'error': f'Failed to get task status: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET']) 
@permission_classes([IsAuthenticated])
def queue_monitoring(request):
    """Get comprehensive queue monitoring information"""
    try:
        from retirementadvisorpro.celery import get_queue_stats
        from celery import current_app
        
        # Get basic queue stats
        queue_stats = get_queue_stats()
        
        # Get worker information
        inspect = current_app.control.inspect()
        worker_stats = {}
        
        if inspect:
            try:
                # Get worker stats
                stats = inspect.stats()
                active = inspect.active()
                registered = inspect.registered()
                
                if stats:
                    for worker, worker_info in stats.items():
                        worker_stats[worker] = {
                            'status': 'online',
                            'pool_writes': worker_info.get('pool', {}).get('writes', {}),
                            'total_tasks': worker_info.get('total', 0),
                            'active_tasks': len(active.get(worker, [])) if active else 0,
                            'registered_tasks': len(registered.get(worker, [])) if registered else 0
                        }
            except Exception as e:
                logger.warning(f"Failed to get detailed worker stats: {str(e)}")
        
        return Response({
            'queue_stats': queue_stats,
            'worker_stats': worker_stats,
            'total_workers': len(worker_stats),
            'timestamp': timezone.now().isoformat()
        })
        
    except Exception as e:
        logger.error(f"Queue monitoring failed: {str(e)}")
        return Response({
            'error': f'Queue monitoring failed: {str(e)}'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


# =============================================================================
# TASK MANAGEMENT API VIEWS
# =============================================================================

class TaskViewSet(viewsets.ModelViewSet):
    """ViewSet for managing tasks with advanced filtering, search, and bulk operations"""
    permission_classes = [IsAuthenticated]
    
    def get_serializer_class(self):
        if self.action == 'create':
            return TaskCreateSerializer
        elif self.action == 'retrieve':
            return TaskDetailSerializer
        return TaskSerializer
    
    def get_queryset(self):
        queryset = Task.objects.filter(
            models.Q(created_by=self.request.user) | 
            models.Q(assigned_to=self.request.user)
        ).select_related(
            'assigned_to', 'created_by', 'client', 'template', 'scenario'
        ).prefetch_related('comments')
        
        # Search functionality
        search = self.request.query_params.get('search')
        if search:
            queryset = queryset.filter(
                models.Q(title__icontains=search) |
                models.Q(description__icontains=search) |
                models.Q(tags__icontains=search) |
                models.Q(client__first_name__icontains=search) |
                models.Q(client__last_name__icontains=search)
            )
        
        # Filter by status
        status_filter = self.request.query_params.get('status')
        if status_filter:
            queryset = queryset.filter(status=status_filter)
        
        # Filter by priority
        priority = self.request.query_params.get('priority')
        if priority:
            queryset = queryset.filter(priority=priority)
        
        # Filter by assigned user
        assigned_to = self.request.query_params.get('assigned_to')
        if assigned_to:
            if assigned_to == 'me':
                queryset = queryset.filter(assigned_to=self.request.user)
            elif assigned_to == 'unassigned':
                queryset = queryset.filter(assigned_to__isnull=True)
            else:
                queryset = queryset.filter(assigned_to_id=assigned_to)
        
        # Filter by client
        client_id = self.request.query_params.get('client')
        if client_id:
            queryset = queryset.filter(client_id=client_id)
        
        # Filter by task type
        task_type = self.request.query_params.get('task_type')
        if task_type:
            queryset = queryset.filter(task_type=task_type)
        
        # Filter by overdue tasks
        if self.request.query_params.get('overdue') == 'true':
            queryset = queryset.filter(
                due_date__lt=timezone.now(),
                status__in=['pending', 'in_progress']
            )
        
        # Filter by due date range
        due_from = self.request.query_params.get('due_from')
        due_to = self.request.query_params.get('due_to')
        if due_from:
            queryset = queryset.filter(due_date__gte=due_from)
        if due_to:
            queryset = queryset.filter(due_date__lte=due_to)
        
        # Ordering
        ordering = self.request.query_params.get('ordering', '-created_at')
        if ordering:
            queryset = queryset.order_by(ordering)
        
        return queryset
    
    def perform_create(self, serializer):
        # Set created_by to current user
        task = serializer.save(created_by=self.request.user)
        
        # Create activity log entry
        ActivityLog.objects.create(
            user=self.request.user,
            activity_type='task_created',
            description=f"Created task: {task.title}",
            client=task.client,
            metadata={
                'task_id': task.id,
                'task_title': task.title,
                'priority': task.priority,
                'assigned_to': task.assigned_to.email if task.assigned_to else None
            }
        )
    
    def perform_update(self, serializer):
        old_task = self.get_object()
        old_status = old_task.status
        old_assigned_to = old_task.assigned_to
        
        task = serializer.save()
        
        # Track status changes
        if task.status != old_status:
            if task.status == 'completed' and old_status != 'completed':
                # Calculate actual duration
                if task.created_at:
                    task.actual_duration_hours = (timezone.now() - task.created_at).total_seconds() / 3600
                    task.completed_at = timezone.now()
                    task.save(update_fields=['actual_duration_hours', 'completed_at'])
                
                # Log completion
                ActivityLog.objects.create(
                    user=self.request.user,
                    activity_type='task_completed',
                    description=f"Completed task: {task.title}",
                    client=task.client,
                    metadata={
                        'task_id': task.id,
                        'task_title': task.title,
                        'actual_duration_hours': task.actual_duration_hours
                    }
                )
            elif old_status == 'completed' and task.status != 'completed':
                # Task reopened
                task.completed_at = None
                task.actual_duration_hours = None
                task.save(update_fields=['completed_at', 'actual_duration_hours'])
                
                ActivityLog.objects.create(
                    user=self.request.user,
                    activity_type='task_reopened',
                    description=f"Reopened task: {task.title}",
                    client=task.client,
                    metadata={'task_id': task.id, 'task_title': task.title}
                )
        
        # Track assignment changes
        if task.assigned_to != old_assigned_to:
            ActivityLog.objects.create(
                user=self.request.user,
                activity_type='task_assigned',
                description=f"Reassigned task: {task.title}",
                client=task.client,
                metadata={
                    'task_id': task.id,
                    'task_title': task.title,
                    'old_assigned_to': old_assigned_to.email if old_assigned_to else None,
                    'new_assigned_to': task.assigned_to.email if task.assigned_to else None
                }
            )
    
    @action(detail=True, methods=['post'])
    def add_comment(self, request, pk=None):
        """Add a comment to a task"""
        task = self.get_object()
        content = request.data.get('content', '').strip()
        
        if not content:
            return Response(
                {'error': 'Comment content is required'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        comment = TaskComment.objects.create(
            task=task,
            user=request.user,
            content=content
        )
        
        # Log comment activity
        ActivityLog.objects.create(
            user=request.user,
            activity_type='task_commented',
            description=f"Commented on task: {task.title}",
            client=task.client,
            metadata={
                'task_id': task.id,
                'task_title': task.title,
                'comment_preview': content[:50] + ('...' if len(content) > 50 else '')
            }
        )
        
        serializer = TaskCommentSerializer(comment)
        return Response(serializer.data, status=status.HTTP_201_CREATED)
    
    @action(detail=False, methods=['get'])
    def stats(self, request):
        """Get task statistics for current user"""
        user_tasks = Task.objects.filter(
            models.Q(created_by=request.user) | 
            models.Q(assigned_to=request.user)
        )
        
        # Basic counts
        total_tasks = user_tasks.count()
        completed_tasks = user_tasks.filter(status='completed').count()
        pending_tasks = user_tasks.filter(status='pending').count()
        in_progress_tasks = user_tasks.filter(status='in_progress').count()
        
        # Priority breakdown
        high_priority = user_tasks.filter(priority='high').count()
        medium_priority = user_tasks.filter(priority='medium').count()
        low_priority = user_tasks.filter(priority='low').count()
        
        # Overdue tasks
        overdue_tasks = user_tasks.filter(
            due_date__lt=timezone.now(),
            status__in=['pending', 'in_progress']
        ).count()
        
        # Due soon (next 7 days)
        due_soon = user_tasks.filter(
            due_date__gte=timezone.now(),
            due_date__lt=timezone.now() + timezone.timedelta(days=7),
            status__in=['pending', 'in_progress']
        ).count()
        
        return Response({
            'total_tasks': total_tasks,
            'completed_tasks': completed_tasks,
            'pending_tasks': pending_tasks,
            'in_progress_tasks': in_progress_tasks,
            'high_priority': high_priority,
            'medium_priority': medium_priority,
            'low_priority': low_priority,
            'overdue_tasks': overdue_tasks,
            'due_soon': due_soon,
            'completion_rate': round((completed_tasks / total_tasks * 100) if total_tasks > 0 else 0, 1)
        })
    
    @action(detail=False, methods=['post'])
    def bulk_update(self, request):
        """Bulk update multiple tasks"""
        task_ids = request.data.get('task_ids', [])
        updates = request.data.get('updates', {})
        
        if not task_ids or not updates:
            return Response(
                {'error': 'task_ids and updates are required'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Only allow updating tasks the user has access to
        tasks = Task.objects.filter(
            id__in=task_ids
        ).filter(
            models.Q(created_by=request.user) | models.Q(assigned_to=request.user)
        )
        
        if not tasks.exists():
            return Response(
                {'error': 'No accessible tasks found'},
                status=status.HTTP_404_NOT_FOUND
            )
        
        # Apply updates
        allowed_fields = ['status', 'priority', 'assigned_to', 'due_date', 'progress']
        update_data = {k: v for k, v in updates.items() if k in allowed_fields}
        
        if not update_data:
            return Response(
                {'error': 'No valid update fields provided'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        updated_count = tasks.update(**update_data)
        
        # Log bulk update
        ActivityLog.objects.create(
            user=request.user,
            activity_type='tasks_bulk_updated',
            description=f"Bulk updated {updated_count} tasks",
            metadata={
                'task_count': updated_count,
                'updates': update_data
            }
        )
        
        return Response({
            'message': f'Successfully updated {updated_count} tasks',
            'updated_count': updated_count
        })


class TaskTemplateViewSet(viewsets.ModelViewSet):
    """ViewSet for managing task templates"""
    serializer_class = TaskTemplateSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        queryset = TaskTemplate.objects.filter(created_by=self.request.user)
        
        # Filter by active status
        if self.request.query_params.get('active_only') == 'true':
            queryset = queryset.filter(is_active=True)
        
        # Filter by trigger type
        trigger_type = self.request.query_params.get('trigger_type')
        if trigger_type:
            queryset = queryset.filter(trigger_type=trigger_type)
        
        # Search
        search = self.request.query_params.get('search')
        if search:
            queryset = queryset.filter(
                models.Q(name__icontains=search) |
                models.Q(description__icontains=search) |
                models.Q(default_title__icontains=search)
            )
        
        return queryset.order_by('name')
    
    def perform_create(self, serializer):
        serializer.save(created_by=self.request.user)
    
    @action(detail=True, methods=['post'])
    def create_task_from_template(self, request, pk=None):
        """Create a task from this template"""
        template = self.get_object()
        
        # Get additional data from request
        client_id = request.data.get('client_id')
        lead_id = request.data.get('lead_id')
        assigned_to_id = request.data.get('assigned_to_id')
        custom_title = request.data.get('title')
        custom_description = request.data.get('description')
        custom_due_date = request.data.get('due_date')
        
        # Create task from template
        task_data = {
            'title': custom_title or template.default_title,
            'description': custom_description or template.default_description,
            'priority': template.default_priority,
            'task_type': 'template_generated',
            'template': template,
            'estimated_duration': template.estimated_duration,
            'created_by': request.user,
        }
        
        # Set relationships
        if client_id:
            try:
                client = Client.objects.get(id=client_id, advisor=request.user)
                task_data['client'] = client
            except Client.DoesNotExist:
                return Response(
                    {'error': 'Invalid client ID'},
                    status=status.HTTP_400_BAD_REQUEST
                )
        
        if lead_id:
            try:
                lead = Lead.objects.get(id=lead_id, advisor=request.user)
                task_data['lead'] = lead
            except Lead.DoesNotExist:
                return Response(
                    {'error': 'Invalid lead ID'},
                    status=status.HTTP_400_BAD_REQUEST
                )
        
        # Set assignment
        if assigned_to_id:
            try:
                User = get_user_model()
                assigned_user = User.objects.get(id=assigned_to_id)
                task_data['assigned_to'] = assigned_user
            except User.DoesNotExist:
                return Response(
                    {'error': 'Invalid assigned user ID'},
                    status=status.HTTP_400_BAD_REQUEST
                )
        elif template.auto_assign_to_creator:
            task_data['assigned_to'] = request.user
        
        # Set due date
        if custom_due_date:
            task_data['due_date'] = custom_due_date
        
        # Create the task
        task = Task.objects.create(**task_data)
        
        # Log activity
        ActivityLog.objects.create(
            user=request.user,
            activity_type='task_created_from_template',
            description=f"Created task from template: {template.name}",
            client=task.client,
            metadata={
                'task_id': task.id,
                'template_id': template.id,
                'template_name': template.name,
                'task_title': task.title
            }
        )
        
        serializer = TaskDetailSerializer(task)
        return Response(serializer.data, status=status.HTTP_201_CREATED)


# =============================================================================
# CALENDAR INTEGRATION API VIEWS
# =============================================================================

class CalendarAccountViewSet(viewsets.ModelViewSet):
    """ViewSet for managing calendar accounts (Google Calendar, Outlook Calendar)"""
    serializer_class = CalendarAccountSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        queryset = CalendarAccount.objects.filter(user=self.request.user)
        
        # Filter by provider
        provider = self.request.query_params.get('provider')
        if provider:
            queryset = queryset.filter(provider=provider)
        
        # Filter by active status
        if self.request.query_params.get('active_only') == 'true':
            queryset = queryset.filter(is_active=True)
        
        return queryset.order_by('-primary_calendar', '-created_at')
    
    def perform_create(self, serializer):
        serializer.save(user=self.request.user)
    
    @action(detail=True, methods=['post'])
    def sync_events(self, request, pk=None):
        """Manually sync events from external calendar"""
        calendar_account = self.get_object()
        
        try:
            from .services.calendar_service import CalendarService
            
            service = CalendarService(request.user, request)
            result = service.sync_calendar_events(calendar_account)
            
            # Log activity
            ActivityLog.objects.create(
                user=request.user,
                activity_type='calendar_sync',
                description=f"Synced calendar: {calendar_account.display_name}",
                metadata={
                    'calendar_account_id': calendar_account.id,
                    'sync_result': result
                }
            )
            
            return Response({
                'message': 'Calendar sync completed successfully',
                'result': result
            })
            
        except Exception as e:
            return Response(
                {'error': f'Calendar sync failed: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
    
    @action(detail=True, methods=['post'])
    def toggle_sync(self, request, pk=None):
        """Enable/disable calendar sync"""
        calendar_account = self.get_object()
        calendar_account.sync_enabled = not calendar_account.sync_enabled
        calendar_account.save(update_fields=['sync_enabled'])
        
        action_text = 'enabled' if calendar_account.sync_enabled else 'disabled'
        
        # Log activity
        ActivityLog.objects.create(
            user=request.user,
            activity_type='calendar_sync_toggled',
            description=f"Calendar sync {action_text}: {calendar_account.display_name}",
            metadata={
                'calendar_account_id': calendar_account.id,
                'sync_enabled': calendar_account.sync_enabled
            }
        )
        
        return Response({
            'message': f'Calendar sync {action_text}',
            'sync_enabled': calendar_account.sync_enabled
        })


class CalendarEventViewSet(viewsets.ModelViewSet):
    """ViewSet for managing calendar events with advanced filtering and search"""
    permission_classes = [IsAuthenticated]
    
    def get_serializer_class(self):
        if self.action == 'create':
            return CalendarEventCreateSerializer
        return CalendarEventSerializer
    
    def get_queryset(self):
        # Only show events from user's calendar accounts
        queryset = CalendarEvent.objects.filter(
            calendar_account__user=self.request.user
        ).select_related(
            'calendar_account', 'client', 'lead', 'task'
        ).prefetch_related('reminders')
        
        # Search functionality
        search = self.request.query_params.get('search')
        if search:
            queryset = queryset.filter(
                models.Q(title__icontains=search) |
                models.Q(description__icontains=search) |
                models.Q(location__icontains=search) |
                models.Q(client__first_name__icontains=search) |
                models.Q(client__last_name__icontains=search) |
                models.Q(lead__first_name__icontains=search) |
                models.Q(lead__last_name__icontains=search)
            )
        
        # Filter by date range
        start_date = self.request.query_params.get('start_date')
        end_date = self.request.query_params.get('end_date')
        if start_date:
            queryset = queryset.filter(start_datetime__gte=start_date)
        if end_date:
            queryset = queryset.filter(start_datetime__lte=end_date)
        
        # Filter by status
        status_filter = self.request.query_params.get('status')
        if status_filter:
            queryset = queryset.filter(status=status_filter)
        
        # Filter by client
        client_id = self.request.query_params.get('client')
        if client_id:
            queryset = queryset.filter(client_id=client_id)
        
        # Filter by lead
        lead_id = self.request.query_params.get('lead')
        if lead_id:
            queryset = queryset.filter(lead_id=lead_id)
        
        # Filter by upcoming events
        if self.request.query_params.get('upcoming') == 'true':
            queryset = queryset.filter(start_datetime__gte=timezone.now())
        
        # Filter by today's events
        if self.request.query_params.get('today') == 'true':
            today = timezone.now().date()
            queryset = queryset.filter(start_datetime__date=today)
        
        # Filter by calendar account
        calendar_account_id = self.request.query_params.get('calendar_account')
        if calendar_account_id:
            queryset = queryset.filter(calendar_account_id=calendar_account_id)
        
        # Ordering
        ordering = self.request.query_params.get('ordering', 'start_datetime')
        queryset = queryset.order_by(ordering)
        
        return queryset
    
    def perform_create(self, serializer):
        # For manual event creation, associate with user's primary calendar
        calendar_account = CalendarAccount.objects.filter(
            user=self.request.user, 
            primary_calendar=True,
            is_active=True
        ).first()
        
        if not calendar_account:
            # Use any active calendar account
            calendar_account = CalendarAccount.objects.filter(
                user=self.request.user,
                is_active=True
            ).first()
        
        if not calendar_account:
            raise ValidationError("No active calendar account found. Please connect a calendar first.")
        
        event = serializer.save(
            calendar_account=calendar_account,
            external_event_id=f"manual_{timezone.now().timestamp()}",
            calendar_id=calendar_account.external_account_id,
            is_synced=False  # Manual events need to be pushed to external calendar
        )
        
        # Log activity
        ActivityLog.objects.create(
            user=self.request.user,
            activity_type='calendar_event_created',
            description=f"Created calendar event: {event.title}",
            client=event.client,
            lead=event.lead,
            metadata={
                'event_id': event.id,
                'event_title': event.title,
                'start_datetime': event.start_datetime.isoformat()
            }
        )
    
    @action(detail=True, methods=['post'])
    def associate_client(self, request, pk=None):
        """Associate event with a client"""
        event = self.get_object()
        client_id = request.data.get('client_id')
        
        if client_id:
            try:
                client = Client.objects.get(id=client_id, advisor=request.user)
                event.client = client
                event.lead = None  # Clear lead if client is set
                event.save(update_fields=['client', 'lead'])
                
                # Log activity
                ActivityLog.objects.create(
                    user=request.user,
                    activity_type='calendar_event_associated',
                    description=f"Associated event '{event.title}' with client {client}",
                    client=client,
                    metadata={
                        'event_id': event.id,
                        'event_title': event.title
                    }
                )
                
                return Response({'message': 'Event associated with client successfully'})
                
            except Client.DoesNotExist:
                return Response(
                    {'error': 'Client not found'},
                    status=status.HTTP_404_NOT_FOUND
                )
        else:
            event.client = None
            event.save(update_fields=['client'])
            return Response({'message': 'Client association removed'})
    
    @action(detail=True, methods=['post'])
    def associate_lead(self, request, pk=None):
        """Associate event with a lead"""
        event = self.get_object()
        lead_id = request.data.get('lead_id')
        
        if lead_id:
            try:
                lead = Lead.objects.get(id=lead_id, advisor=request.user)
                event.lead = lead
                event.client = None  # Clear client if lead is set
                event.save(update_fields=['lead', 'client'])
                
                # Log activity
                ActivityLog.objects.create(
                    user=request.user,
                    activity_type='calendar_event_associated',
                    description=f"Associated event '{event.title}' with lead {lead}",
                    lead=lead,
                    metadata={
                        'event_id': event.id,
                        'event_title': event.title
                    }
                )
                
                return Response({'message': 'Event associated with lead successfully'})
                
            except Lead.DoesNotExist:
                return Response(
                    {'error': 'Lead not found'},
                    status=status.HTTP_404_NOT_FOUND
                )
        else:
            event.lead = None
            event.save(update_fields=['lead'])
            return Response({'message': 'Lead association removed'})
    
    @action(detail=False, methods=['get'])
    def upcoming_today(self, request):
        """Get today's upcoming events"""
        today = timezone.now().date()
        now = timezone.now()
        
        events = CalendarEvent.objects.filter(
            calendar_account__user=request.user,
            start_datetime__date=today,
            start_datetime__gte=now,
            status='confirmed'
        ).order_by('start_datetime')[:10]
        
        serializer = CalendarEventSerializer(events, many=True)
        return Response(serializer.data)
    
    @action(detail=False, methods=['get'])
    def stats(self, request):
        """Get calendar event statistics"""
        user_events = CalendarEvent.objects.filter(calendar_account__user=request.user)
        
        # Date ranges
        today = timezone.now().date()
        this_week_start = today - timezone.timedelta(days=today.weekday())
        this_month_start = today.replace(day=1)
        
        stats = {
            'total_events': user_events.count(),
            'today_events': user_events.filter(start_datetime__date=today).count(),
            'this_week_events': user_events.filter(
                start_datetime__date__gte=this_week_start,
                start_datetime__date__lt=this_week_start + timezone.timedelta(days=7)
            ).count(),
            'this_month_events': user_events.filter(
                start_datetime__date__gte=this_month_start
            ).count(),
            'upcoming_events': user_events.filter(
                start_datetime__gte=timezone.now(),
                status='confirmed'
            ).count(),
            'client_meetings': user_events.filter(client__isnull=False).count(),
            'lead_meetings': user_events.filter(lead__isnull=False).count(),
        }
        
        return Response(stats)


class MeetingTemplateViewSet(viewsets.ModelViewSet):
    """ViewSet for managing meeting templates"""
    serializer_class = MeetingTemplateSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        queryset = MeetingTemplate.objects.filter(user=self.request.user)
        
        # Filter by active status
        if self.request.query_params.get('active_only') == 'true':
            queryset = queryset.filter(is_active=True)
        
        # Search
        search = self.request.query_params.get('search')
        if search:
            queryset = queryset.filter(
                models.Q(name__icontains=search) |
                models.Q(description__icontains=search) |
                models.Q(default_title__icontains=search)
            )
        
        return queryset.order_by('-usage_count', 'name')
    
    def perform_create(self, serializer):
        serializer.save(user=self.request.user)
    
    @action(detail=True, methods=['post'])
    def schedule_meeting(self, request, pk=None):
        """Schedule a meeting using this template"""
        template = self.get_object()
        serializer = MeetingScheduleSerializer(data=request.data, context={'request': request})
        serializer.is_valid(raise_exception=True)
        
        data = serializer.validated_data
        
        # Get template defaults
        title = data.get('title') or template.default_title
        description = data.get('description') or template.default_description
        location = data.get('location') or template.default_location
        duration_minutes = data.get('duration_minutes') or template.default_duration
        
        # Calculate end time
        start_datetime = data['start_datetime']
        end_datetime = start_datetime + timezone.timedelta(minutes=duration_minutes)
        
        # Get calendar account
        calendar_account = CalendarAccount.objects.filter(
            user=request.user,
            is_active=True,
            primary_calendar=True
        ).first()
        
        if not calendar_account:
            calendar_account = CalendarAccount.objects.filter(
                user=request.user,
                is_active=True
            ).first()
        
        if not calendar_account:
            return Response(
                {'error': 'No active calendar account found'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Prepare attendees
        attendees = []
        attendee_emails = data.get('attendee_emails', [])
        
        # Add client/lead email to attendees
        if data.get('client_id'):
            try:
                client = Client.objects.get(id=data['client_id'], advisor=request.user)
                if client.email and client.email not in attendee_emails:
                    attendee_emails.append(client.email)
            except Client.DoesNotExist:
                pass
        
        if data.get('lead_id'):
            try:
                lead = Lead.objects.get(id=data['lead_id'], advisor=request.user)
                if lead.email and lead.email not in attendee_emails:
                    attendee_emails.append(lead.email)
            except Lead.DoesNotExist:
                pass
        
        # Format attendees for calendar
        for email in attendee_emails:
            attendees.append({
                'email': email,
                'displayName': email.split('@')[0].title(),
                'responseStatus': 'needsAction'
            })
        
        # Create calendar event
        event = CalendarEvent.objects.create(
            calendar_account=calendar_account,
            external_event_id=f"template_{timezone.now().timestamp()}",
            calendar_id=calendar_account.external_account_id,
            title=title,
            description=description,
            location=location,
            start_datetime=start_datetime,
            end_datetime=end_datetime,
            attendees=attendees,
            client_id=data.get('client_id'),
            lead_id=data.get('lead_id'),
            is_synced=False  # Will be synced to external calendar
        )
        
        # Create video meeting if requested
        if template.include_video_link:
            try:
                from .services.video_service import VideoConferenceService
                
                video_service = VideoConferenceService(request.user)
                meeting_type = data.get('meeting_type') or template.preferred_meeting_type
                meeting_result = video_service.create_video_meeting(event, meeting_type, template)
                
                # Add meeting URL to location if not already present
                if meeting_result.get('meeting_url') and not event.location:
                    event.location = meeting_result.get('meeting_url')
                    event.save(update_fields=['location'])
            except Exception as e:
                logger.warning(f"Failed to create video meeting: {str(e)}")
                # Continue even if video meeting creation fails
        
        # Create follow-up task if enabled
        if template.create_follow_up_task:
            follow_up_date = start_datetime + timezone.timedelta(days=template.follow_up_task_days)
            
            Task.objects.create(
                title=template.follow_up_task_title,
                description=f"Follow up on meeting: {title}",
                priority='medium',
                task_type='follow_up',
                assigned_to=request.user,
                created_by=request.user,
                client_id=data.get('client_id'),
                lead_id=data.get('lead_id'),
                due_date=follow_up_date
            )
        
        # Increment template usage
        template.increment_usage()
        
        # Log activity
        ActivityLog.objects.create(
            user=request.user,
            activity_type='meeting_scheduled',
            description=f"Scheduled meeting: {title}",
            client_id=data.get('client_id'),
            lead_id=data.get('lead_id'),
            metadata={
                'event_id': event.id,
                'template_id': template.id,
                'template_name': template.name,
                'meeting_time': start_datetime.isoformat()
            }
        )
        
        serializer = CalendarEventSerializer(event)
        return Response(serializer.data, status=status.HTTP_201_CREATED)


# =============================================================================
# CALENDAR OAUTH ENDPOINTS
# =============================================================================

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def google_calendar_auth_url(request):
    """Generate Google Calendar OAuth2 authorization URL"""
    try:
        from .services.calendar_service import CalendarService
        
        service = CalendarService(request.user, request)
        auth_url = service.get_google_calendar_auth_url()
        
        return Response({
            'auth_url': auth_url,
            'provider': 'google'
        })
        
    except Exception as e:
        return Response(
            {'error': f'Failed to generate Google Calendar auth URL: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['GET'])
@permission_classes([AllowAny])
def google_calendar_oauth_callback(request):
    """Handle Google Calendar OAuth2 callback"""
    try:
        code = request.GET.get('code')
        state = request.GET.get('state')
        error = request.GET.get('error')
        
        if error:
            return Response(
                {'error': f'Google Calendar OAuth error: {error}'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        if not code:
            return Response(
                {'error': 'Authorization code not provided'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Extract user ID from state
        user_id = state.split('_')[1] if state and '_' in state else None
        if not user_id:
            return Response(
                {'error': 'Invalid state parameter'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        user = get_user_model().objects.get(id=user_id)
        
        from .services.calendar_service import CalendarService
        service = CalendarService(user, request)
        calendar_account = service.handle_google_calendar_callback(code, state)
        
        # Return success page or redirect
        return Response({
            'success': True,
            'message': 'Google Calendar connected successfully',
            'calendar_account': CalendarAccountSerializer(calendar_account).data
        })
        
    except Exception as e:
        logger.error(f"Google Calendar callback error: {str(e)}")
        return Response(
            {'error': f'Google Calendar connection failed: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def outlook_calendar_auth_url(request):
    """Generate Outlook Calendar OAuth2 authorization URL"""
    try:
        from .services.calendar_service import CalendarService
        
        service = CalendarService(request.user, request)
        auth_url = service.get_outlook_calendar_auth_url()
        
        return Response({
            'auth_url': auth_url,
            'provider': 'outlook'
        })
        
    except Exception as e:
        return Response(
            {'error': f'Failed to generate Outlook Calendar auth URL: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['GET'])
@permission_classes([AllowAny])
def outlook_calendar_oauth_callback(request):
    """Handle Outlook Calendar OAuth2 callback"""
    try:
        code = request.GET.get('code')
        state = request.GET.get('state')
        error = request.GET.get('error')
        
        if error:
            return Response(
                {'error': f'Outlook Calendar OAuth error: {error}'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        if not code:
            return Response(
                {'error': 'Authorization code not provided'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Extract user ID from state
        user_id = state.split('_')[1] if state and '_' in state else None
        if not user_id:
            return Response(
                {'error': 'Invalid state parameter'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        user = get_user_model().objects.get(id=user_id)
        
        from .services.calendar_service import CalendarService
        service = CalendarService(user, request)
        calendar_account = service.handle_outlook_calendar_callback(code, state)
        
        return Response({
            'success': True,
            'message': 'Outlook Calendar connected successfully',
            'calendar_account': CalendarAccountSerializer(calendar_account).data
        })
        
    except Exception as e:
        logger.error(f"Outlook Calendar callback error: {str(e)}")
        return Response(
            {'error': f'Outlook Calendar connection failed: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def calendar_settings_status(request):
    """Get calendar OAuth settings status"""
    try:
        from .services.calendar_service import validate_calendar_settings
        
        status_info = validate_calendar_settings()
        
        # Add user's connected accounts
        user_accounts = CalendarAccount.objects.filter(user=request.user, is_active=True)
        
        accounts_info = {
            'google_connected': user_accounts.filter(provider='google').exists(),
            'outlook_connected': user_accounts.filter(provider='outlook').exists(),
            'total_accounts': user_accounts.count(),
            'accounts': CalendarAccountSerializer(user_accounts, many=True).data
        }
        
        return Response({
            **status_info,
            **accounts_info
        })
        
    except Exception as e:
        logger.error(f"Calendar settings status error: {str(e)}")
        return Response(
            {'error': f'Failed to get calendar settings status: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


# =============================================================================
# VIDEO CONFERENCING API ENDPOINTS
# =============================================================================

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def create_video_meeting(request, event_id):
    """Create a video meeting for a calendar event"""
    try:
        # Get the calendar event
        event = CalendarEvent.objects.get(
            id=event_id,
            calendar_account__user=request.user
        )
        
        # Get meeting type from request
        meeting_type = request.data.get('meeting_type', 'zoom')
        
        # Check if event already has a meeting
        if event.meeting_url:
            return Response(
                {'error': 'Event already has a video meeting'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Create video meeting
        from .services.video_service import VideoConferenceService
        
        service = VideoConferenceService(request.user)
        result = service.create_video_meeting(event, meeting_type)
        
        # Log activity
        ActivityLog.objects.create(
            user=request.user,
            activity_type='video_meeting_created',
            description=f"Created {meeting_type} meeting for: {event.title}",
            client=event.client,
            lead=event.lead,
            metadata={
                'event_id': event.id,
                'meeting_type': meeting_type,
                'meeting_url': result.get('meeting_url')
            }
        )
        
        # Return updated event data
        serializer = CalendarEventSerializer(event)
        return Response({
            'event': serializer.data,
            'meeting_details': result
        })
        
    except CalendarEvent.DoesNotExist:
        return Response(
            {'error': 'Calendar event not found'},
            status=status.HTTP_404_NOT_FOUND
        )
    except Exception as e:
        logger.error(f"Failed to create video meeting: {str(e)}")
        return Response(
            {'error': f'Failed to create video meeting: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['PUT'])
@permission_classes([IsAuthenticated])
def update_video_meeting(request, event_id):
    """Update a video meeting for a calendar event"""
    try:
        # Get the calendar event
        event = CalendarEvent.objects.get(
            id=event_id,
            calendar_account__user=request.user
        )
        
        # Check if event has a meeting
        if not event.meeting_url:
            return Response(
                {'error': 'Event does not have a video meeting'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Update video meeting
        from .services.video_service import VideoConferenceService
        
        service = VideoConferenceService(request.user)
        result = service.update_video_meeting(event)
        
        # Log activity
        ActivityLog.objects.create(
            user=request.user,
            activity_type='video_meeting_updated',
            description=f"Updated {event.meeting_type} meeting for: {event.title}",
            client=event.client,
            lead=event.lead,
            metadata={
                'event_id': event.id,
                'meeting_type': event.meeting_type
            }
        )
        
        return Response(result)
        
    except CalendarEvent.DoesNotExist:
        return Response(
            {'error': 'Calendar event not found'},
            status=status.HTTP_404_NOT_FOUND
        )
    except Exception as e:
        logger.error(f"Failed to update video meeting: {str(e)}")
        return Response(
            {'error': f'Failed to update video meeting: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['DELETE'])
@permission_classes([IsAuthenticated])
def delete_video_meeting(request, event_id):
    """Delete a video meeting from a calendar event"""
    try:
        # Get the calendar event
        event = CalendarEvent.objects.get(
            id=event_id,
            calendar_account__user=request.user
        )
        
        # Check if event has a meeting
        if not event.meeting_url:
            return Response(
                {'error': 'Event does not have a video meeting'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Delete video meeting
        from .services.video_service import VideoConferenceService
        
        service = VideoConferenceService(request.user)
        success = service.delete_video_meeting(event)
        
        if success:
            # Log activity
            ActivityLog.objects.create(
                user=request.user,
                activity_type='video_meeting_deleted',
                description=f"Deleted video meeting for: {event.title}",
                client=event.client,
                lead=event.lead,
                metadata={'event_id': event.id}
            )
            
            return Response({'message': 'Video meeting deleted successfully'})
        else:
            return Response(
                {'error': 'Failed to delete video meeting'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
    except CalendarEvent.DoesNotExist:
        return Response(
            {'error': 'Calendar event not found'},
            status=status.HTTP_404_NOT_FOUND
        )
    except Exception as e:
        logger.error(f"Failed to delete video meeting: {str(e)}")
        return Response(
            {'error': f'Failed to delete video meeting: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_meeting_join_info(request, event_id):
    """Get joining information for a video meeting"""
    try:
        # Get the calendar event
        event = CalendarEvent.objects.get(
            id=event_id,
            calendar_account__user=request.user
        )
        
        # Get meeting join information
        from .services.video_service import VideoConferenceService
        
        service = VideoConferenceService(request.user)
        join_info = service.get_meeting_join_info(event)
        
        return Response(join_info)
        
    except CalendarEvent.DoesNotExist:
        return Response(
            {'error': 'Calendar event not found'},
            status=status.HTTP_404_NOT_FOUND
        )
    except Exception as e:
        logger.error(f"Failed to get meeting join info: {str(e)}")
        return Response(
            {'error': f'Failed to get meeting join info: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def send_meeting_reminder(request, event_id):
    """Send meeting reminder to participants"""
    try:
        # Get the calendar event
        event = CalendarEvent.objects.get(
            id=event_id,
            calendar_account__user=request.user
        )
        
        # Get reminder minutes from request
        reminder_minutes = request.data.get('reminder_minutes', 15)
        immediate = request.data.get('immediate', False)
        
        # Send meeting reminder
        from .services.video_service import MeetingReminderService
        
        service = MeetingReminderService()
        
        if immediate:
            success = service.send_immediate_reminder(event)
            message = 'Meeting reminder sent immediately'
        else:
            success = service.send_meeting_reminder(event, reminder_minutes)
            message = f'Meeting reminder scheduled for {reminder_minutes} minutes before meeting'
        
        if success:
            # Create or update reminder record
            CalendarEventReminder.objects.update_or_create(
                event=event,
                reminder_type='email',
                minutes_before=reminder_minutes,
                defaults={
                    'remind_at': event.start_datetime - timezone.timedelta(minutes=reminder_minutes)
                }
            )
            
            # Log activity
            ActivityLog.objects.create(
                user=request.user,
                activity_type='meeting_reminder_sent',
                description=f"Scheduled reminder for: {event.title}",
                client=event.client,
                lead=event.lead,
                metadata={
                    'event_id': event.id,
                    'reminder_minutes': reminder_minutes,
                    'immediate': immediate
                }
            )
            
            return Response({'message': message})
        else:
            return Response(
                {'error': 'Failed to send meeting reminder'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
    except CalendarEvent.DoesNotExist:
        return Response(
            {'error': 'Calendar event not found'},
            status=status.HTTP_404_NOT_FOUND
        )
    except Exception as e:
        logger.error(f"Failed to send meeting reminder: {str(e)}")
        return Response(
            {'error': f'Failed to send meeting reminder: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def video_settings_status(request):
    """Get video conferencing settings status"""
    try:
        from .services.video_service import validate_video_settings
        
        status_info = validate_video_settings()
        
        return Response(status_info)
        
    except Exception as e:
        logger.error(f"Video settings status error: {str(e)}")
        return Response(
            {'error': f'Failed to get video settings status: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_jump_ai_meeting_insights(request, event_id):
    """Get AI insights and analytics from a completed Jump.ai meeting"""
    try:
        # Get the calendar event
        event = CalendarEvent.objects.get(
            id=event_id,
            calendar_account__user=request.user
        )
        
        # Verify this is a Jump.ai meeting
        if event.meeting_type != 'jump_ai':
            return Response(
                {'error': 'This is not a Jump.ai meeting'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Get insights from Jump.ai
        from .services.video_service import VideoConferenceService
        video_service = VideoConferenceService(request.user)
        
        insights = video_service.get_jump_ai_meeting_insights(event)
        
        if insights.get('success'):
            return Response({
                'meeting_id': event.meeting_id,
                'meeting_title': event.title,
                'insights': insights,
                'event_date': event.start_time,
                'duration_minutes': int((event.end_time - event.start_time).total_seconds() / 60)
            })
        else:
            return Response(
                {'error': insights.get('error', 'Failed to retrieve insights')},
                status=status.HTTP_404_NOT_FOUND
            )
        
    except CalendarEvent.DoesNotExist:
        return Response(
            {'error': 'Calendar event not found'},
            status=status.HTTP_404_NOT_FOUND
        )
    except Exception as e:
        logger.error(f"Jump.ai insights error: {str(e)}")
        return Response(
            {'error': f'Failed to get meeting insights: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


# Temporary mock Report Center endpoints for debugging
from rest_framework.decorators import api_view

@api_view(['GET', 'POST'])
@permission_classes([AllowAny])
def mock_report_templates(request):
    """Temporary mock endpoint for report templates while fixing Report Center imports"""
    if request.method == 'GET':
        # Return mock template data
        return Response({
            'results': [
                {
                    'id': 1,
                    'name': 'Comprehensive Retirement Report',
                    'description': 'Complete financial analysis and recommendations',
                    'template_type': 'comprehensive',
                    'is_public': True,
                    'is_active': True,
                    'created_by': request.user.id if request.user.is_authenticated else 1,
                    'created_at': '2025-01-01T00:00:00Z',
                    'updated_at': '2025-01-01T00:00:00Z',
                    'preview_image': None,
                    'sections': []
                },
                {
                    'id': 2,
                    'name': 'Tax Planning Analysis',
                    'description': 'Tax optimization strategies and projections',
                    'template_type': 'tax_planning',
                    'is_public': True,
                    'is_active': True,
                    'created_by': request.user.id if request.user.is_authenticated else 1,
                    'created_at': '2025-01-01T00:00:00Z',
                    'updated_at': '2025-01-01T00:00:00Z',
                    'preview_image': None,
                    'sections': []
                }
            ],
            'count': 2,
            'next': None,
            'previous': None
        })
    else:
        # Mock POST response
        return Response({
            'id': 3,
            'name': request.data.get('name', 'New Template'),
            'description': request.data.get('description', ''),
            'template_type': request.data.get('template_type', 'comprehensive'),
            'is_public': False,
            'is_active': True,
            'created_by': request.user.id if request.user.is_authenticated else 1,
            'created_at': '2025-01-01T00:00:00Z',
            'updated_at': '2025-01-01T00:00:00Z',
            'sections': []
        }, status=status.HTTP_201_CREATED)

@api_view(['GET', 'POST'])  
@permission_classes([AllowAny])
def mock_reports(request):
    """Temporary mock endpoint for reports"""
    if request.method == 'GET':
        return Response({
            'results': [],
            'count': 0,
            'next': None,
            'previous': None
        })
    else:
        # More complete mock response that matches what the frontend expects
        return Response({
            'id': 1,
            'name': request.data.get('name', 'New Report'),
            'description': request.data.get('description', ''),
            'client_id': request.data.get('client_id'),
            'scenario_id': request.data.get('scenario_id'),  
            'template_id': request.data.get('template_id'),
            'export_format': request.data.get('export_format', 'pdf'),
            'sections': request.data.get('sections', []),
            'status': request.data.get('status', 'draft'),
            'created_by': request.user.id if request.user.is_authenticated else 1,
            'created_at': '2025-01-01T00:00:00Z',
            'updated_at': '2025-01-01T00:00:00Z',
            'client_name': None,  # Would be populated from client_id in real implementation
            'template_name': None,  # Would be populated from template_id in real implementation
        }, status=status.HTTP_201_CREATED)

@api_view(['POST'])
@permission_classes([AllowAny])
def mock_generate_report(request, report_id):
    """Temporary mock endpoint for report generation"""
    format_type = request.data.get('format', 'pdf')
    
    return Response({
        'task_id': f'mock_task_{report_id}_{format_type}',
        'format': format_type,
        'report_id': report_id,
        'status': 'generating',
        'message': f'Report generation started for {format_type.upper()} format',
        'estimated_completion': '2025-01-01T00:05:00Z'  # Mock 5 minutes from now
    })

@api_view(['GET'])
@permission_classes([AllowAny])
def mock_report_status(request, report_id):
    """Temporary mock endpoint for checking report generation status"""
    # Simulate a completed report after some time (for demo purposes)
    return Response({
        'id': report_id,
        'status': 'completed',  # Could be 'generating', 'completed', 'failed'
        'generation_completed_at': '2025-01-01T00:02:00Z',
        'download_url': f'/api/report-center/reports/{report_id}/download/',
        'file_size': '2.4 MB',
        'pages': 12,
        'message': 'Report generation completed successfully'
    })

