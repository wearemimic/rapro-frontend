import os
import sys
import glob
import logging
from PIL import Image
from pptx import Presentation
from io import BytesIO
from core.pptx_utils import check_external_tools, convert_pptx_to_pdf, extract_images_from_pdf, create_enhanced_thumbnails

# Configure logging
logging.basicConfig(level=logging.INFO, 
                    format='%(message)s')

# Add the project directory to the Python path
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# Import our utility function
from core.pptx_utils import extract_slide_image_fallback

def test_with_file(pptx_path):
    """Test thumbnail generation with a local PPTX file"""
    print(f"Testing with file: {pptx_path}")
    
    # Check if file exists
    if not os.path.exists(pptx_path):
        print(f"Error: File {pptx_path} not found")
        return
    
    try:
        # Open the PowerPoint file
        presentation = Presentation(pptx_path)
        print(f"Successfully opened presentation with {len(presentation.slides)} slides")
        
        # Process each slide
        for i in range(len(presentation.slides)):
            try:
                print(f"Processing slide {i+1}...")
                # Try to extract an image from the slide
                thumbnail = extract_slide_image_fallback(presentation, i)
                
                # Save the thumbnail
                output_path = f"slide_{i+1}_thumbnail.png"
                thumbnail.save(output_path)
                print(f"Saved thumbnail to {output_path}")
            except Exception as e:
                print(f"Error creating thumbnail for slide {i+1}: {str(e)}")
    
    except Exception as e:
        print(f"Error processing PowerPoint file: {str(e)}")

def main():
    """Test PowerPoint thumbnail generation"""
    print("Checking for external tools...")
    tools = check_external_tools()
    
    if tools["libreoffice"]:
        print(f"LibreOffice: Available at {tools['libreoffice']}")
    else:
        print("LibreOffice: Not found")
    
    if tools["poppler"]:
        print("pdf2image/poppler: Available")
    else:
        print("pdf2image/poppler: Not available")
    
    # Find a PPTX file to test
    templates_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "media", "templates")
    print(f"Looking for PPTX files in {templates_dir}...")
    
    pptx_files = []
    for root, dirs, files in os.walk(templates_dir):
        for file in files:
            if file.endswith(".pptx"):
                pptx_files.append(os.path.join(root, file))
    
    if not pptx_files:
        print("No PPTX files found in the templates directory.")
        return
    
    # Use the first PPTX file found
    test_pptx = pptx_files[0]
    print(f"Found PPTX file: {test_pptx}")
    
    # Open the file for testing
    with open(test_pptx, 'rb') as f:
        # Create a file-like object that mimics a Django uploaded file
        class MockFile:
            def __init__(self, file_obj):
                self.file = file_obj
                
            def chunks(self):
                self.file.seek(0)
                yield self.file.read()
        
        mock_file = MockFile(f)
        
        # Generate thumbnails
        print("Generating thumbnails...")
        thumbnails = create_enhanced_thumbnails(mock_file)
        
        print(f"Generated {len(thumbnails)} thumbnails")
        
        # Save the thumbnails
        for i, thumbnail in enumerate(thumbnails):
            output_path = f"slide_{i+1}_thumbnail.png"
            thumbnail.save(output_path)
            print(f"Saved thumbnail: {output_path}")
            
            # Print details about slide 2 if this is slide 2
            if i == 1:  # 0-indexed, so 1 is slide 2
                print(f"\nSlide 2 details:")
                print(f"Size: {thumbnail.size}")
                print(f"Mode: {thumbnail.mode}")
                # Sample some pixels to check colors
                center_pixel = thumbnail.getpixel((thumbnail.width // 2, thumbnail.height // 2))
                print(f"Center pixel color: {center_pixel}")
                
                # Check for blue background (should be around RGB(0, 114, 198))
                blue_count = 0
                sample_points = [
                    (50, 50), 
                    (thumbnail.width - 50, 50),
                    (50, thumbnail.height - 50),
                    (thumbnail.width - 50, thumbnail.height - 50)
                ]
                
                for x, y in sample_points:
                    pixel = thumbnail.getpixel((x, y))
                    print(f"Pixel at ({x}, {y}): {pixel}")
                    # Check if pixel is blueish
                    if isinstance(pixel, tuple) and len(pixel) >= 3:
                        r, g, b = pixel[:3]
                        if r < 50 and g > 50 and b > 150:  # Rough check for blue
                            blue_count += 1
                
                print(f"Blue background detection: {blue_count} out of {len(sample_points)} sample points are blue")

if __name__ == "__main__":
    main() 