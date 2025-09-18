"""
PowerPoint Generation Service for Report Center
Handles native PowerPoint (.pptx) export with charts, branding, and professional layouts
"""

import os
import io
import json
import base64
from datetime import datetime
from typing import Dict, List, Any, Optional, Tuple
import logging

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image
import matplotlib
matplotlib.use('Agg')  # Use non-GUI backend
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from io import BytesIO

from django.conf import settings
from django.core.files.base import ContentFile

from .chart_export_service import ChartExportService
from .scenario_integration_service import ScenarioDataIntegrationService

logger = logging.getLogger(__name__)


class PowerPointGeneratorService:
    """Service for generating professional PowerPoint presentations"""
    
    def __init__(self):
        self.chart_service = ChartExportService()
        self.data_service = ScenarioDataIntegrationService()
        
        # PowerPoint styling constants
        self.TITLE_FONT_SIZE = Pt(28)
        self.SUBTITLE_FONT_SIZE = Pt(18)
        self.BODY_FONT_SIZE = Pt(14)
        self.CAPTION_FONT_SIZE = Pt(10)
        
        # Color scheme (can be customized per advisor)
        self.PRIMARY_COLOR = RGBColor(0, 114, 198)  # #0072C6
        self.SECONDARY_COLOR = RGBColor(68, 68, 68)  # #444444
        self.ACCENT_COLOR = RGBColor(255, 127, 39)  # #FF7F27
    
    def generate_powerpoint_report(
        self,
        sections: List[Dict[str, Any]],
        data: Dict[str, Any],
        options: Dict[str, Any] = None
    ) -> Tuple[bytes, str]:
        """
        Generate a complete PowerPoint presentation from sections and data
        
        Args:
            sections: List of section configurations
            data: Report data including client, scenario, and results
            options: Generation options (template, branding, etc.)
            
        Returns:
            Tuple of (file_content_bytes, filename)
        """
        try:
            # Create new presentation
            prs = Presentation()
            
            # Apply branding if provided
            if options and 'branding' in options:
                self._apply_branding(prs, options['branding'])
            
            # Generate slides for each section
            for section in sections:
                self._generate_slide_for_section(prs, section, data, options)
            
            # Save to BytesIO
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            # Generate filename
            client_name = data.get('client', {}).get('last_name', 'Report')
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"{client_name}_Report_{timestamp}.pptx"
            
            return buffer.getvalue(), filename
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint report: {str(e)}")
            raise
    
    def _generate_slide_for_section(
        self,
        prs: Presentation,
        section: Dict[str, Any],
        data: Dict[str, Any],
        options: Dict[str, Any]
    ):
        """Generate a slide for a specific section"""
        section_type = section.get('type', '')
        section_config = section.get('config', {})
        
        if section_type == 'cover':
            self._create_cover_slide(prs, section_config, data)
        elif section_type == 'executive_summary':
            self._create_executive_summary_slide(prs, section_config, data)
        elif section_type == 'scenario_overview':
            self._create_scenario_overview_slide(prs, section_config, data)
        elif section_type == 'financial_timeline':
            self._create_financial_timeline_slide(prs, section_config, data)
        elif section_type == 'asset_allocation':
            self._create_asset_allocation_slide(prs, section_config, data)
        elif section_type == 'monte_carlo':
            self._create_monte_carlo_slide(prs, section_config, data)
        elif section_type == 'tax_analysis':
            self._create_tax_analysis_slide(prs, section_config, data)
        elif section_type == 'irmaa_analysis':
            self._create_irmaa_analysis_slide(prs, section_config, data)
        elif section_type == 'scenario_comparison':
            self._create_scenario_comparison_slide(prs, section_config, data)
        elif section_type == 'recommendations':
            self._create_recommendations_slide(prs, section_config, data)
        else:
            # Generic text slide for unknown types
            self._create_text_slide(prs, section_config, data)
    
    def _create_cover_slide(self, prs: Presentation, config: Dict, data: Dict):
        """Create cover slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_placeholder = slide.shapes.title
        title_placeholder.text = config.get('title', f"Retirement Analysis for {data['client']['first_name']} {data['client']['last_name']}")
        
        # Subtitle
        subtitle_placeholder = slide.placeholders[1]
        subtitle_placeholder.text = f"Prepared by {data.get('advisor', {}).get('name', 'Your Advisor')}\n{datetime.now().strftime('%B %Y')}"
        
        # Add logo if available
        if data.get('advisor', {}).get('branding', {}).get('logo_url'):
            try:
                self._add_logo_to_slide(slide, data['advisor']['branding']['logo_url'])
            except Exception as e:
                logger.warning(f"Could not add logo to cover slide: {str(e)}")
    
    def _create_executive_summary_slide(self, prs: Presentation, config: Dict, data: Dict):
        """Create executive summary slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = "Executive Summary"
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Get scenario summary data
        scenario_data = data.get('scenario_data', {})
        scenario_results = data.get('scenario_results', {})
        
        # Build executive summary points
        summary_points = self._build_executive_summary_points(scenario_data, scenario_results, data)
        
        for point in summary_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = self.BODY_FONT_SIZE
            p.level = 0
    
    def _create_scenario_overview_slide(self, prs: Presentation, config: Dict, data: Dict):
        """Create scenario overview slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "Scenario Overview"
        
        # Create content table
        rows, cols = 6, 2
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Populate table with scenario data
        scenario_data = data.get('scenario_data', {})
        
        table.cell(0, 0).text = "Current Age"
        table.cell(0, 1).text = str(scenario_data.get('current_age', 'N/A'))
        
        table.cell(1, 0).text = "Retirement Age"
        table.cell(1, 1).text = str(scenario_data.get('retirement_age', 'N/A'))
        
        table.cell(2, 0).text = "Current Assets"
        table.cell(2, 1).text = f"${scenario_data.get('total_assets', 0):,.0f}"
        
        table.cell(3, 0).text = "Annual Income Needed"
        table.cell(3, 1).text = f"${scenario_data.get('annual_income_needed', 0):,.0f}"
        
        table.cell(4, 0).text = "Social Security"
        table.cell(4, 1).text = f"${scenario_data.get('social_security_benefit', 0):,.0f}/year"
        
        table.cell(5, 0).text = "Tax Status"
        table.cell(5, 1).text = scenario_data.get('tax_status', 'N/A')
    
    def _create_financial_timeline_slide(self, prs: Presentation, config: Dict, data: Dict):
        """Create financial timeline chart slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "Financial Timeline Projection"
        title_frame.paragraphs[0].font.size = self.TITLE_FONT_SIZE
        
        # Generate timeline chart
        chart_image = self._create_timeline_chart(data.get('scenario_results', []))
        
        if chart_image:
            # Add chart image to slide
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            slide.shapes.add_picture(chart_image, left, top, width, height)
    
    def _create_asset_allocation_slide(self, prs: Presentation, config: Dict, data: Dict):
        """Create asset allocation pie chart slide"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "Asset Allocation"
        title_frame.paragraphs[0].font.size = self.TITLE_FONT_SIZE
        
        # Create pie chart
        chart_data = CategoryChartData()
        
        # Get asset allocation data
        assets = data.get('assets', [])
        if assets:
            categories = []
            values = []
            
            for asset in assets:
                if asset.get('current_balance', 0) > 0:
                    categories.append(asset.get('investment_name', 'Unknown'))
                    values.append(float(asset.get('current_balance', 0)))
            
            if categories and values:
                chart_data.categories = categories
                chart_data.add_series('Asset Values', values)
                
                # Add chart to slide
                left = Inches(2)
                top = Inches(2)
                width = Inches(6)
                height = Inches(4.5)
                
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.PIE, left, top, width, height, chart_data
                ).chart
                
                chart.has_legend = True
                chart.legend.position = 2  # Right side
    
    def _create_monte_carlo_slide(self, prs: Presentation, config: Dict, data: Dict):
        """Create Monte Carlo analysis slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "Monte Carlo Analysis"
        
        # Generate Monte Carlo visualization
        monte_carlo_image = self._create_monte_carlo_chart(data.get('scenario_results', []))
        
        if monte_carlo_image:
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            slide.shapes.add_picture(monte_carlo_image, left, top, width, height)
        
        # Add success probability text
        success_rate = self._calculate_success_probability(data.get('scenario_results', []))
        
        text_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
        text_frame = text_box.text_frame
        text_frame.text = f"Success Probability: {success_rate:.1f}%"
        text_frame.paragraphs[0].font.size = Pt(20)
        text_frame.paragraphs[0].font.bold = True
    
    def _create_tax_analysis_slide(self, prs: Presentation, config: Dict, data: Dict):
        """Create tax analysis slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "Tax Analysis"
        
        # Create tax breakdown content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        scenario_results = data.get('scenario_results', [])
        if scenario_results:
            # Get tax information from first result
            result = scenario_results[0]
            
            p1 = tf.add_paragraph()
            p1.text = f"Federal Tax Rate: {result.get('federal_tax_rate', 0):.1f}%"
            p1.font.size = self.BODY_FONT_SIZE
            
            p2 = tf.add_paragraph()
            p2.text = f"State Tax Rate: {result.get('state_tax_rate', 0):.1f}%"
            p2.font.size = self.BODY_FONT_SIZE
            
            p3 = tf.add_paragraph()
            p3.text = f"Annual Tax Liability: ${result.get('total_taxes', 0):,.0f}"
            p3.font.size = self.BODY_FONT_SIZE
            
            # IRMAA information if applicable
            if result.get('irmaa_bracket_number', 0) > 0:
                p4 = tf.add_paragraph()
                p4.text = f"IRMAA Bracket: {result.get('irmaa_bracket_number', 0)}"
                p4.font.size = self.BODY_FONT_SIZE
                
                p5 = tf.add_paragraph()
                p5.text = f"Additional Medicare Premium: ${result.get('irmaa_premium', 0):,.0f}/year"
                p5.font.size = self.BODY_FONT_SIZE
    
    def _create_irmaa_analysis_slide(self, prs: Presentation, config: Dict, data: Dict):
        """Create IRMAA-specific analysis slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "IRMAA Impact Analysis"
        
        # Generate IRMAA chart
        irmaa_image = self._create_irmaa_chart(data.get('scenario_results', []))
        
        if irmaa_image:
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            slide.shapes.add_picture(irmaa_image, left, top, width, height)
    
    def _create_scenario_comparison_slide(self, prs: Presentation, config: Dict, data: Dict):
        """Create scenario comparison slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "Scenario Comparison"
        
        # Create comparison chart
        comparison_image = self._create_comparison_chart(data)
        
        if comparison_image:
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            slide.shapes.add_picture(comparison_image, left, top, width, height)
    
    def _create_recommendations_slide(self, prs: Presentation, config: Dict, data: Dict):
        """Create recommendations slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "Recommendations"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Generate recommendations based on data
        recommendations = self._generate_recommendations(data)
        
        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = rec
            p.font.size = self.BODY_FONT_SIZE
            p.level = 0
    
    def _create_text_slide(self, prs: Presentation, config: Dict, data: Dict):
        """Create generic text slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = config.get('title', 'Information')
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = config.get('content', 'Content placeholder')
    
    def _apply_branding(self, prs: Presentation, branding: Dict):
        """Apply advisor branding to presentation"""
        try:
            # Update color scheme if provided
            if branding.get('primary_color'):
                # Convert hex to RGB
                hex_color = branding['primary_color'].lstrip('#')
                rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                self.PRIMARY_COLOR = RGBColor(*rgb)
            
            # Additional branding customizations can be added here
            
        except Exception as e:
            logger.warning(f"Could not apply branding: {str(e)}")
    
    def _create_timeline_chart(self, scenario_results: List[Dict]) -> Optional[BytesIO]:
        """Create financial timeline chart as image"""
        if not scenario_results:
            return None
            
        try:
            fig, ax = plt.subplots(figsize=(10, 6))
            
            years = [result.get('year', 0) for result in scenario_results]
            portfolio_values = [result.get('end_portfolio_value', 0) for result in scenario_results]
            
            ax.plot(years, portfolio_values, linewidth=3, color='#0072C6')
            ax.fill_between(years, portfolio_values, alpha=0.3, color='#0072C6')
            
            ax.set_xlabel('Year')
            ax.set_ylabel('Portfolio Value ($)')
            ax.set_title('Portfolio Value Over Time')
            ax.grid(True, alpha=0.3)
            
            # Format y-axis as currency
            ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.0f}'))
            
            plt.tight_layout()
            
            # Save to BytesIO
            img_buffer = BytesIO()
            plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            return img_buffer
            
        except Exception as e:
            logger.error(f"Error creating timeline chart: {str(e)}")
            return None
    
    def _create_monte_carlo_chart(self, scenario_results: List[Dict]) -> Optional[BytesIO]:
        """Create Monte Carlo analysis visualization"""
        if not scenario_results:
            return None
            
        try:
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # Simulate Monte Carlo results (in real implementation, use actual MC data)
            years = list(range(len(scenario_results)))
            
            # Create multiple scenario lines
            for i in range(10):  # Show 10 scenarios
                scenario_values = []
                for result in scenario_results:
                    base_value = result.get('end_portfolio_value', 0)
                    # Add some randomness for illustration
                    import random
                    multiplier = random.uniform(0.7, 1.3)
                    scenario_values.append(base_value * multiplier)
                
                ax.plot(years, scenario_values, alpha=0.3, color='gray', linewidth=1)
            
            # Add median line
            median_values = [result.get('end_portfolio_value', 0) for result in scenario_results]
            ax.plot(years, median_values, linewidth=3, color='#0072C6', label='Median Scenario')
            
            ax.set_xlabel('Years into Retirement')
            ax.set_ylabel('Portfolio Value ($)')
            ax.set_title('Monte Carlo Analysis - Portfolio Projections')
            ax.legend()
            ax.grid(True, alpha=0.3)
            
            plt.tight_layout()
            
            img_buffer = BytesIO()
            plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            return img_buffer
            
        except Exception as e:
            logger.error(f"Error creating Monte Carlo chart: {str(e)}")
            return None
    
    def _create_irmaa_chart(self, scenario_results: List[Dict]) -> Optional[BytesIO]:
        """Create IRMAA impact visualization"""
        if not scenario_results:
            return None
            
        try:
            fig, ax = plt.subplots(figsize=(10, 6))
            
            years = [result.get('year', 0) for result in scenario_results]
            irmaa_premiums = [result.get('irmaa_premium', 0) for result in scenario_results]
            
            # Create bar chart for IRMAA premiums
            bars = ax.bar(years, irmaa_premiums, color='#FF7F27', alpha=0.7)
            
            # Add threshold line if available
            threshold = scenario_results[0].get('irmaa_bracket_threshold', 0) if scenario_results else 0
            if threshold > 0:
                ax.axhline(y=threshold, color='red', linestyle='--', 
                          label=f'IRMAA Threshold: ${threshold:,.0f}')
                ax.legend()
            
            ax.set_xlabel('Year')
            ax.set_ylabel('Annual IRMAA Premium ($)')
            ax.set_title('IRMAA Premium Impact Over Time')
            ax.grid(True, alpha=0.3)
            
            plt.tight_layout()
            
            img_buffer = BytesIO()
            plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            return img_buffer
            
        except Exception as e:
            logger.error(f"Error creating IRMAA chart: {str(e)}")
            return None
    
    def _create_comparison_chart(self, data: Dict) -> Optional[BytesIO]:
        """Create scenario comparison visualization"""
        try:
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # This would compare multiple scenarios if available
            # For now, create a simple comparison placeholder
            
            scenarios = ['Current Plan', 'Optimized Plan']
            success_rates = [75, 85]  # Placeholder data
            
            bars = ax.bar(scenarios, success_rates, color=['#0072C6', '#00B04F'])
            
            # Add value labels on bars
            for bar in bars:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height,
                       f'{height}%', ha='center', va='bottom')
            
            ax.set_ylabel('Success Rate (%)')
            ax.set_title('Scenario Comparison - Success Probability')
            ax.set_ylim(0, 100)
            
            plt.tight_layout()
            
            img_buffer = BytesIO()
            plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            return img_buffer
            
        except Exception as e:
            logger.error(f"Error creating comparison chart: {str(e)}")
            return None
    
    def _build_executive_summary_points(self, scenario_data: Dict, scenario_results: List[Dict], data: Dict) -> List[str]:
        """Build executive summary bullet points"""
        points = []
        
        client_name = data.get('client', {}).get('first_name', 'Client')
        current_age = scenario_data.get('current_age', 0)
        retirement_age = scenario_data.get('retirement_age', 65)
        
        points.append(f"{client_name} is {current_age} years old and plans to retire at age {retirement_age}")
        
        total_assets = scenario_data.get('total_assets', 0)
        if total_assets > 0:
            points.append(f"Current total assets: ${total_assets:,.0f}")
        
        if scenario_results:
            success_rate = self._calculate_success_probability(scenario_results)
            points.append(f"Retirement plan success probability: {success_rate:.1f}%")
            
            # Check for IRMAA impact
            has_irmaa = any(result.get('irmaa_bracket_number', 0) > 0 for result in scenario_results)
            if has_irmaa:
                points.append("IRMAA thresholds may impact Medicare premiums")
        
        return points
    
    def _generate_recommendations(self, data: Dict) -> List[str]:
        """Generate recommendations based on scenario analysis"""
        recommendations = []
        
        scenario_results = data.get('scenario_results', [])
        if not scenario_results:
            return ["Complete scenario analysis to generate recommendations"]
        
        success_rate = self._calculate_success_probability(scenario_results)
        
        if success_rate < 70:
            recommendations.append("Consider increasing retirement savings rate")
            recommendations.append("Explore delaying retirement by 1-2 years")
        
        # Check for IRMAA impact
        has_irmaa = any(result.get('irmaa_bracket_number', 0) > 0 for result in scenario_results)
        if has_irmaa:
            recommendations.append("Evaluate Roth conversion strategies to manage IRMAA impact")
            recommendations.append("Consider tax-loss harvesting in taxable accounts")
        
        # Check for high tax burden
        avg_tax_rate = sum(result.get('federal_tax_rate', 0) + result.get('state_tax_rate', 0) 
                          for result in scenario_results) / len(scenario_results)
        if avg_tax_rate > 25:
            recommendations.append("Explore tax-deferred savings opportunities")
            recommendations.append("Review municipal bond allocation for tax efficiency")
        
        return recommendations
    
    def _calculate_success_probability(self, scenario_results: List[Dict]) -> float:
        """Calculate simple success probability (placeholder for Monte Carlo)"""
        if not scenario_results:
            return 0.0
        
        # Simple heuristic: success if portfolio doesn't go to zero
        successful_years = sum(1 for result in scenario_results 
                             if result.get('end_portfolio_value', 0) > 0)
        
        return (successful_years / len(scenario_results)) * 100
    
    def _add_logo_to_slide(self, slide, logo_url: str):
        """Add advisor logo to slide"""
        try:
            # In a real implementation, you would download and add the logo
            # This is a placeholder for the logo functionality
            pass
        except Exception as e:
            logger.warning(f"Could not add logo: {str(e)}")