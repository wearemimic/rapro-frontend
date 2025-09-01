"""
Report Generator Service
Handles PDF and PowerPoint report generation using ReportLab and python-pptx
"""

import os
import io
import json
import logging
from datetime import datetime
from typing import Dict, List, Optional, Any
from decimal import Decimal

from django.conf import settings
from django.utils import timezone
from django.template.loader import render_to_string

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak
from reportlab.platypus.flowables import HRFlowable
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.linecharts import HorizontalLineChart
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.piecharts import Pie
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ..models import Report, ReportTemplate, ReportSection
from .scenario_integration_service import ScenarioDataIntegrationService
from .template_processor import TemplateProcessor
from .chart_export_service import ChartExportService
from .file_storage_service import ReportFileStorageService
from .powerpoint_generator import PowerPointGeneratorService

logger = logging.getLogger(__name__)


class ReportGeneratorService:
    """
    Main service for generating PDF and PowerPoint reports from templates and scenario data
    """
    
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()
        self.storage_config = getattr(settings, 'REPORT_CENTER_STORAGE', {})
        
        # Initialize services
        self.scenario_service = ScenarioDataIntegrationService()
        self.template_processor = TemplateProcessor()
        self.chart_service = ChartExportService()
        self.storage_service = ReportFileStorageService()
        self.powerpoint_service = PowerPointGeneratorService()
    
    def _setup_custom_styles(self):
        """Setup custom ReportLab styles for professional reports"""
        
        # Title style
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Title'],
            fontSize=24,
            spaceAfter=30,
            textColor=colors.HexColor('#2c3e50'),
            alignment=TA_CENTER
        ))
        
        # Section heading style
        self.styles.add(ParagraphStyle(
            name='SectionHeading',
            parent=self.styles['Heading1'],
            fontSize=16,
            spaceBefore=20,
            spaceAfter=12,
            textColor=colors.HexColor('#34495e'),
            borderWidth=0,
            borderColor=colors.HexColor('#3498db'),
            borderPadding=8
        ))
        
        # Financial data style
        self.styles.add(ParagraphStyle(
            name='FinancialData',
            parent=self.styles['Normal'],
            fontSize=11,
            fontName='Helvetica-Bold',
            textColor=colors.HexColor('#27ae60'),
            alignment=TA_RIGHT
        ))
        
        # Disclaimer style
        self.styles.add(ParagraphStyle(
            name='Disclaimer',
            parent=self.styles['Normal'],
            fontSize=8,
            textColor=colors.grey,
            alignment=TA_CENTER,
            spaceBefore=20
        ))
    
    def generate_pdf_report(self, report: Report) -> Dict[str, Any]:
        """
        Generate PDF report from Report object using integrated services
        """
        try:
            # Process template and get structured content
            processed_content = self.template_processor.process_template(report)
            
            # Get scenario data
            scenario_data = self.scenario_service.get_report_data(report)
            
            # Generate temporary file path
            pdf_dir = self.storage_config.get('PDF_REPORTS', 
                os.path.join(settings.MEDIA_ROOT, 'report_center/generated_reports/pdf/'))
            os.makedirs(pdf_dir, exist_ok=True)
            
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            temp_filename = f"temp_report_{report.id}_{timestamp}.pdf"
            temp_file_path = os.path.join(pdf_dir, temp_filename)
            
            # Create PDF document
            doc = SimpleDocTemplate(
                temp_file_path,
                pagesize=letter,
                rightMargin=inch,
                leftMargin=inch,
                topMargin=inch,
                bottomMargin=inch
            )
            
            # Build content using processed template data
            story = []
            
            # Add header section
            self._add_report_header(story, report)
            
            # Process sections from template
            for section_data in processed_content.get('sections', []):
                self._add_processed_section(story, section_data, scenario_data)
            
            # Add footer/disclaimer
            self._add_report_footer(story, report)
            
            # Build PDF
            doc.build(story)
            
            # Store file using storage service
            storage_result = self.storage_service.store_generated_report(
                report, temp_file_path, 'pdf'
            )
            
            if storage_result['success']:
                # Clean up temporary file
                try:
                    os.remove(temp_file_path)
                except OSError:
                    pass
                
                logger.info(f"PDF report generated and stored successfully for report {report.id}")
                return storage_result
            else:
                logger.error(f"PDF storage failed for report {report.id}: {storage_result.get('error')}")
                return storage_result
            
        except Exception as e:
            logger.error(f"PDF generation failed for report {report.id}: {str(e)}")
            return {
                'success': False,
                'error': str(e),
                'format': 'pdf'
            }
    
    def _add_report_header(self, story: List, report: Report):
        """Add report header with title and branding"""
        
        # Report title
        title = Paragraph(report.title or f"Retirement Planning Report", self.styles['CustomTitle'])
        story.append(title)
        story.append(Spacer(1, 12))
        
        # Client and date info
        header_data = [
            ['Client:', report.client.full_name if report.client else 'N/A'],
            ['Scenario:', report.scenario.name if report.scenario else 'N/A'],
            ['Generated:', timezone.now().strftime('%B %d, %Y')],
            ['Advisor:', report.created_by.get_full_name()]
        ]
        
        header_table = Table(header_data, colWidths=[1.5*inch, 4*inch])
        header_table.setStyle(TableStyle([
            ('FONT', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('TEXTCOLOR', (0, 0), (0, -1), colors.HexColor('#7f8c8d')),
            ('ALIGN', (0, 0), (0, -1), 'RIGHT'),
            ('ALIGN', (1, 0), (1, -1), 'LEFT'),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ]))
        
        story.append(header_table)
        story.append(Spacer(1, 20))
        story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor('#bdc3c7')))
        story.append(Spacer(1, 20))
    
    def _add_client_section(self, story: List, report: Report):
        """Add client information section"""
        if not report.client:
            return
        
        client = report.client
        
        # Section heading
        heading = Paragraph("Client Overview", self.styles['SectionHeading'])
        story.append(heading)
        
        # Client details table
        client_data = [
            ['Name:', client.full_name],
            ['Age:', f"{client.age} years old" if client.age else 'Not specified'],
            ['Retirement Age:', f"{client.retirement_age} years old" if client.retirement_age else 'Not specified'],
            ['Current Income:', f"${client.current_income:,.0f}" if client.current_income else 'Not specified'],
            ['Current Savings:', f"${client.current_savings:,.0f}" if client.current_savings else 'Not specified'],
        ]
        
        client_table = Table(client_data, colWidths=[2*inch, 4*inch])
        client_table.setStyle(TableStyle([
            ('FONT', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 11),
            ('TEXTCOLOR', (0, 0), (0, -1), colors.HexColor('#34495e')),
            ('ALIGN', (0, 0), (0, -1), 'RIGHT'),
            ('ALIGN', (1, 0), (1, -1), 'LEFT'),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('ROWBACKGROUNDS', (0, 0), (-1, -1), [colors.white, colors.HexColor('#f8f9fa')]),
        ]))
        
        story.append(client_table)
        story.append(Spacer(1, 20))
    
    def _add_scenario_section(self, story: List, report: Report):
        """Add scenario information section"""
        if not report.scenario:
            return
        
        scenario = report.scenario
        
        # Section heading
        heading = Paragraph("Scenario Analysis", self.styles['SectionHeading'])
        story.append(heading)
        
        # Scenario overview
        overview_text = f"""
        This analysis examines the retirement planning scenario "{scenario.name}" based on current 
        financial data and projected outcomes. The calculations include Monte Carlo simulations, 
        tax considerations, and various retirement income strategies.
        """
        
        overview = Paragraph(overview_text, self.styles['Normal'])
        story.append(overview)
        story.append(Spacer(1, 15))
        
        # Key scenario parameters
        scenario_data = [
            ['Scenario Name:', scenario.name],
            ['Retirement Age:', f"{scenario.retirement_age} years" if scenario.retirement_age else 'Not set'],
            ['Final Age:', f"{scenario.final_age} years" if scenario.final_age else 'Not set'],
            ['Current Year:', str(scenario.current_year) if scenario.current_year else 'Not set'],
        ]
        
        scenario_table = Table(scenario_data, colWidths=[2*inch, 4*inch])
        scenario_table.setStyle(TableStyle([
            ('FONT', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 11),
            ('TEXTCOLOR', (0, 0), (0, -1), colors.HexColor('#34495e')),
            ('ALIGN', (0, 0), (0, -1), 'RIGHT'),
            ('ALIGN', (1, 0), (1, -1), 'LEFT'),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('ROWBACKGROUNDS', (0, 0), (-1, -1), [colors.white, colors.HexColor('#f8f9fa')]),
        ]))
        
        story.append(scenario_table)
        story.append(Spacer(1, 20))
    
    def _add_processed_section(self, story: List, section_data: Dict, scenario_data: Dict):
        """Add section content from processed template data"""
        
        # Section title
        if section_data.get('title'):
            title = Paragraph(section_data['title'], self.styles['SectionHeading'])
            story.append(title)
        
        # Add content based on section type
        content = section_data.get('content', {})
        
        # Add text content
        if 'summary_text' in content:
            para = Paragraph(content['summary_text'].strip(), self.styles['Normal'])
            story.append(para)
            story.append(Spacer(1, 12))
        
        # Add key findings as bullet points
        if 'key_findings' in content:
            for finding in content['key_findings']:
                bullet_text = f"• {finding}"
                bullet_para = Paragraph(bullet_text, self.styles['Normal'])
                story.append(bullet_para)
                story.append(Spacer(1, 6))
        
        # Add recommendations
        if 'recommendations' in content:
            recommendations = content['recommendations']
            if isinstance(recommendations, list):
                for i, rec in enumerate(recommendations, 1):
                    if isinstance(rec, dict):
                        rec_text = f"{i}. {rec.get('recommendation', 'No recommendation text')}"
                    else:
                        rec_text = f"{i}. {rec}"
                    rec_para = Paragraph(rec_text, self.styles['Normal'])
                    story.append(rec_para)
                    story.append(Spacer(1, 8))
        
        # Generate and add charts for this section
        charts = section_data.get('charts', [])
        for chart_config in charts:
            try:
                chart_result = self.chart_service.export_scenario_chart(
                    chart_config.get('data', {}),
                    chart_config.get('chart_type', 'generic'),
                    chart_config.get('chart_id', 'chart')
                )
                
                if chart_result['success'] and os.path.exists(chart_result['image_path']):
                    # Add chart to PDF
                    chart_image = Image(chart_result['image_path'], width=6*inch, height=3.6*inch)
                    story.append(chart_image)
                    story.append(Spacer(1, 12))
                    
                    # Add chart title if specified
                    if chart_config.get('title'):
                        chart_title = Paragraph(chart_config['title'], self.styles['Normal'])
                        story.append(chart_title)
                        story.append(Spacer(1, 12))
                        
            except Exception as e:
                logger.warning(f"Failed to add chart {chart_config.get('chart_id')}: {str(e)}")
        
        story.append(Spacer(1, 20))
    
    def _add_report_sections(self, story: List, report: Report):
        """Add dynamic report sections based on template configuration (legacy method)"""
        
        # Get sections ordered by their order field
        sections = report.sections.filter(is_enabled=True).order_by('order')
        
        for section in sections:
            self._add_section_content(story, section, report)
    
    def _add_section_content(self, story: List, section: ReportSection, report: Report):
        """Add content for a specific report section"""
        
        # Section title
        if section.title:
            title = Paragraph(section.title, self.styles['SectionHeading'])
            story.append(title)
        
        # Process section based on type
        section_type = section.section_type
        content_config = section.content_config or {}
        
        if section_type == 'executive_summary':
            self._add_executive_summary(story, report, content_config)
        elif section_type == 'financial_projections':
            self._add_financial_projections(story, report, content_config)
        elif section_type == 'monte_carlo_analysis':
            self._add_monte_carlo_analysis(story, report, content_config)
        elif section_type == 'tax_analysis':
            self._add_tax_analysis(story, report, content_config)
        elif section_type == 'recommendations':
            self._add_recommendations(story, report, content_config)
        else:
            # Generic text section
            self._add_generic_section(story, section, content_config)
        
        story.append(Spacer(1, 20))
    
    def _add_executive_summary(self, story: List, report: Report, config: Dict):
        """Add executive summary section"""
        
        summary_text = f"""
        Based on our analysis of {report.client.full_name if report.client else 'your'} retirement planning scenario, 
        we have identified key opportunities and recommendations to optimize your retirement income strategy.
        
        This report presents a comprehensive analysis of your current financial position, projected outcomes,
        and specific action items to enhance your retirement security.
        """
        
        summary = Paragraph(summary_text, self.styles['Normal'])
        story.append(summary)
    
    def _add_financial_projections(self, story: List, report: Report, config: Dict):
        """Add financial projections section with tables and charts"""
        
        # Add placeholder financial data table
        projections_data = [
            ['Year', 'Age', 'Assets', 'Income', 'Expenses'],
            ['2024', '65', '$1,250,000', '$85,000', '$75,000'],
            ['2029', '70', '$1,450,000', '$95,000', '$80,000'],
            ['2034', '75', '$1,580,000', '$105,000', '$85,000'],
            ['2039', '80', '$1,620,000', '$110,000', '$90,000'],
        ]
        
        projections_table = Table(projections_data, colWidths=[0.8*inch, 0.8*inch, 1.2*inch, 1.0*inch, 1.0*inch])
        projections_table.setStyle(TableStyle([
            # Header row
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3498db')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONT', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            
            # Data rows
            ('FONT', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f8f9fa')]),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#dee2e6')),
        ]))
        
        story.append(projections_table)
    
    def _add_monte_carlo_analysis(self, story: List, report: Report, config: Dict):
        """Add Monte Carlo simulation results"""
        
        analysis_text = f"""
        Monte Carlo analysis provides a statistical view of potential outcomes by running thousands 
        of simulations with varying market conditions. This analysis helps understand the probability 
        of achieving retirement goals under different scenarios.
        """
        
        analysis = Paragraph(analysis_text, self.styles['Normal'])
        story.append(analysis)
        story.append(Spacer(1, 12))
        
        # Monte Carlo results table
        mc_data = [
            ['Confidence Level', 'Success Probability', 'Final Portfolio Value'],
            ['10th Percentile', '90%', '$850,000'],
            ['25th Percentile', '75%', '$1,150,000'],
            ['50th Percentile (Median)', '50%', '$1,450,000'],
            ['75th Percentile', '25%', '$1,850,000'],
            ['90th Percentile', '10%', '$2,350,000'],
        ]
        
        mc_table = Table(mc_data, colWidths=[2*inch, 1.5*inch, 1.5*inch])
        mc_table.setStyle(TableStyle([
            # Header row
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#27ae60')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONT', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            
            # Data rows
            ('FONT', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f8f9fa')]),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#dee2e6')),
        ]))
        
        story.append(mc_table)
    
    def _add_tax_analysis(self, story: List, report: Report, config: Dict):
        """Add tax analysis section including IRMAA and Roth conversions"""
        
        tax_text = f"""
        Tax planning is a critical component of retirement strategy. This analysis examines 
        federal and state tax implications, Medicare IRMAA brackets, and Roth conversion opportunities 
        to optimize after-tax retirement income.
        """
        
        tax_para = Paragraph(tax_text, self.styles['Normal'])
        story.append(tax_para)
        story.append(Spacer(1, 12))
        
        # Tax analysis table
        tax_data = [
            ['Tax Component', 'Current Impact', 'Projected Impact'],
            ['Federal Income Tax', '$15,500', '$12,200'],
            ['State Income Tax', '$4,200', '$3,800'],
            ['IRMAA Surcharge', '$0', '$1,680'],
            ['Total Tax Burden', '$19,700', '$17,680'],
        ]
        
        tax_table = Table(tax_data, colWidths=[2*inch, 1.5*inch, 1.5*inch])
        tax_table.setStyle(TableStyle([
            # Header row
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#e74c3c')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONT', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            
            # Data rows
            ('FONT', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f8f9fa')]),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#dee2e6')),
        ]))
        
        story.append(tax_table)
    
    def _add_recommendations(self, story: List, report: Report, config: Dict):
        """Add recommendations section"""
        
        recommendations = [
            "Consider increasing 401(k) contributions to maximize employer matching",
            "Evaluate Roth conversion opportunities during lower-income years", 
            "Review asset allocation to ensure appropriate risk level for time horizon",
            "Monitor IRMAA brackets to optimize Medicare premium costs",
            "Implement tax-loss harvesting strategies in taxable accounts"
        ]
        
        for i, rec in enumerate(recommendations, 1):
            rec_text = f"{i}. {rec}"
            rec_para = Paragraph(rec_text, self.styles['Normal'])
            story.append(rec_para)
            story.append(Spacer(1, 8))
    
    def _add_generic_section(self, story: List, section: ReportSection, config: Dict):
        """Add generic text section"""
        
        content = config.get('content', 'Section content not configured.')
        para = Paragraph(content, self.styles['Normal'])
        story.append(para)
    
    def _add_report_footer(self, story: List, report: Report):
        """Add report footer with disclaimers"""
        
        story.append(PageBreak())
        
        disclaimer_text = """
        IMPORTANT DISCLAIMER: This report is for informational purposes only and does not constitute 
        investment advice. Past performance does not guarantee future results. All projections are 
        hypothetical and based on assumptions that may not reflect actual future conditions. 
        Please consult with your financial advisor before making investment decisions.
        """
        
        disclaimer = Paragraph(disclaimer_text, self.styles['Disclaimer'])
        story.append(disclaimer)
    
    def generate_pptx_report(self, report: Report) -> Dict[str, Any]:
        """
        Generate PowerPoint report from Report object using PowerPointGeneratorService
        """
        try:
            # Process template and get structured content
            processed_content = self.template_processor.process_template(report)
            
            # Get scenario data
            scenario_data = self.scenario_service.get_report_data(report)
            
            # Prepare sections for PowerPoint generation
            sections = self._prepare_powerpoint_sections(processed_content, report)
            
            # Prepare data dictionary
            data = {
                'client': {
                    'first_name': report.client.first_name if report.client else 'Client',
                    'last_name': report.client.last_name if report.client else 'Name',
                    'full_name': report.client.full_name if report.client else 'Client Name',
                    'age': report.client.age if report.client else None,
                    'retirement_age': report.client.retirement_age if report.client else None,
                    'current_income': report.client.current_income if report.client else None,
                    'current_savings': report.client.current_savings if report.client else None,
                },
                'scenario': {
                    'name': report.scenario.name if report.scenario else 'Default Scenario',
                    'retirement_age': report.scenario.retirement_age if report.scenario else None,
                    'final_age': report.scenario.final_age if report.scenario else None,
                    'current_year': report.scenario.current_year if report.scenario else None,
                } if report.scenario else {},
                'report': {
                    'title': report.title or 'Retirement Planning Report',
                    'created_date': timezone.now().strftime('%B %d, %Y'),
                    'advisor_name': report.created_by.get_full_name() if report.created_by else 'Advisor',
                },
                'scenario_results': scenario_data,
            }
            
            # Prepare options for branding and customization
            options = {
                'branding': {
                    'primary_color': '#3498db',
                    'secondary_color': '#2c3e50',
                    'accent_color': '#e74c3c',
                    'logo_path': None,  # Can be added later
                },
                'include_charts': True,
                'chart_style': 'professional',
            }
            
            # Generate PowerPoint using the comprehensive service
            powerpoint_bytes, filename = self.powerpoint_service.generate_powerpoint_report(
                sections, data, options
            )
            
            # Ensure PPTX directory exists
            pptx_dir = self.storage_config.get('PPTX_REPORTS', 
                os.path.join(settings.MEDIA_ROOT, 'report_center/generated_reports/pptx/'))
            os.makedirs(pptx_dir, exist_ok=True)
            
            # Save to file
            file_path = os.path.join(pptx_dir, filename)
            with open(file_path, 'wb') as f:
                f.write(powerpoint_bytes)
            
            # Store file using storage service
            storage_result = self.storage_service.store_generated_report(
                report, file_path, 'pptx'
            )
            
            if storage_result['success']:
                # Clean up temporary file
                try:
                    os.remove(file_path)
                except OSError:
                    pass
                
                logger.info(f"PowerPoint report generated and stored successfully for report {report.id}")
                return storage_result
            else:
                logger.error(f"PowerPoint storage failed for report {report.id}: {storage_result.get('error')}")
                return storage_result
            
        except Exception as e:
            logger.error(f"PowerPoint generation failed for report {report.id}: {str(e)}")
            return {
                'success': False,
                'error': str(e),
                'format': 'pptx'
            }
    
    def _add_pptx_title_slide(self, prs: Presentation, report: Report):
        """Add title slide to PowerPoint"""
        
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = report.title or "Retirement Planning Report"
        subtitle.text = f"Prepared for {report.client.full_name if report.client else 'Client'}\n{timezone.now().strftime('%B %d, %Y')}"
    
    def _add_pptx_client_slide(self, prs: Presentation, report: Report):
        """Add client overview slide"""
        if not report.client:
            return
        
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Client Overview"
        
        content = slide.placeholders[1]
        client = report.client
        
        content_text = f"""Client: {client.full_name}
Age: {client.age if client.age else 'Not specified'}
Retirement Age: {client.retirement_age if client.retirement_age else 'Not specified'}
Current Income: ${client.current_income:,.0f} if client.current_income else 'Not specified'
Current Savings: ${client.current_savings:,.0f} if client.current_savings else 'Not specified'"""
        
        content.text = content_text
    
    def _add_pptx_scenario_slide(self, prs: Presentation, report: Report):
        """Add scenario analysis slide"""
        if not report.scenario:
            return
        
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Scenario Analysis"
        
        content = slide.placeholders[1]
        scenario = report.scenario
        
        content_text = f"""Scenario: {scenario.name}
Retirement Age: {scenario.retirement_age if scenario.retirement_age else 'Not set'}
Planning Horizon: Age {scenario.final_age if scenario.final_age else 'Not set'}
Analysis Year: {scenario.current_year if scenario.current_year else 'Current'}

This comprehensive analysis examines potential retirement outcomes based on current financial data and market projections."""
        
        content.text = content_text
    
    def _add_pptx_sections(self, prs: Presentation, report: Report):
        """Add dynamic PowerPoint slides based on report sections"""
        
        sections = report.sections.filter(is_enabled=True).order_by('order')
        
        for section in sections:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = section.title or section.section_type.replace('_', ' ').title()
            
            content = slide.placeholders[1]
            content_config = section.content_config or {}
            content_text = content_config.get('content', f'Content for {section.section_type} section.')
            content.text = content_text
    
    def _add_pptx_conclusion_slide(self, prs: Presentation, report: Report):
        """Add conclusion slide"""
        
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Next Steps"
        
        content = slide.placeholders[1]
        content.text = """Thank you for reviewing your retirement planning analysis.

Key Actions:
• Schedule follow-up meeting to discuss recommendations
• Review and adjust investment allocations as needed
• Monitor tax planning opportunities throughout the year
• Update projections annually or with major life changes

Contact your advisor to begin implementing these strategies."""
    
    def _prepare_powerpoint_sections(self, processed_content: Dict, report: Report) -> List[Dict[str, Any]]:
        """
        Prepare sections data structure for PowerPointGeneratorService
        """
        sections = []
        
        # Cover slide section
        sections.append({
            'type': 'cover',
            'title': report.title or 'Retirement Planning Report',
            'subtitle': f"Prepared for {report.client.full_name if report.client else 'Client'}\n{timezone.now().strftime('%B %d, %Y')}",
            'options': {'template': 'professional'}
        })
        
        # Executive summary section
        sections.append({
            'type': 'executive_summary',
            'title': 'Executive Summary',
            'content': {
                'summary_text': f"Based on our comprehensive analysis, we have identified key opportunities to optimize {report.client.full_name if report.client else 'your'} retirement income strategy.",
                'key_findings': [
                    'Portfolio positioned for long-term growth',
                    'Tax optimization opportunities identified',
                    'Income replacement strategy aligned with goals'
                ]
            }
        })
        
        # Add sections from processed template content
        for section_data in processed_content.get('sections', []):
            sections.append({
                'type': section_data.get('type', 'generic'),
                'title': section_data.get('title', 'Analysis'),
                'content': section_data.get('content', {}),
                'charts': section_data.get('charts', [])
            })
        
        # Financial timeline section
        if report.scenario:
            sections.append({
                'type': 'financial_timeline',
                'title': 'Financial Projections Timeline',
                'content': {
                    'scenario_name': report.scenario.name,
                    'retirement_age': report.scenario.retirement_age,
                    'final_age': report.scenario.final_age,
                }
            })
        
        # Monte Carlo analysis section
        sections.append({
            'type': 'monte_carlo_analysis',
            'title': 'Monte Carlo Simulation Results',
            'content': {
                'success_probability': '85%',
                'confidence_intervals': {
                    '10th': '$750,000',
                    '50th': '$1,250,000',
                    '90th': '$2,100,000'
                }
            }
        })
        
        # Tax analysis section
        sections.append({
            'type': 'tax_analysis',
            'title': 'Tax Optimization Analysis',
            'content': {
                'current_tax_burden': '$18,500',
                'optimized_tax_burden': '$15,200',
                'annual_savings': '$3,300',
                'irmaa_considerations': True
            }
        })
        
        # Recommendations section
        sections.append({
            'type': 'recommendations',
            'title': 'Strategic Recommendations',
            'content': {
                'recommendations': [
                    {'priority': 'High', 'recommendation': 'Maximize 401(k) contributions to take advantage of employer matching'},
                    {'priority': 'Medium', 'recommendation': 'Consider Roth conversions during lower-income retirement years'},
                    {'priority': 'Medium', 'recommendation': 'Rebalance portfolio to age-appropriate asset allocation'},
                    {'priority': 'Low', 'recommendation': 'Monitor IRMAA thresholds to minimize Medicare surcharges'}
                ]
            }
        })
        
        return sections


class ChartExportService:
    """
    Service for exporting Chart.js visualizations to include in reports
    """
    
    def __init__(self):
        self.chart_storage = getattr(settings, 'REPORT_CENTER_STORAGE', {}).get(
            'CHART_EXPORTS', 
            os.path.join(settings.MEDIA_ROOT, 'report_center/chart_exports/')
        )
        os.makedirs(self.chart_storage, exist_ok=True)
    
    def export_chart_to_image(self, chart_config: Dict, chart_id: str) -> Dict[str, Any]:
        """
        Export Chart.js configuration to image file for inclusion in reports
        """
        try:
            # This will be implemented to work with frontend chart generation
            # For now, return placeholder
            
            placeholder_path = os.path.join(self.chart_storage, f"chart_{chart_id}.png")
            
            # Create placeholder image file
            from PIL import Image, ImageDraw, ImageFont
            
            img = Image.new('RGB', (800, 400), color='white')
            draw = ImageDraw.Draw(img)
            
            # Add placeholder text
            try:
                font = ImageFont.truetype("arial.ttf", 20)
            except:
                font = ImageFont.load_default()
            
            text = f"Chart Placeholder\nChart ID: {chart_id}"
            draw.text((350, 180), text, fill='black', font=font, anchor='mm')
            
            img.save(placeholder_path)
            
            return {
                'success': True,
                'image_path': placeholder_path,
                'chart_id': chart_id
            }
            
        except Exception as e:
            logger.error(f"Chart export failed for {chart_id}: {str(e)}")
            return {
                'success': False,
                'error': str(e),
                'chart_id': chart_id
            }